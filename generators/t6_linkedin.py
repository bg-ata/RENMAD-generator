"""
T6 — LinkedIn / Webinar Promotional Banner (1200×627px)
Full-bleed background photo + content block left + speakers bottom row.
"""
import io
from PIL import Image, ImageDraw, ImageFont
from utils.image_utils import make_circular, hex_to_rgb, load_image, resize_keep_aspect
from config import LANGUAGE_STRINGS

W, H = 1200, 627

# Layout constants
TITLE_BLOCK_RIGHT = 680   # content fills left ~55%
SPEAKER_ROW_TOP = 410
PHOTO_SIZE = 110
PHOTO_SPACING = 10


def _get_font(name: str, size: int):
    """Load brand font via central utility."""
    from utils.fonts import get
    mapping = {
        "heavy":   "montserrat-black",
        "bold":    "montserrat-bold",
        "regular": "inter",
    }
    return get(mapping.get(name, "inter"), size)


def _wrap_text(draw, text, font, max_width):
    words = text.split()
    lines = []
    current = ""
    for word in words:
        test = (current + " " + word).strip()
        bbox = draw.textbbox((0, 0), test, font=font)
        if bbox[2] > max_width and current:
            lines.append(current)
            current = word
        else:
            current = test
    if current:
        lines.append(current)
    return lines


def _darken_region(img: Image.Image, x1, y1, x2, y2, alpha=110) -> Image.Image:
    """Apply a semi-transparent dark overlay over a region."""
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    draw.rectangle([x1, y1, x2, y2], fill=(0, 0, 0, alpha))
    return Image.alpha_composite(img.convert("RGBA"), overlay)


def generate(session: dict, theme: dict, language: str = "es",
             background_img: Image.Image = None,
             renmad_logo: Image.Image = None) -> tuple[bytes, bytes]:
    """
    Returns (pptx_bytes, png_bytes).
    session keys: title, date_str, time_str, speakers=[{name, title, company, is_moderator, photo}]
    theme keys: hex, rgb, label_es, label_en
    """
    strings = LANGUAGE_STRINGS.get(language, LANGUAGE_STRINGS["es"])
    theme_rgb = tuple(theme["rgb"])
    theme_hex = theme["hex"]

    # ── Build PIL image ───────────────────────────────────────────────────────
    if background_img:
        bg = background_img.convert("RGBA").resize((W, H), Image.LANCZOS)
    else:
        bg = Image.new("RGBA", (W, H), (30, 30, 60))

    # Darken left content area so text is legible
    bg = _darken_region(bg, 0, 0, TITLE_BLOCK_RIGHT, H, alpha=140)
    # Lighten bottom speaker strip slightly
    speaker_overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    sod = ImageDraw.Draw(speaker_overlay)
    sod.rectangle([0, SPEAKER_ROW_TOP - 20, W, H], fill=(0, 0, 0, 80))
    bg = Image.alpha_composite(bg, speaker_overlay)

    draw = ImageDraw.Draw(bg)

    # ── "WEBINAR GRATUITO" badge ──────────────────────────────────────────────
    badge_font = _get_font("bold", 20)
    badge_text = strings["free_webinar"]
    badge_pad_x, badge_pad_y = 14, 7
    bbox = draw.textbbox((0, 0), badge_text, font=badge_font)
    badge_w = bbox[2] + badge_pad_x * 2
    badge_h = bbox[3] + badge_pad_y * 2
    badge_x, badge_y = 30, 28
    draw.rectangle([badge_x, badge_y, badge_x + badge_w, badge_y + badge_h], fill=theme_rgb)
    draw.text((badge_x + badge_pad_x, badge_y + badge_pad_y), badge_text,
              font=badge_font, fill=(255, 255, 255))

    # ── Session title ─────────────────────────────────────────────────────────
    title_font = _get_font("heavy", 42)
    title_y = badge_y + badge_h + 22
    max_title_w = TITLE_BLOCK_RIGHT - 50
    lines = _wrap_text(draw, session.get("title", "").upper(), title_font, max_title_w)
    for line in lines[:3]:
        draw.text((30, title_y), line, font=title_font, fill=(255, 255, 255))
        lh = draw.textbbox((0, 0), line, font=title_font)[3]
        title_y += lh + 6

    # ── Date / time pill ──────────────────────────────────────────────────────
    date_font = _get_font("bold", 24)
    date_str = session.get("date_str", "")
    time_str = session.get("time_str", "")
    pill_text = f"{date_str}  |  {time_str}" if date_str else time_str
    pill_y = title_y + 18
    pill_pad_x, pill_pad_y = 16, 8
    bbox = draw.textbbox((0, 0), pill_text, font=date_font)
    pill_w = bbox[2] + pill_pad_x * 2
    pill_h = bbox[3] + pill_pad_y * 2
    draw.rounded_rectangle([30, pill_y, 30 + pill_w, pill_y + pill_h],
                            radius=6, fill=theme_rgb)
    draw.text((30 + pill_pad_x, pill_y + pill_pad_y), pill_text,
              font=date_font, fill=(255, 255, 255))

    # ── RENMAD logo (top-right) ───────────────────────────────────────────────
    if renmad_logo:
        logo = renmad_logo.convert("RGBA")
        logo_h = 72
        ratio = logo.width / logo.height
        logo_w = int(logo_h * ratio)
        logo = logo.resize((logo_w, logo_h), Image.LANCZOS)
        logo_x = W - logo_w - 20
        bg.paste(logo, (logo_x, 18), mask=logo)

    # ── Speaker row ───────────────────────────────────────────────────────────
    speakers = session.get("speakers", [])
    if speakers:
        n = len(speakers)
        total_w = n * PHOTO_SIZE + (n - 1) * PHOTO_SPACING
        # Centre within right portion or full width
        start_x = (W - total_w) // 2

        name_font = _get_font("bold", 14)
        title_font_s = _get_font("regular", 11)
        company_font = _get_font("bold", 12)
        mod_font = _get_font("bold", 10)

        for i, spk in enumerate(speakers):
            cx = start_x + i * (PHOTO_SIZE + PHOTO_SPACING)
            cy = SPEAKER_ROW_TOP

            # Circular photo
            if spk.get("photo"):
                try:
                    photo = make_circular(spk["photo"].convert("RGBA"), PHOTO_SIZE,
                                         border_color=theme_rgb, border_width=3)
                    bg.paste(photo, (cx, cy), mask=photo)
                except Exception:
                    pass

            # Moderator tag
            if spk.get("is_moderator"):
                mod_label = f"[{strings['moderator'].upper()}]"
                mbbox = draw.textbbox((0, 0), mod_label, font=mod_font)
                mx = cx + (PHOTO_SIZE - mbbox[2]) // 2
                draw.text((mx, cy - 16), mod_label, font=mod_font, fill=theme_rgb + (255,))

            text_y = cy + PHOTO_SIZE + 6
            text_cx = cx + PHOTO_SIZE // 2

            # Company logo below photo (if provided)
            if spk.get("company_logo"):
                try:
                    clogo = spk["company_logo"].convert("RGBA")
                    clogo.thumbnail((PHOTO_SIZE, 30), Image.LANCZOS)
                    lx = cx + (PHOTO_SIZE - clogo.width) // 2
                    bg.paste(clogo, (lx, text_y), mask=clogo)
                    text_y += clogo.height + 4
                except Exception:
                    pass

            # Speaker name
            name = spk.get("name", "")
            nbbox = draw.textbbox((0, 0), name, font=name_font)
            nx = text_cx - nbbox[2] // 2
            draw.text((nx, text_y), name, font=name_font, fill=(255, 255, 255))
            text_y += nbbox[3] + 2

            # Job title (truncated)
            job = spk.get("title", "")
            if len(job) > 35:
                job = job[:33] + "…"
            jbbox = draw.textbbox((0, 0), job, font=title_font_s)
            jx = text_cx - jbbox[2] // 2
            draw.text((jx, text_y), job, font=title_font_s, fill=(200, 200, 200))
            text_y += jbbox[3] + 1

            # Company name
            company = spk.get("company", "")
            if len(company) > 25:
                company = company[:23] + "…"
            cbbox = draw.textbbox((0, 0), company, font=company_font)
            cx2 = text_cx - cbbox[2] // 2
            draw.text((cx2, text_y), company, font=company_font, fill=theme_rgb + (255,))

    # ── Convert to RGB and export PNG ────────────────────────────────────────
    final = bg.convert("RGB")
    png_buf = io.BytesIO()
    final.save(png_buf, format="PNG", dpi=(150, 150))
    png_bytes = png_buf.getvalue()

    # ── Build PPTX ───────────────────────────────────────────────────────────
    pptx_bytes = _build_pptx(final, W, H)

    return pptx_bytes, png_bytes


def _build_pptx(img: Image.Image, w: int, h: int) -> bytes:
    """Wrap the rendered PIL image into a single-slide PPTX."""
    from pptx import Presentation
    from pptx.util import Emu
    import io

    DPI = 96
    def px(p): return Emu(int(p / DPI * 914400))

    prs = Presentation()
    prs.slide_width = px(w)
    prs.slide_height = px(h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    img_buf = io.BytesIO()
    img.save(img_buf, format="PNG")
    img_buf.seek(0)
    slide.shapes.add_picture(img_buf, 0, 0, px(w), px(h))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
