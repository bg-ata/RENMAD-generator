"""T3 — Full-colour Announcement Slide (1920×1080px)."""
import io
from PIL import Image, ImageDraw, ImageFont
from config import LANGUAGE_STRINGS

W, H = 1920, 1080

ANNOUNCEMENTS = {
    "break":      {"icon": "☕", "key": "break"},
    "lunch":      {"icon": "🍽", "key": "lunch"},
    "end_of_day": {"icon": "🌅", "key": "end_of_day"},
    "mute_phone": {"icon": "🔇", "key": "mute_phone"},
    "welcome":    {"icon": "👋", "key": "welcome"},
}


def _get_font(name, size):
    import os
    c = {
        "heavy": ["C:/Windows/Fonts/impact.ttf", "C:/Windows/Fonts/arialbd.ttf"],
        "bold":  ["C:/Windows/Fonts/calibrib.ttf", "C:/Windows/Fonts/arialbd.ttf"],
    }
    for p in c.get(name, c["bold"]):
        if os.path.exists(p):
            try: return ImageFont.truetype(p, size)
            except: pass
    return ImageFont.load_default()


def generate(announcement_type: str, theme: dict, language: str = "es",
             custom_text: str = "") -> tuple[bytes, bytes]:
    strings = LANGUAGE_STRINGS.get(language, LANGUAGE_STRINGS["es"])
    theme_rgb = tuple(theme["rgb"])

    canvas = Image.new("RGB", (W, H), theme_rgb)
    draw = ImageDraw.Draw(canvas)

    ann = ANNOUNCEMENTS.get(announcement_type, {"icon": "📢", "key": "welcome"})
    main_text = custom_text or strings.get(ann["key"], announcement_type.upper())

    title_font = _get_font("heavy", 120)
    sub_font = _get_font("bold", 48)

    # Main text
    tb = draw.textbbox((0, 0), main_text.upper(), font=title_font)
    draw.text(((W - tb[2]) // 2, (H - tb[3]) // 2 - 40),
              main_text.upper(), font=title_font, fill=(255, 255, 255))

    # Subtle secondary text at 50% opacity
    sub = "RENMAD"
    sb = draw.textbbox((0, 0), sub, font=sub_font)
    overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    od = ImageDraw.Draw(overlay)
    od.text(((W - sb[2]) // 2, (H - sb[3]) // 2 + 100),
            sub, font=sub_font, fill=(255, 255, 255, 128))
    canvas = Image.alpha_composite(canvas.convert("RGBA"), overlay).convert("RGB")

    png_buf = io.BytesIO()
    canvas.save(png_buf, format="PNG", dpi=(150, 150))
    png_bytes = png_buf.getvalue()
    pptx_bytes = _build_pptx(canvas, W, H)
    return pptx_bytes, png_bytes


def _build_pptx(img, w, h):
    from pptx import Presentation
    from pptx.util import Emu
    import io
    DPI = 96
    def px(p): return Emu(int(p / DPI * 914400))
    prs = Presentation()
    prs.slide_width = px(w)
    prs.slide_height = px(h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    buf = io.BytesIO()
    img.save(buf, format="PNG"); buf.seek(0)
    slide.shapes.add_picture(buf, 0, 0, px(w), px(h))
    out = io.BytesIO(); prs.save(out)
    return out.getvalue()
