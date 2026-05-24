"""
Miniature — Website Card   1200 × 675

Layout
──────
 ┌──────────────────────────────────────────────────────────────┐  y=0
 │  accent bar (5 px)                            [event logo]  │
 │  SESSION TITLE  (centred)                                   │
 │  [date pill]  [time pill]  (centred)                        │
 │                                                              │
 │  [photo] [photo] [photo]   ← fills space down to split     │
 ├── theme colour line (4 px) ──────────────────────────────────┤  y=SPLIT_Y=540
 │  WHITE STRIP — logo pills flush to split line               │
 └──────────────────────────────────────────────────────────────┘  y=675

Template A  →  circular crop  (make_circular)
Template B  →  rembg person cutout  (falls back to coloured silhouette)
Template C  →  rounded-square crop  (ImageOps.fit + mask)

PPTX output: background (everything except logos) + logos as individual
             moveable/resizable picture shapes — so logos can be adjusted
             in PowerPoint without regenerating.

No speaker names · no job titles · no CTA.
"""
import io, math
from PIL import Image, ImageDraw, ImageOps
from utils.image_utils import make_circular
from utils.text_utils import wrap_title
from config import LANGUAGE_STRINGS

W, H            = 1200, 675

WHITE_STRIP_H   = 135           # thick white logo strip
SPLIT_Y         = H - WHITE_STRIP_H   # = 540

BORDER_W        = 4
PHOTO_GAP       = 18
MAX_PHOTO_D     = 230
MIN_PHOTO_D     = 80
PHOTO_R         = 14            # corner radius for square template

LOGO_H_TGT      = 60            # logo content height — large and clear
LOGO_PAD        = 15
LOGO_PILL_H     = LOGO_H_TGT + LOGO_PAD * 2   # = 90
LOGO_Y_OFFSET   = 10            # gap from accent line to top of logo pill
MAX_LOGO_W      = 280           # hard cap on logo content width (any template)

TITLE_TOP       = 68
BODY_KEEP_PCT   = 0.65          # fraction of person to keep (head+torso)

EMU_PX          = 914400 // 96  # 9525 EMU per pixel at 96 DPI


# ── Fonts ─────────────────────────────────────────────────────────────────────

def _f(name, size):
    from utils.fonts import get
    return get({"heavy": "montserrat-black", "bold": "montserrat-bold",
                "reg": "inter"}.get(name, "inter"), size)


def _wrap(draw, text, font, max_w):
    words, lines, cur = text.split(), [], ""
    for w in words:
        t = (cur + " " + w).strip()
        if draw.textbbox((0, 0), t, font=font)[2] > max_w and cur:
            lines.append(cur); cur = w
        else:
            cur = t
    if cur:
        lines.append(cur)
    return lines


# ── Logo helpers ──────────────────────────────────────────────────────────────

def _trim_logo(img):
    """Remove transparent border padding and near-white margins before scaling."""
    from PIL import ImageChops as IC
    rgba = img.convert("RGBA")
    bbox = rgba.split()[3].getbbox()
    if bbox:
        rgba = rgba.crop(bbox)
    r, g, b, a = rgba.split()
    # Pixels where ALL channels < 240 are considered content (not near-white)
    dark = IC.darker(IC.darker(r, g), b)
    wb = dark.point(lambda v: 255 if v < 240 else 0).getbbox()
    if wb:
        rgba = rgba.crop(wb)
    return rgba


def _logo_pill(logo_img, col_budget):
    """White pill with logo box-fitted within col_budget × LOGO_H_TGT.

    All pills share the same LOGO_PILL_H height regardless of logo aspect
    ratio — wide logos are scaled DOWN so they never overflow their column.
    The logo is centred inside the pill both horizontally and vertically.

    Returns a PIL RGBA image, or None if the logo is empty.
    """
    try:
        clogo = _trim_logo(logo_img)
    except Exception:
        clogo = logo_img.convert("RGBA")
    lw, lh = clogo.size
    if lh == 0:
        return None

    # Content box: col_budget wide (minus padding), LOGO_H_TGT tall
    max_cw = max(20, min(col_budget - LOGO_PAD * 2, MAX_LOGO_W))
    max_ch = LOGO_H_TGT

    # Box-fit: use whichever axis is the tighter constraint
    scale = min(max_ch / lh, max_cw / lw)
    nw = max(1, int(lw * scale))
    nh = max(1, int(lh * scale))
    clogo = clogo.resize((nw, nh), Image.LANCZOS)

    # Pill width adapts to content; pill height is always LOGO_PILL_H
    pill_w = nw + LOGO_PAD * 2
    pill   = Image.new("RGBA", (pill_w, LOGO_PILL_H), (0, 0, 0, 0))
    ImageDraw.Draw(pill).rounded_rectangle(
        [0, 0, pill_w - 1, LOGO_PILL_H - 1], radius=12, fill=(255, 255, 255, 255))
    # Centre logo inside the pill
    px = (pill_w - nw) // 2
    py = (LOGO_PILL_H - nh) // 2
    pill.paste(clogo, (px, py), mask=clogo)
    return pill


def _logo_pill_safe(logo_img, col_budget):
    if logo_img is None:
        return None
    try:
        return _logo_pill(logo_img, col_budget)
    except Exception:
        return None


# ── Photo renderers ───────────────────────────────────────────────────────────

# ── Template A: circles ───────────────────────────────────────────────────────

def _photo_circle(photo_img, d, theme_rgb):
    return make_circular(photo_img.convert("RGBA"), d,
                         border_color=theme_rgb, border_width=BORDER_W)


def _sil_circle(d, theme_rgb):
    """Placeholder circle silhouette."""
    total = d + BORDER_W * 2
    cv = Image.new("RGBA", (total, total), (0, 0, 0, 0))
    dr = ImageDraw.Draw(cv)
    dr.ellipse([0, 0, total-1, total-1], fill=(*theme_rgb, 255))
    dr.ellipse([BORDER_W, BORDER_W, BORDER_W+d-1, BORDER_W+d-1], fill=(32, 32, 40, 255))
    cx = BORDER_W + d // 2; r = d // 2
    tr, tg, tb = theme_rgb; col = (tr, tg, tb, 170)
    hr = int(r * 0.22); hy = BORDER_W + int(r * 0.38)
    dr.ellipse([cx-hr, hy-hr, cx+hr, hy+hr], fill=col)
    bwd = int(r * 0.38); by = hy + hr + 2
    dr.ellipse([cx-bwd, by, cx+bwd, by+int(r*0.72)], fill=col)
    return cv


# ── Template B: person cutouts ────────────────────────────────────────────────

def _photo_cutout(photo_img, col_w, max_h):
    """rembg person cutout scaled to fit col_w × max_h. Returns RGBA or None."""
    try:
        from rembg import remove
        cutout = remove(photo_img.convert("RGBA"))
        alpha  = cutout.split()[3]
        bbox   = alpha.point(lambda v: 255 if v > 20 else 0).getbbox()
        if not bbox:
            return None
        cutout = cutout.crop(bbox)
        pw, ph = cutout.size
        # Keep head + torso
        keep_h = max(1, int(ph * BODY_KEEP_PCT))
        cutout = cutout.crop((0, 0, pw, keep_h))
        pw, ph = cutout.size
        # Scale so person fills max_h
        scale  = max_h / ph
        new_w  = max(1, int(pw * scale))
        cutout = cutout.resize((new_w, max_h), Image.LANCZOS)
        # Clamp width to column
        limit  = col_w - 20
        if new_w > limit:
            left   = (new_w - limit) // 2
            cutout = cutout.crop((left, 0, left + limit, max_h))
        return cutout
    except Exception:
        return None


def _sil_person(w, h, theme_rgb):
    """Coloured person-shape silhouette at arbitrary size (rembg fallback)."""
    cv   = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    mask = Image.new("L",    (w, h), 0)
    md   = ImageDraw.Draw(mask)
    cx   = w // 2
    hr   = int(w * 0.18); hcy = hr + 8
    md.ellipse([cx-hr, hcy-hr, cx+hr, hcy+hr], fill=215)
    body_top = hcy + hr + 5
    md.rounded_rectangle([int(w*0.05), body_top, int(w*0.95), h],
                         radius=int(w * 0.28), fill=215)
    col = Image.new("RGBA", (w, h), (*theme_rgb, 255))
    col.putalpha(mask)
    cv.paste(col, (0, 0), mask=col)
    return cv


# ── Template C: rounded squares ───────────────────────────────────────────────

def _photo_square(photo_img, d, theme_rgb):
    total = d + BORDER_W * 2
    out   = Image.new("RGBA", (total, total), (0, 0, 0, 0))
    dr    = ImageDraw.Draw(out)
    dr.rounded_rectangle([0, 0, total-1, total-1],
                         radius=PHOTO_R + BORDER_W, fill=(*theme_rgb, 255))
    raw  = ImageOps.fit(photo_img.convert("RGBA"), (d, d), Image.LANCZOS)
    mask = Image.new("L", (d, d), 0)
    ImageDraw.Draw(mask).rounded_rectangle([0, 0, d-1, d-1], radius=PHOTO_R, fill=255)
    raw.putalpha(mask)
    out.paste(raw, (BORDER_W, BORDER_W), mask=raw.split()[3])
    return out


def _sil_square(d, theme_rgb):
    total = d + BORDER_W * 2
    out   = Image.new("RGBA", (total, total), (0, 0, 0, 0))
    dr    = ImageDraw.Draw(out)
    dr.rounded_rectangle([0, 0, total-1, total-1],
                         radius=PHOTO_R + BORDER_W, fill=(*theme_rgb, 255))
    dr.rounded_rectangle([BORDER_W, BORDER_W, BORDER_W+d-1, BORDER_W+d-1],
                         radius=PHOTO_R, fill=(32, 32, 40, 255))
    tr, tg, tb = theme_rgb; col = (tr, tg, tb, 170)
    cx = total // 2
    hrad = int(d * 0.13); hy = BORDER_W + int(d * 0.28)
    dr.ellipse([cx-hrad, hy-hrad, cx+hrad, hy+hrad], fill=col)
    bw2 = int(d * 0.26); by = hy + hrad + 2
    dr.ellipse([cx-bw2, by, cx+bw2, by + int(d * 0.46)], fill=col)
    return out


# ── PPTX builder ─────────────────────────────────────────────────────────────

def _build_pptx(bg_png_bytes: bytes,
                logo_placements: list[tuple]) -> bytes:
    """Create a PPTX with:
       • Full-slide background PNG (everything except logos, baked in)
       • Each logo as a separate, resizable/moveable picture shape
    Returns raw PPTX bytes.
    """
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width  = Emu(W * EMU_PX)
    prs.slide_height = Emu(H * EMU_PX)

    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout

    # ── Background (full canvas without logos) ────────────────────────────────
    bg_io = io.BytesIO(bg_png_bytes)
    slide.shapes.add_picture(bg_io, 0, 0,
                             prs.slide_width, prs.slide_height)

    # ── Logo pills — individual resizable shapes ──────────────────────────────
    for lx, ly, pill in logo_placements:
        ibuf = io.BytesIO()
        pill.save(ibuf, format="PNG")
        ibuf.seek(0)
        slide.shapes.add_picture(
            ibuf,
            Emu(lx * EMU_PX), Emu(ly * EMU_PX),
            Emu(pill.width * EMU_PX), Emu(pill.height * EMU_PX),
        )

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ── Main generator ────────────────────────────────────────────────────────────

def generate(
    session:     dict,
    theme:       dict,
    language:    str               = "es",
    renmad_logo: Image.Image | None = None,
    bg_image:    Image.Image | None = None,
    ata_logo:    Image.Image | None = None,
    template:    str               = "A",
) -> tuple[bytes, bytes]:
    """Return (pptx_bytes, png_bytes).

    PPTX has the background baked in plus logos as individual moveable shapes.
    PNG is the final flat render (identical to what was shown on screen).
    """
    theme_rgb = tuple(theme["rgb"])
    speakers  = session.get("speakers", [])
    n         = max(len(speakers), 1)

    # ── Canvas + full-image dark overlay ──────────────────────────────────────
    canvas = Image.new("RGB", (W, H), (14, 14, 18))
    if bg_image:
        bg = bg_image.convert("RGB").resize((W, H), Image.LANCZOS)
        canvas.paste(bg)
    else:
        d = ImageDraw.Draw(canvas)
        for y in range(H):
            v = int(20 + math.sin(y / H * math.pi) * 14)
            d.line([(0, y), (W, y)], fill=(v, v, v + 4))

    overlay = Image.new("RGBA", (W, H), (8, 8, 14, 168))
    canvas.paste(Image.new("RGB", (W, H), (8, 8, 14)), (0, 0),
                 mask=overlay.split()[3])
    draw = ImageDraw.Draw(canvas)

    # Top accent bar
    draw.rectangle([0, 0, W, 5], fill=theme_rgb)

    # ── Event logo top-right ──────────────────────────────────────────────────
    if renmad_logo:
        logo = renmad_logo.convert("RGBA")
        lh = 64; lw = int(logo.width * lh / logo.height)
        logo = logo.resize((lw, lh), Image.LANCZOS)
        canvas.paste(logo, (W - lw - 20, 12), mask=logo)
    elif ata_logo:
        al = ata_logo.convert("RGBA")
        alh = 52; alw = int(al.width * alh / al.height)
        al = al.resize((alw, alh), Image.LANCZOS)
        canvas.paste(al, (W - alw - 20, 16), mask=al)

    # ── Title (centred) ───────────────────────────────────────────────────────
    # Auto-shrink so the full title always fits — never truncate.
    # Respect explicit '\n' line breaks (preserved by wrap_title).
    title_text = session.get("title", "").upper()
    title_w    = W - 160
    title_f    = _f("heavy", 52)
    t_lines    = wrap_title(draw, title_text, title_f, title_w, max_lines=4)
    # If we got more lines than wanted (truncation), step the font down until
    # the WHOLE title fits in ≤ 3 lines.
    for _size in (48, 44, 40, 36, 32, 28):
        # `wrap_title` would have truncated at max_lines=4 above. Recompute the
        # untruncated wrap directly to detect the real line count.
        f_try = _f("heavy", _size)
        all_lines = wrap_title(draw, title_text, f_try, title_w, max_lines=20)
        if len(all_lines) <= 3:
            title_f = f_try
            t_lines = all_lines
            break
    else:
        # Last resort: keep min font and render every line we got (no truncation)
        title_f = _f("heavy", 28)
        t_lines = wrap_title(draw, title_text, title_f, title_w, max_lines=20)

    ty = TITLE_TOP
    for line in t_lines:
        lb = draw.textbbox((0, 0), line, font=title_f)
        draw.text(((W - lb[2]) // 2, ty), line, font=title_f, fill=(255, 255, 255))
        ty += lb[3] + 8
    ty += 16

    # ── Date / time pills (centred) ───────────────────────────────────────────
    dt_f       = _f("bold", 26)
    pill_items = []
    if session.get("date_str"):
        pill_items.append((session["date_str"].upper(), (20, 20, 20, 220)))
    if session.get("time_str"):
        pill_items.append((session["time_str"].upper(), theme_rgb))

    if pill_items:
        pad_p   = 13; gap_p = 10
        ref_bb  = draw.textbbox((0, 0), "Ag", font=dt_f)
        pill_h  = ref_bb[3] + pad_p * 2
        pill_ws = [draw.textbbox((0, 0), t, font=dt_f)[2] + pad_p * 2
                   for t, _ in pill_items]
        rx = (W - sum(pill_ws) - gap_p * (len(pill_ws) - 1)) // 2
        for (txt, bg_col), pw in zip(pill_items, pill_ws):
            draw.rounded_rectangle([rx, ty, rx + pw, ty + pill_h],
                                   radius=6, fill=bg_col)
            draw.text((rx + pad_p, ty + pad_p), txt, font=dt_f, fill=(255, 255, 255))
            rx += pw + gap_p
        ty += pill_h
    # ty = bottom edge of date/time area

    photo_top_y = ty + 20          # 20 px breathing room below pills
    avail_h     = SPLIT_Y - photo_top_y - 6    # height for photos

    # logo_placements: collected as (lx, logo_y, pill_img) — drawn after white strip
    logo_placements = []
    logo_y          = SPLIT_Y + 4 + LOGO_Y_OFFSET

    # ══════════════════════════════════════════════════════════════════════════
    # Template B — person cutouts (different layout from A/C)
    # ══════════════════════════════════════════════════════════════════════════
    if template == "B":
        col_w      = W // n
        person_h   = max(60, avail_h)
        col_budget = max(60, col_w - 20)

        logo_pills = [_logo_pill_safe(s.get("company_logo"), col_budget)
                      for s in speakers]

        # Pre-process cutouts (rembg model loads once)
        cutouts = []
        for spk in speakers:
            if spk.get("photo"):
                co = _photo_cutout(spk["photo"], col_w, person_h)
                cutouts.append(co)
            else:
                cutouts.append(None)

        for i, (spk, cutout) in enumerate(zip(speakers, cutouts)):
            cx = i * col_w + col_w // 2

            if cutout is not None:
                px = cx - cutout.width // 2
                py = SPLIT_Y - person_h
                canvas.paste(cutout, (px, py), mask=cutout)
            else:
                sil_w = max(60, col_w - 40)
                sil   = _sil_person(sil_w, person_h, theme_rgb)
                px    = cx - sil.width // 2
                canvas.paste(sil, (px, SPLIT_Y - person_h), mask=sil)

        # White strip + accent line
        draw.rectangle([0, SPLIT_Y, W, H], fill=(255, 255, 255))
        draw.rectangle([0, SPLIT_Y, W, SPLIT_Y + 4], fill=theme_rgb)

        # Collect logo placements (don't paste yet — needed for PPTX background)
        for i, pill in enumerate(logo_pills):
            if pill is None:
                continue
            cx = i * col_w + col_w // 2
            lx = max(8, min(W - pill.width - 8, cx - pill.width // 2))
            logo_placements.append((lx, logo_y, pill))

    # ══════════════════════════════════════════════════════════════════════════
    # Templates A and C — fixed-size circles or squares
    # ══════════════════════════════════════════════════════════════════════════
    else:
        # Compute photo diameter: fill available height, capped by width per column
        d_by_h  = max(0, avail_h - BORDER_W * 2)
        d_by_w  = max(0, (W - 60 - PHOTO_GAP * (n - 1)) // n - BORDER_W * 2)
        photo_d = max(MIN_PHOTO_D, min(MAX_PHOTO_D, d_by_h, d_by_w))

        photo_total = photo_d + BORDER_W * 2
        col_budget  = max(60, photo_total + PHOTO_GAP - 10)

        logo_pills = [_logo_pill_safe(s.get("company_logo"), col_budget)
                      for s in speakers]

        row_w   = n * photo_total + (n - 1) * PHOTO_GAP
        start_x = (W - row_w) // 2
        # Bottom of photos flush with SPLIT_Y - 6
        photo_r  = photo_d // 2
        photo_cy = SPLIT_Y - 6 - photo_r - BORDER_W

        for i, spk in enumerate(speakers):
            cx = start_x + i * (photo_total + PHOTO_GAP) + photo_total // 2

            if template == "C":
                if spk.get("photo"):
                    try:   ph_img = _photo_square(spk["photo"], photo_d, theme_rgb)
                    except: ph_img = _sil_square(photo_d, theme_rgb)
                else:
                    ph_img = _sil_square(photo_d, theme_rgb)
            else:   # "A"
                if spk.get("photo"):
                    try:   ph_img = _photo_circle(spk["photo"], photo_d, theme_rgb)
                    except: ph_img = _sil_circle(photo_d, theme_rgb)
                else:
                    ph_img = _sil_circle(photo_d, theme_rgb)

            px = cx - ph_img.width  // 2
            py = photo_cy - ph_img.height // 2
            canvas.paste(ph_img, (px, py), mask=ph_img)

        # White strip + accent line
        draw.rectangle([0, SPLIT_Y, W, H], fill=(255, 255, 255))
        draw.rectangle([0, SPLIT_Y, W, SPLIT_Y + 4], fill=theme_rgb)

        # Collect logo placements
        for i, pill in enumerate(logo_pills):
            if pill is None:
                continue
            cx = start_x + i * (photo_total + PHOTO_GAP) + photo_total // 2
            lx = max(8, min(W - pill.width - 8, cx - pill.width // 2))
            logo_placements.append((lx, logo_y, pill))

    # ── PPTX: background WITHOUT logos (so logos are separate/editable) ───────
    bg_buf = io.BytesIO()
    canvas.save(bg_buf, format="PNG")
    pptx_bytes = _build_pptx(bg_buf.getvalue(), logo_placements)

    # ── PNG: paste logos onto canvas ──────────────────────────────────────────
    for lx, ly, pill in logo_placements:
        canvas.paste(pill, (lx, ly), mask=pill)

    png_buf = io.BytesIO()
    canvas.save(png_buf, format="PNG", dpi=(96, 96))

    return pptx_bytes, png_buf.getvalue()
