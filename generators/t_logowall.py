"""T_Logowall — Event logo wall  1920 px wide, dynamic height

Layout  (space-filling, biggest at the top)
───────────────────────────────────────────
  WHITE background throughout. Each SECTION has a full-width theme-colour header
  band (white centred label) with its logos below.

  To avoid half-empty rows, sparse sections are packed TWO-PER-ROW (two columns
  under one band, each with its own label), exactly like the brand footers.
  Sections with many logos (and the most important tier + speakers) take the
  full width. Within every card the logos are scaled UP to fill the card width
  (capped per tier so the hierarchy holds), so a tier with one or two logos no
  longer leaves the whole line empty.

  Section order: sponsor tiers (biggest first) → Confirmed Speakers →
  secondary groups (Media / Knowledge / Host partners) at the END.

Logo sizing
───────────
  Each logo is TRIMMED of its surrounding white / transparent margin (so only
  the real artwork counts), then scaled to EQUAL VISUAL AREA within its card
  (geometric-mean normalisation) before the card is filled — so a wide wordmark
  and a square badge look similarly prominent.

Output
──────
  PNG, PDF and a fully-editable PPTX (each logo a separate movable picture, each
  label a real text box) — all from one shared layout plan.

Public API:
  generate(event, tiers, secondary, speakers, theme, language, renmad_logo) -> PNG bytes
  generate_pack(...)          -> {"png": bytes, "pdf": bytes, "pptx": bytes}
  generate_all_variants(...)  -> {"es": pack, "en": pack}
"""

import io
import math
from PIL import Image, ImageDraw, ImageChops

# ── Canvas ────────────────────────────────────────────────────────────────────
W = 1920

# ── Layout constants ──────────────────────────────────────────────────────────
PAD_X          = 60
USABLE_W       = W - PAD_X * 2          # 1800
GUTTER         = 56                     # gap between the two columns
HALF_W         = (USABLE_W - GUTTER) // 2
BAND_H         = 54                     # full-width coloured section header band
BAND_FONT      = 30                     # section header text size (px)
PAD_AFTER_BAND = 30                     # space between a header band and its logos
SECTION_GAP    = 46                     # vertical space between section rows
GAP_X          = 70                     # min horizontal gap between logos in a row
GAP_JUSTIFY_MAX = 260                   # max gap when spreading logos to fill a row
ROW_GAP        = 36                     # vertical gap between wrapped logo rows
BOTTOM_PAD     = 48

WIDE_MIN_LOGOS = 4                      # a section with ≥ this many logos takes full width

# ── Size hierarchy ──────────────────────────────────────────────────────────--
# A "nominal size" (geometric-mean side for equal-area scaling) is assigned to
# each PRESENT tier, in importance order. Logos are then scaled UP to fill their
# card, capped at nominal * MAXH_MULT, so higher tiers stay visibly bigger.
H_FIRST_TIER = 200    # nominal of the most important present tier
TIER_RATIO   = 0.74   # each next tier a bit smaller
H_TIER_MIN   = 120    # floor so low tiers stay substantial
MAXH_MULT    = 1.45   # max logo height = tier nominal * this

H_SEC        = 112    # secondary groups (Media Partner, Host, …)
SEC_MAXH     = 168
H_SPK        = 82     # speaker logos — smallest
SPK_MAXH     = 116

# ── Colours ───────────────────────────────────────────────────────────────────
BG         = (255, 255, 255)
TEXT_WHITE = (255, 255, 255)
DIVIDER    = (230, 230, 234)   # faint vertical rule between two columns

# ── Tier system ───────────────────────────────────────────────────────────────
TIER_ORDER = ["diamond", "platinum", "global", "gold", "silver", "bronze", "standard", "sponsor"]

_SPK_LABEL = {
    "es": "PONENTES CONFIRMADOS",
    "en": "CONFIRMED SPEAKERS",
    "it": "RELATORI CONFERMATI",
    "pl": "POTWIERDZENI PRELEGENCI",
}

# ── Scratch surface for measuring text ──────────────────────────────────────--
_SCRATCH   = Image.new("RGB", (8, 8))
_SCRATCH_D = ImageDraw.Draw(_SCRATCH)

# Pixels → EMU (PowerPoint), at 96 dpi
_PX_EMU = 9525


# ── Font helper ───────────────────────────────────────────────────────────────
def _f(name, size):
    from utils.fonts import get
    return get(
        {"heavy": "montserrat-black",
         "bold":  "montserrat-bold",
         "reg":   "inter"}.get(name, "inter"),
        size,
    )


def _text_size(text: str, font) -> tuple[int, int]:
    b = _SCRATCH_D.textbbox((0, 0), text, font=font)
    return b[2] - b[0], b[3] - b[1]


def _fit_font_size(text: str, region_w: int, start: int = BAND_FONT, lo: int = 16) -> int:
    """Largest header font size (≤ start) whose text fits region_w."""
    size = start
    while size > lo and _text_size(text, _f("heavy", size))[0] > region_w - 24:
        size -= 1
    return size


# ── Tier helpers ──────────────────────────────────────────────────────────────
def _tier_rank(name: str) -> int:
    lower = (name or "").lower()
    for i, t in enumerate(TIER_ORDER):
        if t in lower:
            return i
    return 7


def _tier_sizes(n: int) -> list[int]:
    out, h = [], float(H_FIRST_TIER)
    for _ in range(n):
        out.append(max(H_TIER_MIN, int(round(h))))
        h = h * TIER_RATIO
    return out


# ── Logo helpers ──────────────────────────────────────────────────────────────
def _trim(img: Image.Image) -> Image.Image:
    """Crop away surrounding white / transparent margin so only artwork counts."""
    rgba = img.convert("RGBA")

    alpha = rgba.split()[3]
    abox = alpha.getbbox()
    if abox:
        rgba = rgba.crop(abox)

    try:
        white = Image.new("RGBA", rgba.size, (255, 255, 255, 255))
        comp  = Image.alpha_composite(white, rgba).convert("RGB")
        diff  = ImageChops.difference(comp, Image.new("RGB", comp.size, (255, 255, 255)))
        mask  = diff.convert("L").point(lambda p: 255 if p > 18 else 0)
        cbox  = mask.getbbox()
        if cbox:
            cw, ch = cbox[2] - cbox[0], cbox[3] - cbox[1]
            w, h   = rgba.size
            if cw * ch >= 0.015 * w * h:
                rgba = rgba.crop(cbox)
    except Exception:
        pass

    return rgba


def _paste(canvas: Image.Image, logo: Image.Image, x: int, y: int) -> None:
    if logo.mode != "RGBA":
        logo = logo.convert("RGBA")
    canvas.paste(logo, (x, y), mask=logo)


def _wrap_indices(widths: list, max_w: int) -> list:
    rows, cur, cur_w = [], [], 0
    for i, w in enumerate(widths):
        need = w + (GAP_X if cur else 0)
        if cur and cur_w + need > max_w:
            rows.append(cur)
            cur, cur_w = [i], w
        else:
            cur.append(i)
            cur_w += need
    if cur:
        rows.append(cur)
    return rows


def _card_rows(logos: list, card_w: int, nominal: int, max_h: int) -> list:
    """Equal-area-scale logos, then scale the whole set UP to fill card_w
    (capped at max_h), and wrap into centred rows. Returns list[list[Image]]."""
    trimmed = []
    for lg in logos:
        try:
            t = _trim(lg)
            if t.size[0] > 0 and t.size[1] > 0:
                trimmed.append(t)
        except Exception:
            pass
    if not trimmed:
        return []

    def area_scale(t):
        lw, lh = t.size
        s = nominal / math.sqrt(lw * lh)            # equal visual area
        return min(s, max_h / lh, card_w / lw)      # never exceed caps / card width

    scales  = [area_scale(t) for t in trimmed]
    widths  = [t.size[0] * s for t, s in zip(trimmed, scales)]
    heights = [t.size[1] * s for t, s in zip(trimmed, scales)]

    # Fill: enlarge so the widest row spans card_w, but don't exceed max_h.
    rows_idx  = _wrap_indices(widths, card_w)
    widest    = max(sum(widths[i] for i in r) + GAP_X * (len(r) - 1) for r in rows_idx)
    cur_max_h = max(heights) if heights else 1
    fill = 1.0
    if widest > 0 and cur_max_h > 0:
        fill = min(card_w / widest, max_h / cur_max_h)
    fill = max(fill, 1.0)
    if fill > 1.02:
        scales = [s * fill for s in scales]

    imgs = []
    for t, s in zip(trimmed, scales):
        nw, nh = max(1, int(round(t.size[0] * s))), max(1, int(round(t.size[1] * s)))
        imgs.append(t.resize((nw, nh), Image.LANCZOS))

    rows, cur, cur_w = [], [], 0
    for im in imgs:
        need = im.width + (GAP_X if cur else 0)
        if cur and cur_w + need > card_w:
            rows.append(cur)
            cur, cur_w = [im], im.width
        else:
            cur.append(im)
            cur_w += need
    if cur:
        rows.append(cur)
    return rows


def _rows_block_h(rows: list) -> int:
    if not rows:
        return 0
    return sum(max(im.height for im in r) for r in rows) + ROW_GAP * (len(rows) - 1)


def _emit_card(elements: list, rows: list, x0: int, card_w: int, y: int) -> int:
    """Place logo rows inside [x0, x0+card_w] starting at y. Logos in a row are
    spread out (wider, capped gap) so they fill the width instead of clustering
    in the centre; a single logo is centred."""
    for row in rows:
        n      = len(row)
        total  = sum(im.width for im in row)
        row_h  = max(im.height for im in row)
        if n == 1:
            gap = 0
            x   = x0 + (card_w - total) / 2
        else:
            gap   = max(GAP_X, min((card_w - total) / (n + 1), GAP_JUSTIFY_MAX))
            block = total + gap * (n - 1)
            x     = x0 + (card_w - block) / 2
        for im in row:
            elements.append({"type": "image", "img": im,
                             "x": int(round(x)), "y": int(y + (row_h - im.height) // 2)})
            x += im.width + gap
        y += row_h + ROW_GAP
    return (y - ROW_GAP) if rows else y


# ── Sections ──────────────────────────────────────────────────────────────────
def _build_sections(tiers: list, secondary: list, speakers: list, language: str) -> list:
    active_tiers = [t for t in (tiers or []) if t.get("logos")]
    active_tiers.sort(key=lambda t: _tier_rank(t.get("name", "")))
    sizes = _tier_sizes(len(active_tiers))

    sections = []
    for i, (tier, sz) in enumerate(zip(active_tiers, sizes)):
        logos = tier["logos"]
        sections.append({
            "name":   (tier.get("name", "") or "").upper(),
            "logos":  logos,
            "nominal": sz,
            "max_h":   int(sz * MAXH_MULT),
            # The single most important tier always gets the full width (headline).
            "wide":   (i == 0) or (len(logos) >= WIDE_MIN_LOGOS),
        })

    spk_logos = [lg for s in (speakers or []) if s.get("logos") for lg in s["logos"]]
    if spk_logos:
        sections.append({
            "name":   _SPK_LABEL.get(language, _SPK_LABEL["en"]),
            "logos":  spk_logos,
            "nominal": H_SPK, "max_h": SPK_MAXH, "wide": True,
        })

    for grp in (secondary or []):
        if grp.get("logos"):
            logos = grp["logos"]
            sections.append({
                "name":   (grp.get("name", "") or "").upper(),
                "logos":  logos,
                "nominal": H_SEC, "max_h": SEC_MAXH,
                "wide":   len(logos) >= WIDE_MIN_LOGOS,
            })

    return sections


# ── Layout plan (shared by every output format) ───────────────────────────────
def _band(elements: list, y: int, theme_rgb: tuple) -> None:
    elements.append({"type": "rect", "x": 0, "y": y, "w": W, "h": BAND_H, "color": theme_rgb})


def _band_label(elements: list, text: str, x0: int, region_w: int, y: int) -> None:
    if not text:
        return
    size = _fit_font_size(text, region_w)
    elements.append({"type": "text", "text": text, "x": x0, "y": y,
                     "w": region_w, "h": BAND_H, "size": size,
                     "color": TEXT_WHITE, "anchor": "middle"})


def _build_plan(event, tiers, secondary, speakers, theme, language, renmad_logo=None) -> dict:
    # renmad_logo is intentionally ignored — the wall is pasted where RENMAD
    # branding already exists.
    theme_rgb = tuple(theme["rgb"])
    elements  = []
    y = 0

    sections = _build_sections(tiers, secondary, speakers, language)

    def emit_full(sec, y):
        rows = _card_rows(sec["logos"], USABLE_W, sec["nominal"], sec["max_h"])
        if not rows:
            return y
        _band(elements, y, theme_rgb)
        _band_label(elements, sec["name"], PAD_X, USABLE_W, y)
        yy = y + BAND_H + PAD_AFTER_BAND
        yy = _emit_card(elements, rows, PAD_X, USABLE_W, yy)
        return yy

    def emit_pair(left, right, y):
        rx0   = PAD_X + HALF_W + GUTTER
        lrows = _card_rows(left["logos"],  HALF_W, left["nominal"],  left["max_h"])
        rrows = _card_rows(right["logos"], HALF_W, right["nominal"], right["max_h"])
        _band(elements, y, theme_rgb)
        _band_label(elements, left["name"],  PAD_X, HALF_W, y)
        _band_label(elements, right["name"], rx0,   HALF_W, y)
        top = y + BAND_H + PAD_AFTER_BAND
        yL  = _emit_card(elements, lrows, PAD_X, HALF_W, top)
        yR  = _emit_card(elements, rrows, rx0,   HALF_W, top)
        bottom = max(yL, yR)
        # faint divider between the two columns
        elements.append({"type": "rect", "x": PAD_X + HALF_W + GUTTER // 2 - 1,
                         "y": top, "w": 2, "h": max(1, bottom - top), "color": DIVIDER})
        return bottom

    pending = None
    for sec in sections:
        if sec["wide"]:
            if pending is not None:
                y = emit_full(pending, y) + SECTION_GAP
                pending = None
            y = emit_full(sec, y) + SECTION_GAP
        else:
            if pending is None:
                pending = sec
            else:
                y = emit_pair(pending, sec, y) + SECTION_GAP
                pending = None
    if pending is not None:
        y = emit_full(pending, y) + SECTION_GAP

    if elements:
        y -= SECTION_GAP
    canvas_h = max(120, y + BOTTOM_PAD)
    return {"canvas_h": canvas_h, "elements": elements, "theme_rgb": theme_rgb}


# ── Renderer: raster image (PNG / PDF source) ─────────────────────────────────
def _render_image(plan: dict) -> Image.Image:
    canvas = Image.new("RGB", (W, plan["canvas_h"]), BG)
    draw   = ImageDraw.Draw(canvas)
    for el in plan["elements"]:
        t = el["type"]
        if t == "rect":
            draw.rectangle([el["x"], el["y"], el["x"] + el["w"], el["y"] + el["h"]],
                           fill=el["color"])
        elif t == "image":
            _paste(canvas, el["img"], el["x"], el["y"])
        elif t == "text":
            f = _f("heavy", el["size"])
            tw, th = _text_size(el["text"], f)
            tx = el["x"] + (el["w"] - tw) // 2
            ty = el["y"] + (el["h"] - th) // 2 if el.get("anchor") == "middle" else el["y"]
            draw.text((tx, ty), el["text"], font=f, fill=el["color"])
    return canvas


# ── Renderer: fully-editable PPTX ─────────────────────────────────────────────
def _render_pptx(plan: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width  = Emu(W * _PX_EMU)
    prs.slide_height = Emu(plan["canvas_h"] * _PX_EMU)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    try:
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    except Exception:
        pass

    for el in plan["elements"]:
        t = el["type"]
        if t == "rect":
            sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Emu(el["x"] * _PX_EMU), Emu(el["y"] * _PX_EMU),
                                        Emu(el["w"] * _PX_EMU), Emu(el["h"] * _PX_EMU))
            sh.fill.solid()
            sh.fill.fore_color.rgb = RGBColor(*el["color"])
            sh.line.fill.background()
            try:
                sh.shadow.inherit = False
            except Exception:
                pass
        elif t == "image":
            b = io.BytesIO()
            el["img"].convert("RGBA").save(b, "PNG")
            b.seek(0)
            slide.shapes.add_picture(b, Emu(el["x"] * _PX_EMU), Emu(el["y"] * _PX_EMU),
                                     Emu(el["img"].width * _PX_EMU), Emu(el["img"].height * _PX_EMU))
        elif t == "text":
            tb = slide.shapes.add_textbox(Emu(el["x"] * _PX_EMU), Emu(el["y"] * _PX_EMU),
                                          Emu(el["w"] * _PX_EMU), Emu(el["h"] * _PX_EMU))
            tf = tb.text_frame
            tf.word_wrap = False
            tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
            if el.get("anchor") == "middle":
                try:
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                except Exception:
                    pass
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = el["text"]
            r.font.bold = True
            r.font.size = Pt(el["size"] * 0.75)
            r.font.name = "Montserrat"
            r.font.color.rgb = RGBColor(*el["color"])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Public API ────────────────────────────────────────────────────────────────
def generate(event, tiers, secondary, speakers, theme,
             language="es", renmad_logo=None) -> bytes:
    plan = _build_plan(event, tiers, secondary, speakers, theme, language, renmad_logo)
    buf = io.BytesIO()
    _render_image(plan).save(buf, format="PNG", dpi=(150, 150))
    return buf.getvalue()


def generate_pack(event, tiers, secondary, speakers, theme,
                  language="es", renmad_logo=None) -> dict:
    plan = _build_plan(event, tiers, secondary, speakers, theme, language, renmad_logo)
    img  = _render_image(plan)

    png = io.BytesIO()
    img.save(png, format="PNG", dpi=(150, 150))
    pdf = io.BytesIO()
    img.convert("RGB").save(pdf, format="PDF", resolution=150.0)

    out = {"png": png.getvalue(), "pdf": pdf.getvalue()}
    try:
        out["pptx"] = _render_pptx(plan)
    except Exception:
        out["pptx"] = None
    return out


def generate_all_variants(event, tiers, secondary, speakers, theme,
                          language="es", renmad_logo=None) -> dict:
    return {
        lang: generate_pack(
            event=event, tiers=tiers, secondary=secondary,
            speakers=speakers, theme=theme, language=lang,
            renmad_logo=renmad_logo,
        )
        for lang in ("es", "en")
    }
