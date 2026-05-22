"""
T8 — My ATA Registration Card (1280×720px)
Left: RENMAD stamp + theme bar + date/time + diagonal stripe.
Right: speaker photos (square) + company logos + names.
"""
import io
from PIL import Image, ImageDraw, ImageFont
from utils.image_utils import make_square_crop, hex_to_rgb, load_image
from config import LANGUAGE_STRINGS

W, H = 1280, 720

LEFT_W = 380          # width of left info panel
STRIPE_W = 40         # diagonal stripe visual width
RIGHT_START = LEFT_W + STRIPE_W


def _get_font(name: str, size: int):
    import os
    candidates = {
        "heavy": ["C:/Windows/Fonts/impact.ttf",
                  "C:/Windows/Fonts/FRADM.TTF",
                  "C:/Windows/Fonts/arialbd.ttf"],
        "bold":  ["C:/Windows/Fonts/calibrib.ttf",
                  "C:/Windows/Fonts/arialbd.ttf"],
        "regular": ["C:/Windows/Fonts/calibri.ttf",
                    "C:/Windows/Fonts/arial.ttf"],
    }
    for path in candidates.get(name, candidates["regular"]):
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except Exception:
                pass
    return ImageFont.load_default()


def generate(session: dict, theme: dict, language: str = "es",
             background_img: Image.Image = None,
             renmad_logo: Image.Image = None,
             ata_logo: Image.Image = None) -> tuple[bytes, bytes]:
    """
    Returns (pptx_bytes, png_bytes).
    """
    strings = LANGUAGE_STRINGS.get(language, LANGUAGE_STRINGS["es"])
    theme_rgb = tuple(theme["rgb"])
    theme_label = theme.get("label_es" if language == "es" else "label_en", "")

    # ── Background ────────────────────────────────────────────────────────────
    if background_img:
        bg = background_img.convert("RGBA").resize((W, H), Image.LANCZOS)
    else:
        bg = Image.new("RGBA", (W, H), (20, 30, 50))

    draw = ImageDraw.Draw(bg)

    # ── Left dark overlay ─────────────────────────────────────────────────────
    overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    od = ImageDraw.Draw(overlay)
    od.rectangle([0, 0, LEFT_W, H], fill=(0, 0, 0, 200))
    bg = Image.alpha_composite(bg, overlay)
    draw = ImageDraw.Draw(bg)

    # ── Diagonal stripe (brand colour) ───────────────────────────────────────
    stripe = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    sd = ImageDraw.Draw(stripe)
    sd.polygon([
        (LEFT_W, 0),
        (LEFT_W + STRIPE_W, 0),
        (LEFT_W, H),
    ], fill=theme_rgb + (255,))
    bg = Image.alpha_composite(bg, stripe)
    draw = ImageDraw.Draw(bg)

    # ── Right white card ──────────────────────────────────────────────────────
    card_x = RIGHT_START + 10
    card_w = W - card_x - 20
    card_y = 60
    card_h = H - 80
    draw.rounded_rectangle([card_x, card_y, card_x + card_w, card_y + card_h],
                            radius=12, fill=(255, 255, 255, 230))

    # ── RENMAD wordmark ───────────────────────────────────────────────────────
    if renmad_logo:
        logo = renmad_logo.convert("RGBA")
        logo_h = 80
        logo_w = int(logo.width * logo_h / logo.height)
        logo = logo.resize((logo_w, logo_h), Image.LANCZOS)
        lx = (LEFT_W - logo_w) // 2
        bg.paste(logo, (lx, 40), mask=logo)
        logo_bottom = 40 + logo_h
    else:
        # Text fallback
        wm_font = _get_font("heavy", 54)
        draw.text((20, 30), "RENMAD", font=wm_font, fill=(255, 255, 255))
        logo_bottom = 90

    # ── Theme colour bar with label ───────────────────────────────────────────
    bar_y = logo_bottom + 8
    bar_h = 38
    draw.rectangle([0, bar_y, LEFT_W, bar_y + bar_h], fill=theme_rgb)
    lbl_font = _get_font("bold", 20)
    lbl = theme_label.upper()
    lb = draw.textbbox((0, 0), lbl, font=lbl_font)
    draw.text(((LEFT_W - lb[2]) // 2, bar_y + (bar_h - lb[3]) // 2),
              lbl, font=lbl_font, fill=(255, 255, 255))

    # ── Date / time ───────────────────────────────────────────────────────────
    info_font = _get_font("bold", 18)
    info_y = bar_y + bar_h + 22
    date_str = session.get("date_str", "")
    time_str = session.get("time_str", "")
    if date_str:
        draw.text((20, info_y), f"{strings['date_label']}: {date_str.upper()}",
                  font=info_font, fill=(255, 255, 255))
        info_y += 30
    if time_str:
        draw.text((20, info_y), f"{strings['time_label']}: {time_str.upper()}",
                  font=info_font, fill=(255, 255, 255))

    # ── ATA logo (top-right of full image) ───────────────────────────────────
    if ata_logo:
        ata = ata_logo.convert("RGBA")
        ata_h = 60
        ata_w = int(ata.width * ata_h / ata.height)
        ata = ata.resize((ata_w, ata_h), Image.LANCZOS)
        bg.paste(ata, (W - ata_w - 16, 16), mask=ata)

    # ── Speaker photos in right card ─────────────────────────────────────────
    speakers = session.get("speakers", [])
    if speakers:
        n = len(speakers)
        available_w = card_w - 40
        photo_size = min(140, (available_w - (n - 1) * 16) // n)
        total_photos_w = n * photo_size + (n - 1) * 16
        photo_start_x = card_x + (card_w - total_photos_w) // 2
        photo_y = card_y + 20

        name_font = _get_font("bold", 13)
        title_font_s = _get_font("regular", 10)
        company_font = _get_font("bold", 11)

        for i, spk in enumerate(speakers):
            px = photo_start_x + i * (photo_size + 16)

            # Square photo
            if spk.get("photo"):
                try:
                    sq = make_square_crop(spk["photo"].convert("RGBA"),
                                          photo_size, photo_size)
                    bg.paste(sq, (px, photo_y), mask=sq)
                except Exception:
                    draw.rectangle([px, photo_y, px + photo_size, photo_y + photo_size],
                                   fill=(200, 200, 200))

            ty = photo_y + photo_size + 6

            # Company logo
            if spk.get("company_logo"):
                try:
                    clogo = spk["company_logo"].convert("RGBA")
                    clogo.thumbnail((photo_size, 28), Image.LANCZOS)
                    lx2 = px + (photo_size - clogo.width) // 2
                    bg.paste(clogo, (lx2, ty), mask=clogo)
                    ty += clogo.height + 4
                except Exception:
                    pass

            # Speaker name
            name = spk.get("name", "")
            nbbox = draw.textbbox((0, 0), name, font=name_font)
            nx = px + (photo_size - nbbox[2]) // 2
            draw.text((nx, ty), name, font=name_font, fill=(50, 50, 50))
            ty += nbbox[3] + 2

            # Job title
            job = spk.get("title", "")
            if len(job) > 30:
                job = job[:28] + "…"
            jbbox = draw.textbbox((0, 0), job, font=title_font_s)
            jx = px + (photo_size - jbbox[2]) // 2
            draw.text((jx, ty), job, font=title_font_s, fill=(100, 100, 100))
            ty += jbbox[3] + 1

            # Company
            company = spk.get("company", "")
            if len(company) > 20:
                company = company[:18] + "…"
            cbbox = draw.textbbox((0, 0), company, font=company_font)
            cx2 = px + (photo_size - cbbox[2]) // 2
            draw.text((cx2, ty), company, font=company_font, fill=theme_rgb)

    # ── Export ────────────────────────────────────────────────────────────────
    final = bg.convert("RGB")
    png_buf = io.BytesIO()
    final.save(png_buf, format="PNG", dpi=(150, 150))
    png_bytes = png_buf.getvalue()
    pptx_bytes = _build_pptx(final, W, H)
    return pptx_bytes, png_bytes


def _build_pptx(img: Image.Image, w: int, h: int) -> bytes:
    from pptx import Presentation
    from pptx.util import Emu
    import io

    DPI = 96
    def px(p): return Emu(int(p / DPI * 914400))

    prs = Presentation()
    prs.slide_width = px(w)
    prs.slide_height = px(h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    img_buf = io.BytesIO()
    img.save(img_buf, format="PNG")
    img_buf.seek(0)
    slide.shapes.add_picture(img_buf, 0, 0, px(w), px(h))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
