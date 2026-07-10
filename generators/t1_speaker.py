"""T1 — Single Speaker Slide  1920 × 1080  (dark dramatic style)

Two-phase rendering:
  Phase 1: background, shapes (→ pptx_canvas, used as PPTX background)
  Phase 2: photo, logo pill and text drawn on top (→ PNG only)
  PPTX:    pptx_canvas + photo as a movable/re-croppable picture shape +
           logo as a picture shape + editable text boxes
"""
import io, math
from PIL import Image, ImageDraw, ImageFilter
from utils.image_utils import make_circular
from utils.gender import moderator_label
from config import LANGUAGE_STRINGS

W, H = 1920, 1080
PHOTO_D  = 480
BORDER_W = 6


def _font(name, size):
    from utils.fonts import get
    return get({"heavy": "montserrat-black",
                "bold":  "montserrat-bold",
                "reg":   "inter"}.get(name, "inter"), size)


def generate(session: dict, theme: dict, language: str = "es",
             renmad_logo: Image.Image | None = None,
             bg_image:    Image.Image | None = None,
             ata_logo:    Image.Image | None = None) -> tuple[bytes, bytes]:

    strings   = LANGUAGE_STRINGS.get(language, LANGUAGE_STRINGS["es"])
    theme_rgb = tuple(theme["rgb"])

    speakers = session.get("speakers", [])
    spk = next((s for s in speakers if not s.get("is_moderator")), speakers[0] if speakers else {})

    # text_draws: [(x, y, text, font, fill)] drawn in Phase 2 (PNG only)
    # text_els:   [{...}] written as PPTX text boxes
    text_draws = []
    text_els   = []
    photo_els  = []
    logo_els   = []

    # ── Background ────────────────────────────────────────────────────────────
    canvas = Image.new("RGB", (W, H), (15, 15, 20))

    if bg_image:
        bg = bg_image.convert("RGB").resize((W, H), Image.LANCZOS)
        bg = bg.filter(ImageFilter.GaussianBlur(radius=3))
        canvas.paste(bg)
        overlay = Image.new("RGB", (W, H), (10, 10, 15))
        canvas.paste(overlay, mask=Image.new("L", (W, H), 160))
    else:
        draw_tmp = ImageDraw.Draw(canvas)
        for y in range(H):
            r = int(15 + math.sin(y / H * math.pi) * 12)
            draw_tmp.line([(0, y), (W, y)], fill=(r, r, r + 5))

    draw = ImageDraw.Draw(canvas)

    # ── Top accent bar ────────────────────────────────────────────────────────
    draw.rectangle([0, 0, W, 8], fill=theme_rgb)

    # ── Left column: photo (Phase 2 for PNG; movable shape in PPTX) ───────────
    col_split  = W // 2   # photo on left half
    cx_photo   = col_split // 2
    cy_photo   = H // 2 - 30

    if spk.get("photo"):
        try:
            photo = make_circular(spk["photo"].convert("RGBA"), PHOTO_D,
                                  border_color=theme_rgb, border_width=BORDER_W)
            sd = PHOTO_D + BORDER_W   # outline centred on edge → outer Ø matches PNG
            photo_els.append({'pil': photo,
                              'px': cx_photo - photo.width // 2,
                              'py': cy_photo - photo.height // 2,
                              'x': cx_photo - sd / 2, 'y': cy_photo - sd / 2,
                              'w': sd, 'h': sd, 'src': spk["photo"]})
        except Exception:
            pass

    # Vertical accent line between columns
    draw.line([(col_split, 60), (col_split, H - 100)], fill=theme_rgb, width=3)

    # ── Right column: text ────────────────────────────────────────────────────
    cx_text = col_split + (W - col_split) // 2
    ty      = H // 2 - 160  # starting y for text block

    # Company logo — grey pill, separate picture shape in PPTX
    if spk.get("company_logo"):
        try:
            clogo = spk["company_logo"].convert("RGBA")
            clogo.thumbnail((240, 64), Image.LANCZOS)
            pad  = 10
            pill = Image.new("RGBA", (clogo.width + pad * 2, clogo.height + pad * 2),
                             (230, 230, 230, 255))
            pill.paste(clogo, (pad, pad), mask=clogo)
            logo_els.append({'pil': pill,
                             'x': cx_text - clogo.width // 2 - pad, 'y': ty - pad,
                             'w': pill.width, 'h': pill.height})
            ty += clogo.height + pad * 2 + 18
        except Exception:
            pass

    # Moderator badge — pill shape in Phase 1, text deferred
    if spk.get("is_moderator"):
        mod_label = moderator_label(spk.get("name", ""), strings)
        mf        = _font("bold", 22)
        mb        = draw.textbbox((0, 0), mod_label, font=mf)
        pad       = 10
        draw.rounded_rectangle(
            [cx_text - mb[2]//2 - pad, ty - 4,
             cx_text + mb[2]//2 + pad, ty + mb[3] + 4],
            radius=6, fill=theme_rgb,
        )
        text_draws.append((cx_text - mb[2] // 2, ty, mod_label, mf, (255, 255, 255)))
        text_els.append({'x': cx_text - mb[2]//2 - 30, 'y': ty,
                         'w': mb[2] + 60, 'h': mb[3] + 10,
                         'text': mod_label, 'size': 22, 'bold': True,
                         'color': (255, 255, 255), 'center': True})
        ty += mb[3] + 20

    # Name
    name = spk.get("name", "")
    nf   = _font("bold", 56)
    nb   = draw.textbbox((0, 0), name, font=nf)
    if name:
        text_draws.append((cx_text - nb[2] // 2, ty, name, nf, (255, 255, 255)))
        text_els.append({'x': cx_text - 400, 'y': ty, 'w': 800, 'h': nb[3] + 16,
                         'text': name, 'size': 56, 'bold': True,
                         'color': (255, 255, 255), 'center': True})
    ty += nb[3] + 8

    # Accent rule under name
    draw.line([(cx_text - 120, ty), (cx_text + 120, ty)], fill=theme_rgb, width=3)
    ty += 16

    # Job title
    role = spk.get("title", "")
    rf   = _font("reg", 34)
    rb   = draw.textbbox((0, 0), role, font=rf)
    if role:
        text_draws.append((cx_text - rb[2] // 2, ty, role, rf, (200, 200, 200)))
        text_els.append({'x': cx_text - 400, 'y': ty, 'w': 800, 'h': rb[3] + 12,
                         'text': role, 'size': 34,
                         'color': (200, 200, 200), 'center': True})
    ty += rb[3] + 8

    # Company
    comp = spk.get("company", "")
    cf   = _font("bold", 32)
    cb   = draw.textbbox((0, 0), comp, font=cf)
    if comp:
        text_draws.append((cx_text - cb[2] // 2, ty, comp, cf, theme_rgb))
        text_els.append({'x': cx_text - 400, 'y': ty, 'w': 800, 'h': cb[3] + 12,
                         'text': comp, 'size': 32, 'bold': True,
                         'color': list(theme_rgb), 'center': True})

    # ── Footer band ───────────────────────────────────────────────────────────
    cta      = session.get("cta_url", "").strip()
    footer_h = 120 if cta else 90
    fy       = H - footer_h

    foot = Image.new("RGBA", (W, footer_h), (*theme_rgb, 230))
    canvas.paste(foot.convert("RGB"), (0, fy))
    draw.line([(0, fy), (W, fy)], fill=(255, 255, 255), width=2)

    # RENMAD logo bottom-right (branding — stays in the background)
    logo_reserve = 0
    if renmad_logo:
        logo = renmad_logo.convert("RGBA")
        lh   = 70
        lw   = int(logo.width * lh / logo.height)
        logo = logo.resize((lw, lh), Image.LANCZOS)
        canvas.paste(logo, (W - lw - 24, fy + (footer_h - lh) // 2), mask=logo)
        logo_reserve = lw + 40

    cx_footer = (W - logo_reserve) // 2

    # Session title
    bar_f      = _font("heavy", 34)
    title_text = session.get("title", "").upper()
    title_y    = fy + 14
    tb         = draw.textbbox((0, 0), title_text, font=bar_f)
    if title_text:
        text_draws.append((cx_footer - tb[2] // 2, title_y, title_text, bar_f,
                           (255, 255, 255)))
        text_els.append({'x': max(0, cx_footer - tb[2] // 2 - 40), 'y': title_y,
                         'w': tb[2] + 80, 'h': tb[3] + 12,
                         'text': title_text, 'size': 34, 'bold': True,
                         'color': (255, 255, 255), 'center': True})

    # Event details: date | time | location
    parts = [p for p in [session.get("date_str",""), session.get("time_str",""),
                          session.get("location","")] if p]
    detail_y = title_y + 42
    if parts:
        detail = "   |   ".join(parts).upper()
        df     = _font("reg", 22)
        db     = draw.textbbox((0, 0), detail, font=df)
        text_draws.append((cx_footer - db[2] // 2, detail_y, detail, df,
                           (230, 230, 230)))
        text_els.append({'x': max(0, cx_footer - db[2] // 2 - 40), 'y': detail_y,
                         'w': db[2] + 80, 'h': db[3] + 10,
                         'text': detail, 'size': 22,
                         'color': (230, 230, 230), 'center': True})
        detail_y += db[3] + 6

    # CTA / URL — translated label + bare URL, hyperlinked in the PPTX
    cta_rect = None
    if cta:
        href     = cta
        label    = strings.get("cta_label", "Regístrate gratis hoy")
        cta_line = f"{label}  →  {href}"
        uf = _font("reg", 20)
        ub = draw.textbbox((0, 0), cta_line, font=uf)
        text_draws.append((cx_footer - ub[2] // 2, detail_y, cta_line, uf,
                           (255, 255, 200)))
        cta_rect = (href, cta_line, cx_footer - ub[2] // 2, detail_y,
                    ub[2] + 60, ub[3] + 10)

    # ── PPTX snapshot: shapes only — no photo, no logo, no session text ──────
    pptx_canvas = canvas.copy()

    # ── Phase 2: photo, logo pill + text for PNG ─────────────────────────────
    for pe in photo_els:
        canvas.paste(pe['pil'], (pe['px'], pe['py']), mask=pe['pil'])
    for le in logo_els:
        canvas.paste(le['pil'], (le['x'], le['y']), mask=le['pil'])
    for x, y, txt, fnt, col in text_draws:
        draw.text((x, y), txt, font=fnt, fill=col)

    png_buf = io.BytesIO()
    canvas.save(png_buf, format="PNG", dpi=(150, 150))
    return (_pptx(pptx_canvas, cta_rect, text_els, logo_els, photo_els, theme_rgb),
            png_buf.getvalue())


def _pptx(img, cta_rect=None, text_els=None, logo_els=None,
          photo_els=None, theme_rgb=None):
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    def px(p): return Emu(int(p / 96 * 914400))
    prs = Presentation()
    prs.slide_width  = px(W); prs.slide_height = px(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    buf = io.BytesIO(); img.save(buf, format="PNG"); buf.seek(0)
    slide.shapes.add_picture(buf, 0, 0, px(W), px(H))

    # Speaker photo — movable picture shape; Crop tool re-frames the face
    from utils.pptx_editable import add_shaped_picture, add_picture_el
    for pe in (photo_els or []):
        try:
            add_shaped_picture(slide, pe['src'], pe['x'], pe['y'],
                               pe['w'], pe['h'], shape="ellipse",
                               border_rgb=theme_rgb, border_px=BORDER_W)
        except Exception:
            pass

    # Company logo pill — separate picture shape
    for le in (logo_els or []):
        try:
            add_picture_el(slide, le['pil'], le['x'], le['y'], le['w'], le['h'])
        except Exception:
            pass

    for el in (text_els or []):
        _w, _h = el.get('w', 400), el.get('h', 60)
        if _w <= 0 or _h <= 0 or not el.get('text'):
            continue
        txBox = slide.shapes.add_textbox(px(el['x']), px(el['y']), px(_w), px(_h))
        txBox.fill.background(); txBox.line.fill.background()
        tf = txBox.text_frame
        tf.word_wrap = el.get('wrap', False)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if el.get('center') else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = el['text']
        # canvas px → pt at 96 dpi, so PPTX text matches the PNG layout
        run.font.size = Pt(el.get('size', 20) * 0.75)
        run.font.bold = el.get('bold', False)
        run.font.name = "Montserrat" if el.get('bold') else "Inter"
        c = el.get('color', (255, 255, 255))
        run.font.color.rgb = RGBColor(int(c[0]), int(c[1]), int(c[2]))

    # CTA hyperlink text box
    if cta_rect:
        url, line, cx, cy, cw, ch = cta_rect
        txBox = slide.shapes.add_textbox(px(cx), px(cy), px(cw), px(ch))
        txBox.fill.background(); txBox.line.fill.background()
        tf = txBox.text_frame; tf.word_wrap = False
        run = tf.paragraphs[0].add_run()
        run.text = line
        run.font.size = Pt(15); run.font.name = "Inter"
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xC8)
        try:
            rPr = run._r.get_or_add_rPr()
            hl  = etree.SubElement(rPr, qn('a:hlinkClick'))
            rId = slide.part.relate_to(
                url,
                'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
                is_external=True)
            hl.set(qn('r:id'), rId)
        except Exception:
            pass

    out = io.BytesIO(); prs.save(out)
    return out.getvalue()
