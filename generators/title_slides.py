# -*- coding: utf-8 -*-
"""
Title-slide generator for RENMAD/ATA events.

Builds **native, editable** PowerPoint slides (theme band + text boxes + pictures)
from a canonical agenda dict (see utils/parse_agenda.py). Native shapes mean a
last-minute name/title fix is editable directly in PowerPoint and survives a
Canva import. A companion exporter (utils/pptx_render.py) flattens the same .pptx
to PNGs for LinkedIn assets, so there is a single source of truth.

This module currently implements the SINGLE-SPEAKER card (vertical slice).
Panel / cover / section / utility renderers are added in a later step.
"""
from __future__ import annotations

import io
import os
import re
import tempfile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls

from config import THEMES
from utils import svg_icons

# ── Geometry (16:9 widescreen) ────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
EMU_PER_IN = 914400

# ── Brand fonts (installed user-scope; PowerPoint renders them directly) ──────
F_BLACK = "Montserrat Black"
F_BOLD = "Montserrat SemiBold"
F_REG = "Montserrat"

CHARCOAL = RGBColor(0x1C, 0x25, 0x29)
MIDGREY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

_HERE = os.path.dirname(__file__)
_ASSETS_LOGOS = os.path.normpath(os.path.join(_HERE, "..", "assets", "logos"))

_TMP = tempfile.mkdtemp(prefix="renmad_title_")


# ── Helpers ───────────────────────────────────────────────────────────────────

def _theme(theme_key: str) -> dict:
    return THEMES.get(theme_key, THEMES["datacenters"])


def _theme_rgb(theme_key: str) -> RGBColor:
    r, g, b = _theme(theme_key)["rgb"]
    return RGBColor(r, g, b)


def _initials(name: str) -> str:
    parts = [p for p in (name or "").split() if p[:1].isalpha()]
    if not parts:
        return "?"
    if len(parts) == 1:
        return parts[0][:2].upper()
    return (parts[0][0] + parts[-1][0]).upper()


def _circle_png(photo: Image.Image | None, name: str, theme_key: str,
                px: int = 720) -> str:
    """
    Circular avatar PNG with a thin near-black ring and a soft drop shadow
    (matches the E044 reference). If `photo` is None, draw a neutral grey disc
    with the speaker's initials so the layout can be reviewed before real
    photos are assigned. Returns a temp file path.
    """
    from PIL import ImageFilter, ImageOps

    pad = int(px * 0.06)                 # room for the shadow
    full = px + pad * 2
    canvas = Image.new("RGBA", (full, full), (0, 0, 0, 0))

    # ── soft drop shadow ─────────────────────────────────────────────────────
    shadow = Image.new("RGBA", (full, full), (0, 0, 0, 0))
    sd = ImageDraw.Draw(shadow)
    sd.ellipse([pad, pad + int(pad * 0.5), pad + px, pad + px + int(pad * 0.5)],
               fill=(0, 0, 0, 70))
    shadow = shadow.filter(ImageFilter.GaussianBlur(int(pad * 0.7)))
    canvas.alpha_composite(shadow)

    # ── disc (photo or initials placeholder) ─────────────────────────────────
    disc = Image.new("RGBA", (px, px), (0, 0, 0, 0))
    mask = Image.new("L", (px, px), 0)
    ImageDraw.Draw(mask).ellipse([0, 0, px - 1, px - 1], fill=255)
    if photo is not None:
        # Top-biased crop so the crop keeps heads/foreheads in frame instead of
        # cutting them off (centering cy=0.5 cropped equally top+bottom → "corta
        # la cabeza"; 0.32 keeps more of the top).
        fitted = ImageOps.fit(photo.convert("RGB"), (px, px),
                              method=Image.LANCZOS, centering=(0.5, 0.32))
        disc.paste(fitted, (0, 0), mask)
    else:
        base = Image.new("RGB", (px, px), (0xE9, 0xEC, 0xEF))
        disc.paste(base, (0, 0), mask)
        d = ImageDraw.Draw(disc)
        try:
            font = ImageFont.truetype(
                os.path.join(_HERE, "..", "assets", "fonts", "Montserrat-Bold.ttf"),
                int(px * 0.32))
        except Exception:
            font = ImageFont.load_default()
        txt = _initials(name)
        bb = d.textbbox((0, 0), txt, font=font)
        d.text(((px - (bb[2] - bb[0])) / 2 - bb[0], (px - (bb[3] - bb[1])) / 2 - bb[1]),
               txt, font=font, fill=(0x9A, 0xA3, 0xAB))
    canvas.alpha_composite(disc, (pad, pad))

    # ── thin near-black ring ─────────────────────────────────────────────────
    ring = max(2, int(px * 0.006))
    ImageDraw.Draw(canvas).ellipse(
        [pad + ring, pad + ring, pad + px - 1 - ring, pad + px - 1 - ring],
        outline=(0x2A, 0x2E, 0x31, 255), width=ring)

    path = os.path.join(_TMP, f"circ_{abs(hash((name, theme_key, id(photo)))) % 10**9}.png")
    canvas.save(path, "PNG")
    return path


def _square_top(photo: Image.Image, bias: float = 0.2) -> Image.Image:
    """Crop a photo to a square, biased toward the top so heads stay in frame.
    This is only the INITIAL framing — in PowerPoint the user can re-crop freely."""
    w, h = photo.size
    s = min(w, h)
    left = (w - s) // 2
    top = int((h - s) * bias)
    return photo.convert("RGB").crop((left, top, left + s, top + s))


def _skin_centroid(img: Image.Image) -> "tuple[float, float] | None":
    """Rough (x_frac, y_frac) centre of skin-coloured pixels — a dependency-free
    stand-in for face detection so an off-centre subject isn't cropped out. Works
    on a downscaled copy for speed. Returns None if too little skin is found."""
    small = img.convert("RGB")
    small.thumbnail((160, 160))
    sw, sh = small.size
    px = small.load()
    sx = sy = n = 0
    for y in range(sh):
        for x in range(sw):
            r, g, b = px[x, y]
            mx, mn = max(r, g, b), min(r, g, b)
            if (r > 95 and g > 40 and b > 20 and r > g and r > b
                    and (mx - mn) > 15 and abs(r - g) > 10):
                sx += x; sy += y; n += 1
    if n < sw * sh * 0.02:          # < 2% skin → unreliable, let caller fall back
        return None
    return (sx / n) / sw, (sy / n) / sh


def _smart_square(photo: Image.Image) -> Image.Image:
    """Crop a photo to a square centred on the speaker's face. Uses a skin-tone
    centroid to recover off-centre subjects (e.g. a person standing on the left of
    a wide shot); falls back to a top-biased centre crop when unsure. Only the
    initial framing — the user can still re-crop freely in PowerPoint."""
    img = photo.convert("RGB")
    w, h = img.size
    s = min(w, h)
    if w == h:
        return img
    cen = _skin_centroid(img)
    if w > h:                       # landscape → choose the horizontal window
        cx = cen[0] if cen else 0.5
        left = int(round(cx * w - s / 2))
        left = max(0, min(w - s, left))
        return img.crop((left, 0, left + s, s))
    # portrait → choose the vertical window, biased so the face sits a bit high
    cy = cen[1] if cen else 0.34
    top = int(round(cy * h - s * 0.45))
    top = max(0, min(h - s, top))
    return img.crop((0, top, s, top + s))


def _crop_to_ellipse(pic):
    """Set a native picture's geometry to an ellipse → circular crop the user can
    still pan/zoom inside via PowerPoint's Crop tool. Exactly ONE prstGeom (two
    would corrupt the file); spPr here holds only a:xfrm, so appending is valid."""
    spPr = pic._element.spPr
    geom = spPr.find(qn("a:prstGeom"))
    if geom is None:
        geom = spPr.makeelement(qn("a:prstGeom"), {})
        geom.append(spPr.makeelement(qn("a:avLst"), {}))
        spPr.append(geom)
    geom.set("prst", "ellipse")


def _soft_shadow(shape, blur_in=0.10, dist_in=0.05, alpha=32):
    """Attach a subtle outer drop shadow to a shape/picture via the DrawingML
    effect list (keeps the E044 look now that the shadow isn't baked in)."""
    spPr = shape._element.spPr
    # reuse an existing (possibly empty) a:effectLst — e.g. the one created by
    # shape.shadow.inherit = False — so we never end up with TWO (→ corrupt file).
    eff = spPr.find(qn("a:effectLst"))
    if eff is None:
        eff = spPr.makeelement(qn("a:effectLst"), {})
        spPr.append(eff)
    else:
        for ch in list(eff):
            eff.remove(ch)
    sh = spPr.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(blur_in * EMU_PER_IN)),
        "dist":    str(int(dist_in * EMU_PER_IN)),
        "dir":     "5400000",          # straight down
        "rotWithShape": "0",
    })
    clr = spPr.makeelement(qn("a:srgbClr"), {"val": "000000"})
    clr.append(spPr.makeelement(qn("a:alpha"), {"val": str(alpha * 1000)}))
    sh.append(clr); eff.append(sh)


def _lock_shape(shape, no_move=False, no_resize=True, no_aspect=True):
    """Add a:spLocks so the circle FRAME can't be accidentally resized/distorted
    (the image inside is still freely adjustable via the picture fill)."""
    sp = shape._element
    nvSpPr = sp.find(qn("p:nvSpPr"))
    cNvSpPr = nvSpPr.find(qn("p:cNvSpPr"))
    locks = cNvSpPr.find(qn("a:spLocks"))
    if locks is None:
        locks = cNvSpPr.makeelement(qn("a:spLocks"), {})
        cNvSpPr.append(locks)
    if no_move:
        locks.set("noMove", "1")
    if no_resize:
        locks.set("noResize", "1")
    if no_aspect:
        locks.set("noChangeAspect", "1")


def _fill_shape_with_photo(shape, photo):
    """Replace a shape's fill with the speaker photo as a PICTURE FILL, clipped to
    the shape. The shape is the fixed window; the image is the fill the user
    pans/zooms via Format Picture → Fill (Offsets), always clipped to the circle."""
    bio = io.BytesIO()
    _square_top(photo).save(bio, "PNG")
    bio.seek(0)
    image_part, rId = shape.part.get_or_add_image_part(bio)

    spPr = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill",
                "a:pattFill", "a:grpFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    blipFill = spPr.makeelement(qn("a:blipFill"), {})
    blip = spPr.makeelement(qn("a:blip"), {qn("r:embed"): rId})
    blipFill.append(blip)
    stretch = spPr.makeelement(qn("a:stretch"), {})
    stretch.append(spPr.makeelement(qn("a:fillRect"), {}))
    blipFill.append(stretch)
    # fill must sit after a:prstGeom and before a:ln
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(blipFill)
    else:
        spPr.append(blipFill)


def _add_circle_window(slide, fx, fy, fw, fh, hcx, hcy, r, fill_rgb=WHITE):
    """A white rectangle with a CIRCULAR HOLE punched out (custom geometry), used
    as a mask: placed on top of a freely-editable photo so only the circle shows.
    Outer rect is wound clockwise, the hole counter-clockwise → the circle is
    empty (lets the photo beneath show through). Coords in EMU; hole centre
    (hcx, hcy) is local to the rectangle."""
    fx, fy = max(0, int(fx)), max(0, int(fy))
    fw, fh, r = int(fw), int(fh), int(r)
    hcx, hcy = int(hcx), int(hcy)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, fx, fy, fw, fh)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background(); shape.shadow.inherit = False
    xml = (
        f'<a:custGeom {nsdecls("a")}>'
        '<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="0" t="0" r="{fw}" b="{fh}"/>'
        '<a:pathLst>'
        f'<a:path w="{fw}" h="{fh}">'
        '<a:moveTo><a:pt x="0" y="0"/></a:moveTo>'
        f'<a:lnTo><a:pt x="{fw}" y="0"/></a:lnTo>'
        f'<a:lnTo><a:pt x="{fw}" y="{fh}"/></a:lnTo>'
        f'<a:lnTo><a:pt x="0" y="{fh}"/></a:lnTo>'
        '<a:close/>'
        f'<a:moveTo><a:pt x="{hcx + r}" y="{hcy}"/></a:moveTo>'
        f'<a:arcTo wR="{r}" hR="{r}" stAng="0" swAng="-5400000"/>'
        f'<a:arcTo wR="{r}" hR="{r}" stAng="-5400000" swAng="-5400000"/>'
        f'<a:arcTo wR="{r}" hR="{r}" stAng="-10800000" swAng="-5400000"/>'
        f'<a:arcTo wR="{r}" hR="{r}" stAng="-16200000" swAng="-5400000"/>'
        '<a:close/>'
        '</a:path></a:pathLst></a:custGeom>'
    )
    spPr = shape._element.spPr
    spPr.replace(spPr.find(qn("a:prstGeom")), parse_xml(xml))
    return shape


def _add_circle_photo(slide, photo, cx_emu, top_emu, dia_in, theme_key, name,
                      frame=None):
    """
    Place a speaker's photo so the picture stays FREELY editable (normal handles)
    while only a fixed circular window of it shows:
      • real photo → 3 stacked layers:
          1. the photo as an ordinary, movable/resizable picture (back),
          2. a white mask = rectangle with a circular HOLE on top of it (hides
             everything outside the circle — invisible on the white slide),
          3. a thin ring outline on top.
        Resize/move the photo with the normal handles; the circle never changes.
      • no photo   → the baked grey initials placeholder (nothing to edit).
    `frame` = (x, y, w, h) EMU rect the white mask covers (defaults to a generous
    box around the circle). The hole is centred on the photo's circle.
    """
    dia = Inches(dia_in)
    x = int(cx_emu - dia / 2)
    top = int(top_emu)
    if photo is None:
        circ = _circle_png(None, name, theme_key)
        return _add_picture_fit(slide, circ, cx_emu, top_emu, dia_in, dia_in)

    # 1) photo as a normal editable picture (initially fills the circle),
    #    face-centred so off-centre subjects aren't cropped out of the circle
    bio = io.BytesIO()
    _smart_square(photo).save(bio, "PNG")
    bio.seek(0)
    pic = slide.shapes.add_picture(bio, x, top, width=dia, height=dia)

    # 2) white mask with a circular hole over the photo
    if frame is None:
        pad = int(Inches(0.45))
        frame = (x - pad, top - pad, int(dia) + 2 * pad, int(dia) + 2 * pad)
    fx, fy, fw, fh = (int(v) for v in frame)
    cy = top + int(dia) // 2
    _add_circle_window(slide, fx, fy, fw, fh,
                       hcx=cx_emu - fx, hcy=cy - fy, r=int(dia) // 2)

    # 3) thin ring on top (+ subtle shadow on the outline)
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, top, dia, dia)
    ring.fill.background()
    ring.line.color.rgb = RGBColor(0x2A, 0x2E, 0x31)
    ring.line.width = Pt(1.25)
    ring.shadow.inherit = False
    _soft_shadow(ring)
    return pic


def _no_line(shape):
    shape.line.fill.background()


def _add_band(slide, theme_key: str, height: Emu, top: Emu = Emu(0)):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SLIDE_W, height)
    band.fill.solid()
    band.fill.fore_color.rgb = _theme_rgb(theme_key)
    _no_line(band)
    band.shadow.inherit = False
    return band


def _add_divider(slide, cx: Emu, cy: Emu, width_in=2.6, rgb=WHITE):
    """Thin horizontal rule centred on (cx, cy) — separates the two title languages."""
    w = Inches(width_in)
    h = Emu(int(0.022 * EMU_PER_IN))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 int(cx - w / 2), int(cy - h / 2), w, h)
    bar.fill.solid(); bar.fill.fore_color.rgb = rgb; _no_line(bar)
    bar.shadow.inherit = False
    return bar


# ── Title measurement (PIL-measured so nothing overflows the band) ───────────
_PX_PER_IN = 144.0   # 1920 px / 13.333 in
_LINE_FACTOR = 1.24  # ~PowerPoint single-line height for Montserrat Black (a bit
                     # generous so the band never ends up shorter than the text)
# Width used to MEASURE wrapping — deliberately narrower than the rendered text
# box (≈11.37in usable) so PIL never under-counts lines vs PowerPoint's slightly
# wider real font metrics (→ the band is sized for ≥ the lines PPT will render).
_BAND_MEASURE_W = 10.7
_FONT_FILES = {
    F_BLACK: "Montserrat-Black.ttf",
    F_BOLD:  "Montserrat-SemiBold.ttf",
    F_REG:   "Montserrat-Regular.ttf",
}


def _pil_font(font_name, size_pt):
    path = os.path.join(_HERE, "..", "assets", "fonts",
                        _FONT_FILES.get(font_name, _FONT_FILES[F_REG]))
    return ImageFont.truetype(path, max(1, int(round(size_pt * 2))))  # pt→px @144dpi


def _wrap_px(text, font, max_w_px):
    words, lines, cur = (text or "").split(), [], ""
    for w in words:
        t = (cur + " " + w).strip()
        if font.getlength(t) > max_w_px and cur:
            lines.append(cur); cur = w
        else:
            cur = t
    if cur:
        lines.append(cur)
    return lines or [""]


def _measure(text, font_name, size_pt, max_w_px):
    f = _pil_font(font_name, size_pt)
    lines = _wrap_px(text, f, max_w_px)
    line_h_in = (size_pt * 2 * _LINE_FACTOR) / _PX_PER_IN
    longest = max((f.getlength(l) for l in lines), default=0) / _PX_PER_IN
    return lines, len(lines) * line_h_in, longest


def _fit_size(text, font_name, max_w_in, max_h_in, hi, lo, max_lines=3):
    max_w_px = max_w_in * _PX_PER_IN
    for s in range(hi, lo - 1, -1):
        lines, h_in, _ = _measure(text, font_name, s, max_w_px)
        if h_in <= max_h_in and len(lines) <= max_lines:
            return s, lines
    lines, _, _ = _measure(text, font_name, lo, max_w_px)
    return lo, lines


def _band_pos(layout, band_pos):
    """Resolve where the title band sits. Marketing is always top; event defaults
    to bottom but the builder can choose. `band_pos` ('top'/'bottom') wins."""
    if band_pos in ("top", "bottom"):
        return band_pos
    return "top" if layout == "marketing" else "bottom"


def _add_title_band(slide, theme_key, session, layout="marketing",
                    title_fit="uniform", band_pos=None, top_offset=Emu(0)):
    """
    Coloured title band with the bilingual session title.
      band_pos  'top' | 'bottom' (defaults from layout: marketing→top, event→bottom)
      title_fit 'uniform'   → fixed band height, title auto-shrinks to fit
                'flexible'  → title kept large, band grows/shrinks to fit it
    Returns (content_top_emu, content_bottom_emu) — the free area for speakers.
    """
    bp = _band_pos(layout, band_pos)
    title1 = (session.get("title") or "").strip()
    title2 = (session.get("title_2") or "").strip()
    max_w_in = _BAND_MEASURE_W
    pad_in = 0.24

    if title_fit == "uniform":
        band_h_in = 1.95 if bp == "top" else 1.6
        avail = (band_h_in - 2 * pad_in)
        if title2:
            en_pt, en_lines = _fit_size(title1, F_BLACK, max_w_in, avail * 0.56,
                                        hi=32, lo=13)
            it_pt, it_lines = _fit_size(title2, F_BOLD, max_w_in, avail * 0.40,
                                        hi=max(13, en_pt - 4), lo=11)
        else:
            en_pt, en_lines = _fit_size(title1, F_BLACK, max_w_in, avail,
                                        hi=34, lo=15)
            it_pt = it_lines = None
    else:  # flexible — keep the title big, grow the band to fit it (shrinking the
           # font only if it would otherwise overflow the band's max height)
        BAND_MAX = 3.0
        en_pt = 30 if bp == "top" else 28
        while True:
            en_lines = _measure(title1, F_BLACK, en_pt, max_w_in * _PX_PER_IN)[0]
            it_pt = max(15, int(en_pt * 0.72)) if title2 else None
            it_lines = (_measure(title2, F_BOLD, it_pt, max_w_in * _PX_PER_IN)[0]
                        if title2 else None)
            en_h = len(en_lines) * (en_pt * 2 * _LINE_FACTOR) / _PX_PER_IN
            it_h = (len(it_lines) * (it_pt * 2 * _LINE_FACTOR) / _PX_PER_IN) if title2 else 0
            div_in = 0.20 if title2 else 0
            band_h_in = 2 * pad_in + en_h + div_in + it_h
            if (len(en_lines) <= 3 and band_h_in <= BAND_MAX) or en_pt <= 16:
                break
            en_pt -= 2
        band_h_in = max(1.2, min(BAND_MAX, band_h_in))

    band_h = Inches(band_h_in)
    if bp == "bottom":
        band_top = SLIDE_H - band_h
        content_top, content_bottom = Emu(0), band_top
    else:
        band_top = top_offset
        content_top, content_bottom = band_top + band_h, SLIDE_H
    _add_band(slide, theme_key, band_h, top=band_top)

    # ── place the measured text blocks, vertically centred in the band ────────
    en_h = len(en_lines) * (en_pt * 2 * _LINE_FACTOR) / _PX_PER_IN
    it_h = (len(it_lines) * (it_pt * 2 * _LINE_FACTOR) / _PX_PER_IN) if title2 else 0
    div_in = 0.18 if title2 else 0
    total_in = en_h + div_in + it_h
    start_in = max(pad_in, (band_h_in - total_in) / 2)
    y = band_top + Inches(start_in)

    en_runs = [[(l, F_BLACK, en_pt, WHITE, True)] for l in en_lines]
    _add_text(slide, Inches(0.93), y, Inches(11.47), Inches(en_h + 0.05),
              en_runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    if title2:
        _, _, longest_en = _measure(title1, F_BLACK, en_pt, max_w_in * _PX_PER_IN)
        div_w = min(4.0, max(1.4, longest_en * 0.45))
        div_cy = y + Inches(en_h + div_in / 2)
        _add_divider(slide, SLIDE_W // 2, div_cy, width_in=div_w)
        it_runs = [[(l, F_BOLD, it_pt, RGBColor(0xF2, 0xF2, 0xF2), False)] for l in it_lines]
        _add_text(slide, Inches(0.93), y + Inches(en_h + div_in),
                  Inches(11.47), Inches(it_h + 0.05),
                  it_runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    return content_top, content_bottom


def _add_text(slide, x, y, w, h, runs, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    """
    runs: list of paragraphs; each paragraph is list of (text, font, size_pt, rgb, bold).
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.0
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        for (text, font, size, rgb, bold) in para:
            r = p.add_run()
            r.text = text
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = rgb
            r.font.bold = bold
    return tb


def _add_picture_fit(slide, path_or_img, cx_emu, top_emu, max_w_in, max_h_in,
                     bottom_emu=None):
    """Add a picture centred horizontally on cx_emu, fit within max_w x max_h.
    If bottom_emu is given, the picture is bottom-aligned to that y (used to seat
    logos of differing heights on a common baseline)."""
    if isinstance(path_or_img, Image.Image):
        bio = io.BytesIO()
        path_or_img.convert("RGBA").save(bio, "PNG")
        bio.seek(0)
        src = bio
        w0, h0 = path_or_img.size
    else:
        src = path_or_img
        with Image.open(path_or_img) as im:
            w0, h0 = im.size
    scale = min(max_w_in * EMU_PER_IN / w0, max_h_in * EMU_PER_IN / h0)
    w = int(w0 * scale)
    h = int(h0 * scale)
    y = int(bottom_emu - h) if bottom_emu is not None else int(top_emu)
    return slide.shapes.add_picture(src, int(cx_emu - w / 2), y, width=w, height=h)


# ── Single-speaker card ───────────────────────────────────────────────────────

# ── Shared slide chrome ───────────────────────────────────────────────────────

def _white_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; _no_line(bg); bg.shadow.inherit = False


def _theme_logo_corner(slide, theme_key, layout="marketing", w_in=1.6, band_pos=None):
    """Small theme logo in the corner OPPOSITE the title band, so neither the band
    nor the speaker photos collide with it (band bottom → logo top-right; band
    top → logo bottom-right)."""
    logo_fn = _theme(theme_key).get("logo_filename")
    if not logo_fn:
        return
    lp = os.path.join(_ASSETS_LOGOS, logo_fn)
    if not os.path.exists(lp):
        return
    with Image.open(lp) as im:
        w0, h0 = im.size
    tw = Inches(w_in)
    th = int(tw * h0 / w0)
    x = SLIDE_W - tw - Inches(0.3)
    y = Inches(0.3) if _band_pos(layout, band_pos) == "bottom" \
        else SLIDE_H - Emu(th) - Inches(0.25)
    slide.shapes.add_picture(lp, x, y, width=tw, height=Emu(th))


# ── Optional top space for hand-added sponsor (SPX) logos ─────────────────────
# Small events have no LED background carrying the sponsor logos, so the slide
# itself needs a blank strip at the top to paste them into, with a slim theme
# band beneath (the "SPX types" line, left blank). Reserving it pushes the
# speaker block down. See Belén's reference slides 2026-06-16.
_SPX_LOGO_H = Inches(1.4)   # blank space at the very top for logos (by hand)
_SPX_BAND_H = Inches(0.3)   # slim theme-colour band under it (blank)
_SPX_RESERVE = _SPX_LOGO_H + _SPX_BAND_H


def _add_spx_strip(slide, theme_key):
    """Draw the slim theme-colour band beneath the (blank) hand-logo space.
    Returns the reserved bottom y (EMU) — content must start below this."""
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, _SPX_LOGO_H, SLIDE_W, _SPX_BAND_H)
    band.fill.solid(); band.fill.fore_color.rgb = _theme_rgb(theme_key)
    _no_line(band); band.shadow.inherit = False
    return _SPX_RESERVE


# ── Single-speaker card ───────────────────────────────────────────────────────

def add_single_speaker_slide(prs: Presentation, theme_key: str, session: dict,
                             speaker: dict, photo: Image.Image | None = None,
                             logo: Image.Image | None = None,
                             lang_strings: dict | None = None,
                             layout: str = "marketing", title_fit: str = "flexible",
                             band_pos: str | None = None, spx_space: bool = False):
    """
    One title card for an individual speaker: circular photo, company logo,
    name (bold) + role, optional Moderator label, plus the bilingual session
    title in a coloured band (top or bottom — see `band_pos`).
    `spx_space` reserves a blank top strip for hand-added sponsor logos.
    """
    bp = _band_pos(layout, band_pos)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _white_bg(slide)

    _top_off = _SPX_RESERVE if (spx_space and bp == "top") else Emu(0)
    content_top, content_bottom = _add_title_band(
        slide, theme_key, session, title_fit=title_fit, band_pos=bp, top_offset=_top_off)
    if spx_space:
        _add_spx_strip(slide, theme_key)
        content_top = max(int(content_top), int(_SPX_RESERVE))

    cx = SLIDE_W // 2
    ls = lang_strings or {}

    # Circle as big as the free area allows (the photo is centred horizontally, so
    # it never reaches the corner logo — no clearance needed here).
    below_in = 1.95                          # logo zone + name block under the circle
    content_h = content_bottom - content_top
    avail_in = content_h / EMU_PER_IN
    dia_in = max(1.5, min(3.1, avail_in - below_in - 0.3))
    block_h = Inches(dia_in + below_in)
    y = content_top + max(Inches(0.1), Emu(int((content_h - block_h) / 2)))
    _fy = max(int(content_top), int(y - Inches(0.4)))
    _frame = (int(Inches(2.0)), _fy, int(Inches(9.33)),
              int(y + Inches(dia_in) + Inches(0.1)) - _fy)
    _add_circle_photo(slide, photo, cx, y, dia_in, theme_key,
                      speaker.get("name", ""), frame=_frame)
    y = y + Inches(dia_in) + Inches(0.18)

    # company logo in a fixed-height zone (so the name sits at a stable height)
    logo_zone_in = 0.95
    if logo is not None:
        _add_picture_fit(slide, logo, cx, y, 2.6, logo_zone_in,
                         bottom_emu=y + Inches(logo_zone_in))
        y = y + Inches(logo_zone_in + 0.14)
    else:
        company = (speaker.get("company") or "").strip()
        if company:
            _add_text(slide, Inches(1.5), y, Inches(10.333), Inches(logo_zone_in),
                      [[(company.upper(), F_BOLD, 18, CHARCOAL, True)]],
                      anchor=MSO_ANCHOR.BOTTOM)
            y = y + Inches(logo_zone_in + 0.14)

    # name / role / (moderator) stacked — moderator last so it never shifts name
    name = (speaker.get("name") or "").strip()
    role = (speaker.get("role") or "").strip()
    name_paras = [[(name, F_BLACK, 24, CHARCOAL, True)]]
    if role:
        name_paras.append([(role, F_REG, 15, MIDGREY, False)])
    if speaker.get("is_moderator"):
        name_paras.append([(ls.get("moderator", "Moderator").upper(), F_BOLD, 13,
                            _theme_rgb(theme_key), True)])
    _add_text(slide, Inches(1.0), y, Inches(11.333), Inches(1.0), name_paras)

    if not spx_space:
        _theme_logo_corner(slide, theme_key, band_pos=bp)
    return slide


# ── Combined panel slide (row of speakers) ────────────────────────────────────

def add_panel_slide(prs: Presentation, theme_key: str, session: dict,
                    photos: dict | None = None, logos: dict | None = None,
                    lang_strings: dict | None = None, layout: str = "marketing",
                    title_fit: str = "flexible", band_pos: str | None = None,
                    spx_space: bool = False):
    """One slide for a whole panel: bilingual title band + a row of speakers
    (photo / company logo / name / role / Moderator chip). `band_pos` puts the
    title band 'top' or 'bottom'. `spx_space` reserves a blank top strip for
    hand-added sponsor logos."""
    photos = photos or {}
    logos = logos or {}
    ls = lang_strings or {}
    bp = _band_pos(layout, band_pos)
    speakers = [s for s in session.get("speakers", []) if (s.get("name") or "").strip()]
    if not speakers:
        return None

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _white_bg(slide)

    _top_off = _SPX_RESERVE if (spx_space and bp == "top") else Emu(0)
    content_top, content_bottom = _add_title_band(
        slide, theme_key, session, title_fit=title_fit, band_pos=bp, top_offset=_top_off)
    if spx_space:
        _add_spx_strip(slide, theme_key)
        content_top = max(int(content_top), int(_SPX_RESERVE))

    n = len(speakers)
    # circle diameter shrinks as the panel grows (bumped up — "bigger where possible")
    dia_in = {1: 3.0, 2: 2.85, 3: 2.5, 4: 2.1, 5: 1.78}.get(n, 1.6)
    name_pt = {1: 22, 2: 20, 3: 18, 4: 16, 5: 14}.get(n, 13)
    role_pt = max(10, name_pt - 5)
    comp_pt = max(11, name_pt - 4)

    margin = Inches(0.55)
    usable = SLIDE_W - 2 * margin
    col_w = usable / n
    below_in = 0.14 + 0.82 + 1.4             # logo zone + name block under each circle
    logo_zone_in = 0.82

    # Inset the photo block away from the corner logo (which sits opposite the
    # band) so a circle never overlaps it: band bottom → logo top, clear the top;
    # band top → logo bottom, clear the bottom.
    clear = Emu(0) if spx_space else Inches(1.1)   # no corner logo when SPX space is on
    ct = content_top + (clear if bp == "bottom" else Emu(0))
    cb = content_bottom - (clear if bp == "top" else Emu(0))
    ch = cb - ct
    block_h = Inches(dia_in + below_in)
    photo_top = ct + max(Inches(0.2), Emu(int((ch - block_h) / 2)))

    # Cap the circle so circle + text still fit within the usable region.
    avail_below_in = (cb - photo_top) / EMU_PER_IN - below_in - 0.1
    if dia_in > avail_below_in:
        dia_in = max(1.4, avail_below_in)
    max_logo_w = min(col_w / EMU_PER_IN - 0.25, 1.8)

    # ── fixed vertical zones so every column's logo / name / role line up ─────
    photo_bottom  = photo_top + Inches(dia_in)
    logo_zone_top = photo_bottom + Inches(0.14)
    logo_zone_bot = logo_zone_top + Inches(logo_zone_in)
    name_y        = logo_zone_bot + Inches(0.06)
    max_logo_w    = min(col_w / EMU_PER_IN - 0.25, 1.8)

    for i, sp in enumerate(speakers):
        col_cx = margin + col_w * i + col_w // 2
        col_left = Emu(int(col_cx - col_w // 2))
        _cfy = max(int(content_top), int(photo_top - Inches(0.35)))
        _cframe = (int(col_left), _cfy, int(col_w),
                   int(logo_zone_top - Inches(0.05)) - _cfy)
        _add_circle_photo(slide, photos.get(sp.get("name", "")), col_cx, photo_top,
                          dia_in, theme_key, sp.get("name", ""), frame=_cframe)

        # company logo (bottom-aligned in the fixed zone) or company name text
        logo = logos.get((sp.get("company") or "").strip())
        if logo is not None:
            _add_picture_fit(slide, logo, col_cx, logo_zone_top,
                             max_logo_w, logo_zone_in, bottom_emu=logo_zone_bot)
        else:
            company = (sp.get("company") or "").strip()
            if company:
                _add_text(slide, col_left, logo_zone_top, col_w, Inches(logo_zone_in),
                          [[(company.upper(), F_BOLD, comp_pt, CHARCOAL, True)]],
                          anchor=MSO_ANCHOR.BOTTOM)

        # name / role / (moderator) — all start at the SAME name_y across columns;
        # the moderator label is the LAST line so it never shifts the name.
        paras = [[(sp.get("name", ""), F_BLACK, name_pt, CHARCOAL, True)]]
        role = (sp.get("role") or "").strip()
        if role:
            paras.append([(role, F_REG, role_pt, MIDGREY, False)])
        if sp.get("is_moderator"):
            paras.append([(ls.get("moderator", "Moderator").upper(), F_BOLD,
                           max(10, role_pt), _theme_rgb(theme_key), True)])
        _add_text(slide, col_left, name_y, col_w, Inches(1.5), paras,
                  anchor=MSO_ANCHOR.TOP)

    if not spx_space:
        _theme_logo_corner(slide, theme_key, band_pos=bp)
    return slide


# ── Cover / welcome slide ─────────────────────────────────────────────────────

def add_cover_slide(prs: Presentation, theme_key: str, event: dict,
                    logo_wall: "Image.Image | str | None" = None,
                    lang_strings: dict | None = None):
    """Opening slide: RENMAD theme logo + Welcome kicker + event title +
    date·venue + the sponsor/speaker logo wall along the bottom. All text native."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    ls = lang_strings or {}
    cx = SLIDE_W // 2
    y = Inches(0.5)

    # RENMAD theme logo, centred
    logo_fn = _theme(theme_key).get("logo_filename")
    lp = os.path.join(_ASSETS_LOGOS, logo_fn) if logo_fn else None
    if lp and os.path.exists(lp):
        with Image.open(lp) as im:
            w0, h0 = im.size
        tw = Inches(3.0); th = int(tw * h0 / w0)
        slide.shapes.add_picture(lp, int(cx - tw / 2), int(y), width=tw, height=Emu(th))
        y = y + Emu(th) + Inches(0.3)

    # Welcome kicker
    welcome = ls.get("welcome", "Welcome").upper()
    _add_text(slide, Inches(1.0), y, Inches(11.33), Inches(0.5),
              [[(welcome, F_BOLD, 22, _theme_rgb(theme_key), True)]],
              anchor=MSO_ANCHOR.TOP)
    y = y + Inches(0.6)

    # Sponsor / speaker logo wall along the bottom — placed first so we know its
    # top edge and can seat the date·venue line just above it.
    wall_top = SLIDE_H - Inches(0.4)
    if logo_wall is not None:
        src = logo_wall
        if isinstance(logo_wall, Image.Image):
            bio = io.BytesIO(); logo_wall.save(bio, "PNG"); bio.seek(0)
            src = bio; w0, h0 = logo_wall.size
        else:
            with Image.open(logo_wall) as im:
                w0, h0 = im.size
        max_w, max_h = Inches(12.0), Inches(2.1)
        scale = min(int(max_w) / w0, int(max_h) / h0)
        ww, wh = int(w0 * scale), int(h0 * scale)
        wall_top = int(SLIDE_H - Inches(0.4) - wh)
        slide.shapes.add_picture(src, int(cx - ww / 2), wall_top, width=ww, height=wh)

    # Event title (auto-fit, up to 3 lines)
    title = (event.get("title") or "").strip()
    t_pt, t_lines = _fit_size(title, F_BLACK, 11.4, 1.7, hi=34, lo=22, max_lines=3)
    t_h = len(t_lines) * (t_pt * 2 * _LINE_FACTOR) / _PX_PER_IN
    _add_text(slide, Inches(0.95), y, Inches(11.47), Inches(t_h + 0.1),
              [[(l, F_BLACK, t_pt, CHARCOAL, True)] for l in t_lines],
              anchor=MSO_ANCHOR.TOP)
    y = y + Inches(t_h + 0.12)

    # Date · venue — vertically centred in the gap between the title and the wall
    dv = "   ·   ".join(x for x in [(event.get("date_str") or "").strip(),
                                    (event.get("venue") or "").strip()] if x)
    if dv:
        gap_h = max(Inches(0.4), wall_top - Inches(0.12) - y)
        _add_text(slide, Inches(1.0), y, Inches(11.33), gap_h,
                  [[(dv.upper(), F_BOLD, 19, _theme_rgb(theme_key), True)]],
                  anchor=MSO_ANCHOR.MIDDLE)
    return slide


# ── Section / break divider slides ────────────────────────────────────────────

_EMOJI_RE = re.compile(
    "[\U0001F000-\U0001FAFF\U00002600-\U000027BF\U0001F900-\U0001F9FF"
    "\U00002190-\U000021FF\U00002B00-\U00002BFF️‍]+")


def _strip_emoji(s: str) -> str:
    return _EMOJI_RE.sub("", s or "").strip()


_BREAK_ICON_NAME = {
    "registration": "checkin", "break": "coffee", "lunch": "meal",
    "cocktail": "cocktail", "networking": "people",
    "welcome": "clock", "closing": "clock",
}

# Split-panel section/utility layout (chosen by Belén 2026-06-16):
#   left charcoal panel with a white icon · right theme-colour area with a
#   left-aligned title, a cream divider rule, and an optional detail line.
_PANEL_RGB = RGBColor(0x1C, 0x25, 0x29)
_CREAM = RGBColor(0xFF, 0xE2, 0xD4)
_PANEL_W = Inches(4.55)


def _section_base(slide, theme_key, icon_name, icon_in=2.5):
    """Theme-colour background + left charcoal panel with a centred white icon.
    Returns the x where the right-hand content column should start (EMU)."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = _theme_rgb(theme_key)
    _no_line(bg); bg.shadow.inherit = False

    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, _PANEL_W, SLIDE_H)
    panel.fill.solid(); panel.fill.fore_color.rgb = _PANEL_RGB
    _no_line(panel); panel.shadow.inherit = False

    icon = svg_icons.render(icon_name, 600, (255, 255, 255)) if icon_name else None
    if icon is not None:
        top = (SLIDE_H - Inches(icon_in)) // 2
        _add_picture_fit(slide, icon, _PANEL_W // 2, top, icon_in, icon_in)
    return _PANEL_W + Inches(0.7)


def _section_text(slide, x0, title, title2=None, detail_runs=None):
    """Vertically-centred, left-aligned title (+ optional 2nd language) with a
    cream divider rule, and optional detail line(s) under it."""
    content_w_in = (SLIDE_W - x0) / EMU_PER_IN - 0.55
    en_pt, en_lines = _fit_size(title, F_BLACK, content_w_in, 3.4,
                                hi=52, lo=24, max_lines=3)
    en_h = len(en_lines) * (en_pt * 2 * _LINE_FACTOR) / _PX_PER_IN
    it_h = it_pt = 0
    it_lines = []
    has2 = bool(title2 and title2 != title)
    if has2:
        it_pt = max(20, int(en_pt * 0.52))
        it_lines, it_h, _ = _measure(title2, F_BOLD, it_pt, content_w_in * _PX_PER_IN)
    det_h = (len(detail_runs) * 0.42) if detail_runs else 0
    total = en_h + (0.10 + it_h if has2 else 0) + 0.26 + (0.30 + det_h if detail_runs else 0)
    y_in = max(0.7, (SLIDE_H / EMU_PER_IN - total) / 2)
    cw = SLIDE_W - x0 - Inches(0.55)

    en_runs = [[(l, F_BLACK, en_pt, WHITE, True)] for l in en_lines]
    _add_text(slide, x0, Inches(y_in), cw, Inches(en_h + 0.1), en_runs,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    y_in += en_h
    if has2:
        y_in += 0.10
        it_runs = [[(l, F_BOLD, it_pt, _CREAM, False)] for l in it_lines]
        _add_text(slide, x0, Inches(y_in), cw, Inches(it_h + 0.1), it_runs,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        y_in += it_h
    y_in += 0.16
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, Inches(y_in),
                                  Inches(min(2.8, content_w_in * 0.55)),
                                  Emu(int(0.05 * EMU_PER_IN)))
    rule.fill.solid(); rule.fill.fore_color.rgb = _CREAM
    _no_line(rule); rule.shadow.inherit = False
    y_in += 0.10
    if detail_runs:
        y_in += 0.30
        _add_text(slide, x0, Inches(y_in), cw, Inches(det_h + 0.2), detail_runs,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    return Inches(y_in)


def add_break_slide(prs: Presentation, theme_key: str, session: dict,
                    lang_strings: dict | None = None, layout: str = "marketing"):
    """Split-panel section/break divider (coffee, lunch, cocktail, registration,
    welcome…): charcoal icon panel + the (bilingual) label. No time shown —
    events run late and there's rarely a chance to correct it. All native."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    x0 = _section_base(slide, theme_key,
                       _BREAK_ICON_NAME.get(session.get("break_kind"), "clock"))
    title = _strip_emoji(session.get("title", "")).upper()
    title2 = _strip_emoji(session.get("title_2", "")).upper()
    _section_text(slide, x0, title, title2=title2)
    return slide


# ── Utility slides (mute phone · WiFi · scan-to-register) ─────────────────────

_SCAN_LABEL = {"en": "Scan to register", "es": "Escanea para registrarte",
               "it": "Scansiona per registrarti", "pl": "Zeskanuj, aby się zarejestrować"}
_NETWORK_LABEL = {"en": "Network", "es": "Red", "it": "Rete", "pl": "Sieć"}
_PASSWORD_LABEL = {"en": "Password", "es": "Contraseña", "it": "Password", "pl": "Hasło"}


_MUTE_SUB = {
    "en": "Please silence your devices during the sessions.",
    "es": "Silencia tus dispositivos durante las sesiones.",
    "it": "Silenzia i tuoi dispositivi durante le sessioni.",
    "pl": "Prosimy o wyciszenie urządzeń podczas sesji.",
}


def add_utility_slide(prs: Presentation, theme_key: str, kind: str,
                      lang_strings: dict | None = None, data: dict | None = None,
                      lang: str = "en"):
    """A housekeeping slide (kind = 'mute' | 'wifi' | 'qr') in the same split-panel
    style as the breaks: charcoal icon panel + left-aligned title/detail. All
    native so a colleague can edit the wording / WiFi / link or delete it."""
    ls = lang_strings or {}
    data = data or {}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    icon_name = {"mute": "volume_off", "wifi": "wifi", "qr": "qr"}.get(kind, kind)
    x0 = _section_base(slide, theme_key, icon_name)
    cw = SLIDE_W - x0 - Inches(0.55)

    if kind == "mute":
        title = ls.get("mute_phone", "Mute your phone").upper()
        sub = _MUTE_SUB.get(lang, _MUTE_SUB["en"])
        _section_text(slide, x0, title,
                      detail_runs=[[(sub, F_BOLD, 22, _CREAM, False)]])

    elif kind == "wifi":
        net = (data.get("network") or "____________").strip()
        pwd = (data.get("password") or "____________").strip()
        nlab = _NETWORK_LABEL.get(lang, "Network").upper()
        plab = _PASSWORD_LABEL.get(lang, "Password").upper()
        _section_text(slide, x0, "WIFI", detail_runs=[
            [(f"{nlab}   ", F_REG, 24, _CREAM, False), (net, F_BOLD, 24, WHITE, True)],
            [(f"{plab}   ", F_REG, 24, _CREAM, False), (pwd, F_BOLD, 24, WHITE, True)]])

    elif kind == "qr":
        title = _SCAN_LABEL.get(lang, "Scan to register").upper()
        _add_text(slide, x0, Inches(1.35), cw, Inches(1.7),
                  [[(title, F_BLACK, 40, WHITE, True)]],
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        box = Inches(2.3)
        ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    int(x0), Inches(3.45), box, box)
        ph.fill.solid(); ph.fill.fore_color.rgb = WHITE
        ph.line.color.rgb = WHITE; ph.shadow.inherit = False
        ph.text_frame.text = "QR"
        _p = ph.text_frame.paragraphs[0]; _p.alignment = PP_ALIGN.CENTER
        _r = _p.runs[0]; _r.font.name = F_BOLD; _r.font.size = Pt(28)
        _r.font.color.rgb = _theme_rgb(theme_key)
        url = (data.get("url") or "").strip()
        if url:
            _add_text(slide, x0, Inches(6.0), cw, Inches(0.7),
                      [[(url, F_BOLD, 20, _CREAM, True)]],
                      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    return slide


# ── Full event deck: combined panel slides + individual cards ─────────────────

def build_event_deck(agenda: dict, theme_key: str, out_path: str,
                     photos: dict | None = None, logos: dict | None = None,
                     lang_strings: dict | None = None,
                     include_cards: bool = True, layout: str = "marketing",
                     title_fit: str = "flexible", include_breaks: bool = True,
                     cover: bool = True, lang: str = "en",
                     utility_kinds: "list | None" = None,
                     event_band: str = "bottom", spx_space: bool = False) -> str:
    """
    Optional cover slide, then for each session in agenda order:
      • break  → full-bleed section/break divider (coffee, lunch, cocktail…)
      • panel  → combined panel slide, then (optional) one card per panelist
      • presentation/speaker → one speaker card
    layout: 'event' or 'marketing'. For event decks the title band position is
    `event_band` ('top' or 'bottom'); marketing is always top.
    """
    photos = photos or {}
    logos = logos or {}
    band_pos = "top" if layout == "marketing" else event_band
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    if cover and agenda.get("event", {}).get("title"):
        add_cover_slide(prs, theme_key, agenda["event"],
                        agenda.get("logo_wall"), lang_strings)

    if utility_kinds:
        ev = agenda.get("event", {})
        for k in utility_kinds:
            add_utility_slide(prs, theme_key, k, lang_strings,
                              data={"url": ev.get("register_url")}, lang=lang)

    for sess in agenda.get("sessions", []):
        stype = sess.get("type")
        if stype == "break":
            if include_breaks:
                add_break_slide(prs, theme_key, sess, lang_strings, layout)
            continue
        named = [s for s in sess.get("speakers", []) if (s.get("name") or "").strip()]
        if not named:
            continue
        if stype == "panel":
            add_panel_slide(prs, theme_key, sess, photos, logos, lang_strings,
                            layout=layout, title_fit=title_fit, band_pos=band_pos,
                            spx_space=spx_space)
            if include_cards:
                for sp in named:
                    add_single_speaker_slide(
                        prs, theme_key, sess, sp,
                        photo=photos.get(sp.get("name", "")),
                        logo=logos.get((sp.get("company") or "").strip()),
                        lang_strings=lang_strings, layout=layout,
                        title_fit=title_fit, band_pos=band_pos, spx_space=spx_space)
        elif stype in ("presentation", "speaker"):
            for sp in named:
                add_single_speaker_slide(
                    prs, theme_key, sess, sp,
                    photo=photos.get(sp.get("name", "")),
                    logo=logos.get((sp.get("company") or "").strip()),
                    lang_strings=lang_strings, layout=layout,
                    title_fit=title_fit, band_pos=band_pos, spx_space=spx_space)

    prs.save(out_path)
    return out_path


# ── Deck builder (vertical slice: individual cards for every named speaker) ───

def build_speaker_cards_deck(agenda: dict, theme_key: str, out_path: str,
                             photos: dict | None = None,
                             logos: dict | None = None,
                             lang_strings: dict | None = None) -> str:
    """
    Build a .pptx with ONE single-speaker card per named speaker across all
    speaker/presentation/panel sessions. Skips sessions with no named speakers.
    photos: {speaker_name: PIL.Image}; logos: {company: PIL.Image}.
    """
    photos = photos or {}
    logos = logos or {}
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for sess in agenda.get("sessions", []):
        if sess.get("type") not in ("panel", "presentation", "speaker"):
            continue
        for sp in sess.get("speakers", []):
            name = sp.get("name", "")
            add_single_speaker_slide(
                prs, theme_key, sess, sp,
                photo=photos.get(name),
                logo=logos.get((sp.get("company") or "").strip()),
                lang_strings=lang_strings,
            )

    prs.save(out_path)
    return out_path
