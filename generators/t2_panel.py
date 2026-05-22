"""
T2 — Panel Title Slide  1920 × 1080
Dark dramatic style: optional photo background with overlay,
bold header, centered speaker cards, event-detail footer.
"""
import io, math
from PIL import Image, ImageDraw, ImageFilter, ImageFont
from utils.image_utils import make_circular
from config import LANGUAGE_STRINGS

W, H = 1920, 1080


def _font(name: str, size: int):
    from utils.fonts import get
    return get({"heavy": "montserrat-black",
                "bold":  "montserrat-bold",
                "reg":   "inter"}.get(name, "inter"), size)


def _wrap(draw, text, font, max_w):
    words = text.split()
    lines, cur = [], ""
    for w in words:
        candidate = (cur + " " + w).strip()
        if draw.textbbox((0, 0), candidate, font=font)[2] > max_w and cur:
            lines.append(cur); cur = w
        else:
            cur = candidate
    if cur:
        lines.append(cur)
    return lines


def _draw_text_centered(draw, text, font, cx, y, fill):
    """Draw horizontally centred text; returns new y below the text."""
    bb = draw.textbbox((0, 0), text, font=font)
    draw.text((cx - bb[2] // 2, y), text, font=font, fill=fill)
    return y + bb[3]


def _build_canvas(bg_image: Image.Image | None, theme_rgb: tuple) -> Image.Image:
    """Create base canvas: photo bg with overlay, or dark gradient."""
    canvas = Image.new("RGB", (W, H), (15, 15, 20))

    if bg_image:
        bg = bg_image.convert("RGB").resize((W, H), Image.LANCZOS)
        # Slight blur so it doesn't compete with text
        bg = bg.filter(ImageFilter.GaussianBlur(radius=3))
        canvas.paste(bg)
        # Dark overlay  (~65% opacity)
        overlay = Image.new("RGB", (W, H), (10, 10, 15))
        mask    = Image.new("L",   (W, H), 165)
        canvas.paste(overlay, mask=mask)
    else:
        # Subtle vignette gradient: very dark edges, slightly lighter centre
        draw = ImageDraw.Draw(canvas)
        for y in range(H):
            ratio = math.sin(y / H * math.pi)          # 0 at edges, 1 at centre
            r = int(15 + ratio * 12)
            g = int(15 + ratio * 12)
            b = int(20 + ratio * 15)
            draw.line([(0, y), (W, y)], fill=(r, g, b))

    return canvas


def _draw_header(canvas: Image.Image, draw: ImageDraw.ImageDraw,
                 title: str, theme_rgb: tuple, renmad_logo: Image.Image | None):
    """Coloured accent bar + bold session title + logo."""
    accent_h = 8
    draw.rectangle([0, 0, W, accent_h], fill=theme_rgb)

    # Semi-transparent dark header zone
    header_zone_h = 160
    hz = Image.new("RGBA", (W, header_zone_h), (*theme_rgb, 40))
    canvas.paste(Image.alpha_composite(
        canvas.crop((0, accent_h, W, accent_h + header_zone_h)).convert("RGBA"), hz
    ).convert("RGB"), (0, accent_h))

    # Logo
    logo_w = 0
    if renmad_logo:
        logo = renmad_logo.convert("RGBA")
        lh   = 100
        lw   = int(logo.width * lh / logo.height)
        logo = logo.resize((lw, lh), Image.LANCZOS)
        logo_x = W - lw - 30
        logo_y = accent_h + (header_zone_h - lh) // 2
        canvas.paste(logo, (logo_x, logo_y), mask=logo)
        logo_w = lw + 50

    # Session title — up to 2 lines, centred in the available width
    tf          = _font("heavy", 58)
    text        = title.upper()
    avail_w     = W - logo_w - 60
    cx_title    = avail_w // 2          # centre of the available zone
    lines       = _wrap(draw, text, tf, avail_w)[:2]
    total_h     = sum(draw.textbbox((0, 0), l, font=tf)[3] + 4 for l in lines)
    ty = accent_h + (header_zone_h - total_h) // 2
    for line in lines:
        lh2 = draw.textbbox((0, 0), line, font=tf)[3]
        lw2 = draw.textbbox((0, 0), line, font=tf)[2]
        draw.text((cx_title - lw2 // 2, ty), line, font=tf, fill=(255, 255, 255))
        ty += lh2 + 4

    return accent_h + header_zone_h   # bottom of header zone


def _draw_footer(canvas: Image.Image, draw: ImageDraw.ImageDraw,
                 theme_rgb: tuple, date_str: str, time_str: str,
                 location: str, cta_url: str = ""):
    """Coloured footer band with event details and optional CTA/URL."""
    has_url  = bool(cta_url.strip())
    footer_h = 110 if has_url else 90
    fy       = H - footer_h

    foot_img = Image.new("RGBA", (W, footer_h), (*theme_rgb, 230))
    canvas.paste(foot_img.convert("RGB"), (0, fy))
    draw.line([(0, fy), (W, fy)], fill=(255, 255, 255), width=2)

    parts = [p for p in [date_str, time_str, location] if p]
    if has_url:
        # Two lines: event details top, URL bottom
        detail_y = fy + 14
        url_y    = fy + footer_h - 34
        if parts:
            detail = "   |   ".join(parts).upper()
            df     = _font("bold", 26)
            bb     = draw.textbbox((0, 0), detail, font=df)
            draw.text((W // 2 - bb[2] // 2, detail_y), detail, font=df, fill=(255, 255, 255))
        uf = _font("reg", 22)
        ub = draw.textbbox((0, 0), cta_url, font=uf)
        draw.text((W // 2 - ub[2] // 2, url_y), cta_url, font=uf,
                  fill=(255, 255, 200))
    else:
        if parts:
            detail = "   |   ".join(parts).upper()
            df     = _font("bold", 28)
            bb     = draw.textbbox((0, 0), detail, font=df)
            cy     = fy + footer_h // 2
            draw.text((W // 2 - bb[2] // 2, cy - bb[3] // 2),
                      detail, font=df, fill=(255, 255, 255))

    return fy


def generate(session: dict, theme: dict, language: str = "es",
             renmad_logo: Image.Image | None = None,
             bg_image:    Image.Image | None = None) -> tuple[bytes, bytes]:

    strings   = LANGUAGE_STRINGS.get(language, LANGUAGE_STRINGS["es"])
    theme_rgb = tuple(theme["rgb"])

    canvas = _build_canvas(bg_image, theme_rgb)
    draw   = ImageDraw.Draw(canvas)

    header_bottom = _draw_header(canvas, draw,
                                 session.get("title", ""), theme_rgb, renmad_logo)

    footer_top = _draw_footer(canvas, draw, theme_rgb,
                               session.get("date_str", ""),
                               session.get("time_str", ""),
                               session.get("location",  ""),
                               session.get("cta_url",   ""))

    # ── Speaker area ──────────────────────────────────────────────────────────
    speakers  = session.get("speakers", [])
    n         = max(len(speakers), 1)
    area_top  = header_bottom + 20
    area_bot  = footer_top   - 20
    area_h    = area_bot - area_top
    col_w     = W // n

    has_photos = any(spk.get("photo") for spk in speakers)
    photo_size = min(220, int(area_h * 0.42)) if has_photos else 0

    # Fonts
    name_f    = _font("bold",  32)
    jobtitle_f = _font("reg",  24)
    company_f = _font("bold",  24)
    mod_f     = _font("bold",  20)

    for i, spk in enumerate(speakers):
        cx = i * col_w + col_w // 2

        # Vertical centering: total block height
        block_h = photo_size
        if has_photos and photo_size:
            block_h += 16   # gap photo → text
        # estimate text lines: mod label + name + job + company + logo
        block_h += 30 + 34 + 28 + 28  # approximate
        if spk.get("company_logo"):
            block_h += 56

        block_start_y = area_top + (area_h - block_h) // 2
        ty = block_start_y

        # Photo circle
        if has_photos:
            if spk.get("photo"):
                try:
                    photo = make_circular(spk["photo"].convert("RGBA"), photo_size,
                                          border_color=theme_rgb, border_width=4)
                    px = cx - photo.width // 2
                    canvas.paste(photo, (px, ty), mask=photo)
                except Exception:
                    pass
            ty += photo_size + 16

        # Thin separator line above speaker info
        draw.line([(cx - 80, ty), (cx + 80, ty)], fill=theme_rgb, width=2)
        ty += 10

        # Company logo
        if spk.get("company_logo"):
            try:
                clogo = spk["company_logo"].convert("RGBA")
                clogo.thumbnail((160, 44), Image.LANCZOS)
                lx = cx - clogo.width // 2
                # White-ish background pill behind logo for contrast
                pad = 8
                pill = Image.new("RGBA", (clogo.width + pad*2, clogo.height + pad*2), (240,240,240,210))
                canvas.paste(pill.convert("RGB"), (lx - pad, ty - pad))
                canvas.paste(clogo, (lx, ty), mask=clogo)
                ty += clogo.height + pad * 2 + 6
            except Exception:
                pass

        # Moderator label
        if spk.get("is_moderator"):
            label = f"— {strings['moderator'].upper()} —"
            lb    = draw.textbbox((0, 0), label, font=mod_f)
            draw.text((cx - lb[2] // 2, ty), label, font=mod_f, fill=theme_rgb)
            ty += lb[3] + 6

        # Name
        name = spk.get("name", "")
        nb   = draw.textbbox((0, 0), name, font=name_f)
        draw.text((cx - nb[2] // 2, ty), name, font=name_f, fill=(255, 255, 255))
        ty += nb[3] + 5

        # Job title (wrap if very long)
        job = spk.get("title", "")
        if len(job) > 50:
            job = job[:48] + "…"
        jb = draw.textbbox((0, 0), job, font=jobtitle_f)
        draw.text((cx - jb[2] // 2, ty), job, font=jobtitle_f, fill=(200, 200, 200))
        ty += jb[3] + 4

        # Company
        company = spk.get("company", "")
        cb = draw.textbbox((0, 0), company, font=company_f)
        draw.text((cx - cb[2] // 2, ty), company, font=company_f,
                  fill=(*theme_rgb, 255) if len(theme_rgb) == 3 else theme_rgb)

        # Vertical divider between speakers (except last)
        if i < n - 1:
            div_x = (i + 1) * col_w
            draw.line([(div_x, area_top + 30), (div_x, area_bot - 30)],
                      fill=(80, 80, 80), width=1)

    png_buf = io.BytesIO()
    canvas.save(png_buf, format="PNG", dpi=(150, 150))
    return _build_pptx(canvas, W, H), png_buf.getvalue()


def _build_pptx(img, w, h):
    from pptx import Presentation
    from pptx.util import Emu
    DPI = 96
    def px(p): return Emu(int(p / DPI * 914400))
    prs = Presentation()
    prs.slide_width  = px(w)
    prs.slide_height = px(h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    buf = io.BytesIO()
    img.save(buf, format="PNG"); buf.seek(0)
    slide.shapes.add_picture(buf, 0, 0, px(w), px(h))
    out = io.BytesIO(); prs.save(out)
    return out.getvalue()
