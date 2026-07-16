"""
Template A — "Circles"  1920 × 1080

Two-phase rendering:
  Phase 1: background, shapes, photos (→ pptx_canvas, used as PPTX image layer)
  Phase 2: session text drawn on top (→ PNG only)
  PPTX:    pptx_canvas (no text) + editable text boxes = zero doubling
"""
import io, math
from PIL import Image, ImageDraw, ImageFilter, ImageOps
from utils.image_utils import make_circular
from utils.icons import pin, arrow, paste_icon_text
from utils.gender import moderator_label
from utils.text_utils import wrap_title
from config import LANGUAGE_STRINGS
import re

W, H = 1920, 1080
SPLIT_Y    = 640
LEFT_COVER = 1150
CIRCLE_D   = 280
CIRCLE_R   = CIRCLE_D // 2
BORDER_W   = 5

LOGO_H_TARGET = 54
LOGO_PILL_H   = LOGO_H_TARGET + 32   # = 86

_HREF_RE = re.compile(
    r'(?:https?://\S+|www\.\S+|\S+\.(?:com|es|org|net|io|eu)/\S+)', re.I)


def _extract_url(text):
    m = _HREF_RE.search(text)
    return m.group(0) if m else text


def _f(name, size):
    from utils.fonts import get
    return get({"heavy": "montserrat-black", "bold": "montserrat-bold",
                "reg": "inter"}.get(name, "inter"), size)


def _wrap(draw, text, font, max_w):
    words, lines, cur = text.split(), [], ""
    for w in words:
        t = (cur + " " + w).strip()
        if draw.textbbox((0, 0), t, font=font)[2] > max_w and cur:
            lines.append(cur); cur = w
        else:
            cur = t
    if cur:
        lines.append(cur)
    return lines


def _build_bg(bg_image, theme_rgb):
    canvas = Image.new("RGB", (W, H), (14, 14, 18))
    if bg_image:
        bg = bg_image.convert("RGB").resize((W, H), Image.LANCZOS)
        canvas.paste(bg)
    else:
        d = ImageDraw.Draw(canvas)
        for y in range(H):
            v = int(20 + math.sin(y / H * math.pi) * 12)
            d.line([(0, y), (W, y)], fill=(v, v, v + 4))
    return canvas


def _dark_gradient(w, h, alpha_max=210):
    grad = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d    = ImageDraw.Draw(grad)
    for x in range(w):
        t = x / w
        a = int(alpha_max * (1 - t ** 0.65))
        d.line([(x, 0), (x, h - 1)], fill=(8, 8, 14, a))
    return grad


def _silhouette(size, theme_rgb, border_color=None, bw=4):
    total = size + (bw * 2 if border_color else 0)
    cv = Image.new("RGBA", (total, total), (0, 0, 0, 0))
    cd = ImageDraw.Draw(cv)
    if border_color:
        cd.ellipse([0, 0, total-1, total-1], fill=(*border_color, 255))
        off = bw
    else:
        off = 0
    cd.ellipse([off, off, off+size-1, off+size-1], fill=(32, 32, 40, 255))
    cx = off + size // 2; r = size // 2
    tr, tg, tb = theme_rgb; col = (tr, tg, tb, 180)
    hrad = int(r * 0.22); hy = off + int(r * 0.38)
    cd.ellipse([cx-hrad, hy-hrad, cx+hrad, hy+hrad], fill=col)
    bwd = int(r * 0.38); by = hy + hrad + 2
    cd.ellipse([cx-bwd, by, cx+bwd, by+int(r*0.72)], fill=col)
    return cv


def _trim_logo(img):
    """Remove transparent border padding before scaling.
    Ensures the logo content (not blank padding) fills LOGO_H_TARGET."""
    rgba = img.convert("RGBA")
    alpha = rgba.split()[3]
    bbox  = alpha.getbbox()          # tight bbox of non-transparent pixels
    if bbox:
        rgba = rgba.crop(bbox)
    # If image appears to have a white/light background with no real alpha,
    # try a brightness-based trim so padded white doesn't inflate the logo size.
    r, g, b, a = rgba.split()
    import PIL.ImageChops as IC
    # Build a mask of "non-white" pixels (any channel < 240)
    dark = IC.darker(IC.darker(r, g), b)
    dark_thresh = dark.point(lambda v: 255 if v < 240 else 0)
    wb = dark_thresh.getbbox()
    if wb:
        rgba = rgba.crop(wb)
    return rgba


def _logo_pill(logo_img, col_budget=260):
    """White pill; logo box-fitted within col_budget × LOGO_H_TARGET.

    Wide logos are scaled DOWN (both axes) so they never exceed the column
    width and never block adjacent logos.  Pill height is always LOGO_PILL_H;
    logo is centred inside both axes.
    """
    try:
        clogo = _trim_logo(logo_img)
    except Exception:
        clogo = logo_img.convert("RGBA")
    lw_orig, lh_orig = clogo.size
    if lh_orig == 0:
        return None
    pad    = 16
    max_cw = max(20, col_budget - pad * 2)
    max_ch = LOGO_H_TARGET
    scale  = min(max_ch / lh_orig, max_cw / lw_orig)
    new_w  = max(1, int(lw_orig * scale))
    new_h  = max(1, int(lh_orig * scale))
    clogo  = clogo.resize((new_w, new_h), Image.LANCZOS)
    pill_w = new_w + pad * 2
    pill   = Image.new("RGBA", (pill_w, LOGO_PILL_H), (0, 0, 0, 0))
    ImageDraw.Draw(pill).rounded_rectangle(
        [0, 0, pill_w - 1, LOGO_PILL_H - 1], radius=12, fill=(255, 255, 255, 255))
    px = (pill_w - new_w) // 2
    py = (LOGO_PILL_H - new_h) // 2
    pill.paste(clogo, (px, py), mask=clogo)
    return pill


def _pill_badge(draw, text, font, x, y, bg_color,
                text_color=(255, 255, 255), pad=14, shape_only=False):
    """Filled rounded-rectangle badge.  Returns (right_x, bottom_y).
    shape_only=True draws only the background rectangle (text deferred to Phase 2)."""
    bb = draw.textbbox((0, 0), text, font=font)
    bw, bh = bb[2] + pad * 2, bb[3] + pad * 2
    draw.rounded_rectangle([x, y, x + bw, y + bh], radius=6, fill=bg_color)
    if not shape_only:
        draw.text((x + pad, y + pad), text, font=font, fill=text_color)
    return x + bw, y + bh


def generate(session: dict, theme: dict, language: str = "es",
             renmad_logo: Image.Image | None = None,
             bg_image:    Image.Image | None = None,
             ata_logo:    Image.Image | None = None,
             variant:     str = "normal",
             hide_speakers: bool = False) -> tuple[bytes, bytes]:

    strings   = LANGUAGE_STRINGS.get(language, LANGUAGE_STRINGS["es"])
    theme_rgb = tuple(theme["rgb"])

    # ── Phase 1: build all visual elements (NO session text) ─────────────────
    # text_draws: [(x, y, text, font, fill)] drawn in Phase 2 (PNG only)
    # text_els:   [{...}] written as PPTX text boxes
    text_draws = []
    text_els   = []
    logo_els   = []   # company logos → separate resizable PPTX picture shapes
    photo_els  = []   # speaker photos → movable, re-croppable PPTX picture shapes

    canvas = _build_bg(bg_image, theme_rgb)
    # No speakers → dark wash covers the whole canvas (no white speaker strip)
    dark_h = H if hide_speakers else SPLIT_Y
    grad   = _dark_gradient(W, dark_h, alpha_max=215)
    canvas.paste(Image.new("RGB", (W, dark_h), (8, 8, 14)),
                 (0, 0), mask=grad.split()[3])
    draw = ImageDraw.Draw(canvas)

    draw.rectangle([0, 0, W, 6], fill=theme_rgb)   # accent bar

    pad_l = 52; ty = 28

    # "WEBINAR GRATUITO" badge — branding label, not editable, draw in Phase 1
    # "Grabación disponible" variant swaps the badge text.
    wb_f  = _f("bold", 28)
    badge = strings.get("free_webinar", "WEBINAR GRATUITO")
    if variant == "recording_available":
        badge = strings.get("recording_available", "Grabación disponible").upper()
    _, ty = _pill_badge(draw, badge, wb_f, pad_l, ty,
                        bg_color=theme_rgb, text_color=(255, 255, 255))
    ty += 22

    # ── Title: measure positions, defer text ─────────────────────────────────
    title_f      = _f("heavy", 70)
    title_text   = session.get("title", "").upper()
    title_w      = W - pad_l - 220
    t_lines      = wrap_title(draw, title_text, title_f, title_w, max_lines=3)
    title_start_y = ty
    for line in t_lines:
        lb = draw.textbbox((0, 0), line, font=title_f)
        text_draws.append((pad_l, ty, line, title_f, (255, 255, 255)))
        ty += lb[3] + 8
    ty += 20
    if title_text:
        text_els.append({'x': pad_l, 'y': title_start_y,
                         'w': title_w, 'h': max(ty - title_start_y, 80),
                         'text': title_text, 'size': 70, 'bold': True,
                         'color': (255, 255, 255), 'wrap': True})

    # ── Date / time / location — pill shapes in Phase 1, text deferred ───────
    dt_f  = _f("bold", 34)
    loc_f = _f("bold", 28)
    rx    = pad_l

    ref_bb     = draw.textbbox((0, 0), "Ag", font=dt_f)
    pill_h_ref = ref_bb[3] + 14 * 2
    mid_y      = ty + pill_h_ref // 2

    if session.get("date_str"):
        date_text = session["date_str"].upper()
        date_x0   = rx
        rx, _     = _pill_badge(draw, date_text, dt_f, rx, ty,
                                  bg_color=(20, 20, 20, 200), shape_only=True)
        text_draws.append((date_x0 + 14, ty + 14, date_text, dt_f, (255, 255, 255)))
        text_els.append({'x': date_x0 + 14, 'y': ty + 14,
                         'w': max(rx - date_x0 - 14, 50), 'h': pill_h_ref,
                         'text': date_text, 'size': 34, 'bold': True,
                         'color': (255, 255, 255)})
        rx += 14

    if session.get("time_str") or variant == "starting_soon":
        # "Empezamos pronto" variant replaces the time source before .upper()
        time_src  = (strings.get("starting_soon", "Empezamos pronto")
                     if variant == "starting_soon" else session.get("time_str", ""))
        time_text = time_src.upper()
        time_x0   = rx
        rx, _     = _pill_badge(draw, time_text, dt_f, rx, ty,
                                  bg_color=theme_rgb, shape_only=True)
        text_draws.append((time_x0 + 14, ty + 14, time_text, dt_f, (255, 255, 255)))
        text_els.append({'x': time_x0 + 14, 'y': ty + 14,
                         'w': max(rx - time_x0 - 14, 50), 'h': pill_h_ref,
                         'text': time_text, 'size': 34, 'bold': True,
                         'color': (255, 255, 255)})
        rx += 14

    if session.get("location"):
        ico      = pin(28, (255, 255, 255))
        loc_bb   = draw.textbbox((0, 0), "Ag", font=loc_f)
        icon_y   = mid_y - ico.height // 2
        text_y   = mid_y - loc_bb[3] // 2
        canvas.paste(ico, (rx, icon_y), mask=ico)   # icon → Phase 1 (image)
        loc_text = session["location"].upper()
        lb_loc   = draw.textbbox((0, 0), loc_text, font=loc_f)
        text_draws.append((rx + ico.width + 8, text_y, loc_text, loc_f, (220, 220, 220)))
        text_els.append({'x': rx + ico.width + 8, 'y': text_y,
                         'w': lb_loc[2] + 20, 'h': loc_bb[3] + 10,
                         'text': loc_text, 'size': 28, 'bold': True,
                         'color': (220, 220, 220)})

    # ── RENMAD / ATA logo top-right ───────────────────────────────────────────
    if renmad_logo:
        logo = renmad_logo.convert("RGBA")
        lh = 96; lw = int(logo.width * lh / logo.height)
        logo = logo.resize((lw, lh), Image.LANCZOS)
        canvas.paste(logo, (W - lw - 28, 16), mask=logo)
    if ata_logo and not renmad_logo:
        al = ata_logo.convert("RGBA")
        alh = 72; alw = int(al.width * alh / al.height)
        al = al.resize((alw, alh), Image.LANCZOS)
        canvas.paste(al, (W - alw - 28, 20), mask=al)

    # ── White speaker strip ───────────────────────────────────────────────────
    if not hide_speakers:
        draw.rectangle([0, SPLIT_Y, W, H], fill=(255, 255, 255))
        draw.rectangle([0, SPLIT_Y, W, SPLIT_Y + 5], fill=theme_rgb)

    # ── Speaker photos, logo pills, dividers ──────────────────────────────────
    speakers = session.get("speakers", [])
    n        = max(len(speakers), 1)
    col_w    = W // n

    col_budget = max(80, col_w - 40)
    logo_pills = []
    for spk in speakers:
        if spk.get("company_logo"):
            try:    logo_pills.append(_logo_pill(spk["company_logo"], col_budget))
            except: logo_pills.append(None)
        else:
            logo_pills.append(None)

    name_f = _f("bold", 30); role_f = _f("reg", 23); mod_f = _f("bold", 20)

    max_name_h = max((draw.textbbox((0, 0), s.get("name",  ""), font=name_f)[3]
                      for s in speakers), default=36)
    max_role_h = max((draw.textbbox((0, 0), s.get("title", ""), font=role_f)[3]
                      for s in speakers), default=28)

    below_circle = SPLIT_Y + CIRCLE_R + BORDER_W + 14
    logo_y = below_circle
    name_y = logo_y + LOGO_PILL_H + 12
    role_y = name_y + max_name_h + 6
    mod_y  = role_y + max_role_h * 2 + 8

    for i, (spk, pill) in enumerate(zip(speakers, logo_pills)):
        cx = i * col_w + col_w // 2
        tw = col_w - 60; tx = cx - tw // 2

        # Photo / silhouette. Real photos are NOT baked into the background —
        # they are pasted in Phase 2 for the PNG and become movable,
        # re-croppable picture shapes in the PPTX.
        photo = None
        if spk.get("photo"):
            try:
                photo = make_circular(spk["photo"].convert("RGBA"), CIRCLE_D,
                                      border_color=theme_rgb, border_width=BORDER_W)
            except Exception:
                photo = None
        if photo is not None:
            sd = CIRCLE_D + BORDER_W   # outline centred on edge → outer Ø matches PNG
            photo_els.append({'pil': photo,
                              'px': cx - photo.width // 2,
                              'py': SPLIT_Y - photo.height // 2,
                              'x': cx - sd / 2, 'y': SPLIT_Y - sd / 2,
                              'w': sd, 'h': sd, 'src': spk["photo"]})
        else:
            sil = _silhouette(CIRCLE_D, theme_rgb, border_color=theme_rgb, bw=BORDER_W)
            canvas.paste(sil,
                         (cx - sil.width // 2, SPLIT_Y - sil.height // 2),
                         mask=sil)

        # Logo pill — collected here; pasted in Phase 2 (PNG) / added as PPTX shape
        if pill is not None:
            lx = cx - pill.width // 2
            ibuf = io.BytesIO(); pill.save(ibuf, format="PNG")
            logo_els.append({'x': lx, 'y': logo_y,
                             'w': pill.width, 'h': pill.height,
                             'img_bytes': ibuf.getvalue(), 'pill': pill})

        # Column divider
        if i < n - 1:
            dx2 = (i + 1) * col_w
            draw.line([(dx2, SPLIT_Y + 5), (dx2, H - 64)],
                      fill=(210, 210, 210), width=1)

        # Defer speaker text ──────────────────────────────────────────────────
        name = spk.get("name", "")
        if name:
            nb = draw.textbbox((0, 0), name, font=name_f)
            text_draws.append((cx - nb[2] // 2, name_y, name, name_f, (20, 20, 20)))
            text_els.append({'x': tx, 'y': name_y, 'w': tw, 'h': max_name_h + 10,
                             'text': name, 'size': 30, 'bold': True,
                             'color': (20, 20, 20), 'center': True})

        role = spk.get("title", "")
        if role:
            rl_lines = _wrap(draw, role, role_f, col_w - 40)[:2]
            ry = role_y
            for rl in rl_lines:
                rb = draw.textbbox((0, 0), rl, font=role_f)
                text_draws.append((cx - rb[2] // 2, ry, rl, role_f, (80, 80, 80)))
                ry += rb[3] + 3
            text_els.append({'x': tx, 'y': role_y, 'w': tw, 'h': max_role_h * 2 + 10,
                             'text': role, 'size': 23, 'color': (80, 80, 80),
                             'center': True, 'wrap': True})

        if spk.get("is_moderator"):
            label = f"— {moderator_label(spk.get('name', ''), strings)} —"
            mb    = draw.textbbox((0, 0), label, font=mod_f)
            text_draws.append((cx - mb[2] // 2, mod_y, label, mod_f, theme_rgb))
            text_els.append({'x': tx, 'y': mod_y, 'w': tw, 'h': 32,
                             'text': label, 'size': 20, 'bold': True,
                             'color': list(theme_rgb), 'center': True})

    # ── CTA bar shape + arrow icon (Phase 1); text added via hyperlink in PPTX ─
    cta_url_raw = session.get("cta_url", "").strip()
    cta_rect    = None
    if cta_url_raw:
        href     = _extract_url(cta_url_raw)
        label    = strings.get("cta_label", "Regístrate gratis hoy")
        cta_text = f"{label}  {href}"
        cta_h    = 52
        draw.rectangle([0, H - cta_h, W, H], fill=(*theme_rgb,))
        cta_f  = _f("bold", 26)
        ico    = arrow(26, (255, 255, 255))
        cb     = draw.textbbox((0, 0), cta_text, font=cta_f)
        tot_w  = ico.width + 10 + cb[2]
        cta_x  = W // 2 - tot_w // 2
        cta_y  = H - cta_h + (cta_h - max(ico.height, cb[3])) // 2
        canvas.paste(ico, (cta_x, cta_y), mask=ico)   # icon in Phase 1
        text_draws.append((cta_x + ico.width + 10, cta_y, cta_text,
                           cta_f, (255, 255, 255)))    # full label + URL for PNG
        cta_rect = (href, cta_x, H - cta_h, tot_w, cta_h)

    # ── PPTX snapshot: all shapes + photos, ZERO session text ────────────────
    pptx_canvas = canvas.copy()

    # ── Phase 2: speaker photos, logo pills + session text for PNG ───────────
    for pe in photo_els:
        canvas.paste(pe['pil'], (pe['px'], pe['py']), mask=pe['pil'])
    for le in logo_els:
        canvas.paste(le['pill'], (le['x'], le['y']), mask=le['pill'])
    for x, y, txt, fnt, col in text_draws:
        draw.text((x, y), txt, font=fnt, fill=col)

    png_buf = io.BytesIO()
    canvas.save(png_buf, format="PNG", dpi=(150, 150))
    return (_pptx(pptx_canvas, cta_rect, text_els, logo_els, photo_els, theme_rgb),
            png_buf.getvalue())


def _pptx(img, cta_rect=None, text_els=None, logo_els=None,
          photo_els=None, theme_rgb=None):
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    def px(p): return Emu(int(p / 96 * 914400))
    prs = Presentation()
    prs.slide_width  = px(W); prs.slide_height = px(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background: clean image with NO baked-in session text
    buf = io.BytesIO(); img.save(buf, format="PNG"); buf.seek(0)
    slide.shapes.add_picture(buf, 0, 0, px(W), px(H))

    # Speaker photos — movable picture shapes; Crop tool re-frames the face
    from utils.pptx_editable import add_shaped_picture
    for pe in (photo_els or []):
        try:
            add_shaped_picture(slide, pe['src'], pe['x'], pe['y'],
                               pe['w'], pe['h'], shape="ellipse",
                               border_rgb=theme_rgb, border_px=BORDER_W)
        except Exception:
            pass

    # Editable text boxes — text only appears here (no image doubling)
    for el in (text_els or []):
        _w, _h = el.get('w', 400), el.get('h', 60)
        if _w <= 0 or _h <= 0 or not el.get('text'):
            continue
        txBox = slide.shapes.add_textbox(px(el['x']), px(el['y']), px(_w), px(_h))
        txBox.fill.background()
        txBox.line.fill.background()
        tf = txBox.text_frame
        tf.word_wrap = el.get('wrap', False)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if el.get('center') else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = el['text']
        # canvas px → pt at 96 dpi, so PPTX text matches the PNG layout
        run.font.size = Pt(el.get('size', 20) * 0.75)
        run.font.bold = el.get('bold', False)
        run.font.name = "Montserrat" if el.get('bold') else "Inter"
        c = el.get('color', (255, 255, 255))
        run.font.color.rgb = RGBColor(int(c[0]), int(c[1]), int(c[2]))

    # Company logos — independent picture shapes (selectable, resizable, replaceable)
    for le in (logo_els or []):
        buf = io.BytesIO(le['img_bytes'])
        slide.shapes.add_picture(buf, px(le['x']), px(le['y']),
                                 px(le['w']), px(le['h']))

    # CTA hyperlink text box
    if cta_rect:
        url, cx, cy, cw, ch = cta_rect
        txBox = slide.shapes.add_textbox(px(cx), px(cy), px(cw + 40), px(ch))
        txBox.fill.background(); txBox.line.fill.background()
        tf = txBox.text_frame; tf.word_wrap = False
        run = tf.paragraphs[0].add_run()
        run.text = "  →  " + url
        run.font.size = Pt(20); run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        try:
            rPr = run._r.get_or_add_rPr()
            hl  = etree.SubElement(rPr, qn('a:hlinkClick'))
            rId = slide.part.relate_to(
                url,
                'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
                is_external=True)
            hl.set(qn('r:id'), rId)
        except Exception:
            pass

    out = io.BytesIO(); prs.save(out); return out.getvalue()
