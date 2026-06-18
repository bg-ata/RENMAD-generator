# -*- coding: utf-8 -*-
"""Design-proposal renderer for the ATA webinar report — EU-PROJECTS LOOK.

Light grey canvas + bold solid-orange blocks + charcoal text/dividers, matching
ATA_Insights_Partner_Pitch_v5.pptx. Real stats from the Phase-1 engine + manual
marketing data, rendered to PNG (Pillow) and assembled into a PPTX.
i18n: every visible label is in STRINGS[lang] (en/es/it/pl).
"""
from __future__ import annotations
import os, sys
from PIL import Image, ImageDraw, ImageFont, ImageOps

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from engine.ingest import build_stats, load_registrations

# ---------------------------------------------------------------- paths (portable)
import tempfile
_HERE = os.path.dirname(os.path.abspath(__file__))          # report_core/design
_REPO = os.path.dirname(os.path.dirname(_HERE))             # renmad_generator
FONTS = os.path.join(_REPO, "assets", "fonts")
ATA_LOGO = os.path.join(_REPO, "assets", "logos", "logo_ata_insights.png")
ASSETS = os.path.join(tempfile.gettempdir(), "ata_report_assets")  # per-run; overridden by generate_report
OUT = os.path.join(tempfile.gettempdir(), "ata_report_out")
os.makedirs(ASSETS, exist_ok=True)
os.makedirs(OUT, exist_ok=True)
ANNOTATE = False   # when True, add a "highlight / why it matters" note to each stat slide
SUFFIX = ""        # output filename suffix (e.g. "_annotated")
COVER_PHOTOS = []  # native editable speaker photos for the cover slide (px coords)
LINK_INFO = None   # clickable hyperlink rect for the engagement (YouTube) slide
NATIVE_NUMS = []   # editable text overlays (so EVERY label can be typed over in PowerPoint)
_CUR_SLIDE = 0     # deck index of the slide currently being drawn (set by each slide fn)

# native-text overlay support: text is NOT baked into the slide PNG; instead every
# label is added as a real, editable PowerPoint text box (build_pptx) and painted
# back onto the on-disk PNG only for the in-app preview (_bake_preview).
_KIND2FILE = {"black": "Montserrat-Black.ttf", "bold": "Montserrat-Bold.ttf",
              "semi": "Montserrat-SemiBold.ttf", "reg": "Inter-Regular.ttf"}
_KIND2PPTX = {"black": ("Montserrat Black", False), "bold": ("Montserrat", True),
              "semi": ("Montserrat SemiBold", False), "reg": ("Inter", False)}
_FILE2KIND = {"montserrat-black.ttf": "black", "montserrat-bold.ttf": "bold",
              "montserrat-semibold.ttf": "semi", "inter-regular.ttf": "reg"}
_MEASURE = ImageDraw.Draw(Image.new("RGB", (8, 8)))

def _slide(n):
    global _CUR_SLIDE; _CUR_SLIDE = n

def _kind_of(font):
    base = os.path.basename(getattr(font, "path", "") or "").lower()
    if base in _FILE2KIND:
        return _FILE2KIND[base]
    if "black" in base: return "black"
    if "semibold" in base: return "semi"
    if "bold" in base: return "bold"
    return "reg"

def _px_of(font):
    return max(1, int(round(getattr(font, "size", 24 * SCALE) / SCALE)))

def _font1(kind, px):
    return ImageFont.truetype(os.path.join(FONTS, _KIND2FILE[kind]), int(px))

def nnum(text, x, y, px, rgb, anchor, slide, kind="black", url=None):
    """Register an EDITABLE figure (not drawn on the baked canvas). `px` is the
    pixel size on the 1920-wide canvas; `slide` is the deck index; `anchor` is the
    PIL anchor used for the original layout (la / mm / rm / lm)."""
    NATIVE_NUMS.append({"slide": slide, "x": x, "y": y, "px": px, "text": str(text),
                        "rgb": tuple(rgb), "anchor": anchor, "kind": kind, "url": url, "spacing": 4})

def _num_bbox(e):
    """Tight pixel bbox (l,t,r,b) of the figure on the 1920-wide canvas."""
    f = _font1(e["kind"], e["px"])
    anch = e["anchor"]
    txt = e["text"]
    if "\n" in txt and anch[1:] not in ("a", "m", "s"):
        anch = anch[0] + "a"
    bb = _MEASURE.textbbox((0, 0), txt, font=f, anchor=anch, spacing=e.get("spacing", 4))
    return (e["x"] + bb[0], e["y"] + bb[1], e["x"] + bb[2], e["y"] + bb[3])

def _num_bbox(e):
    """Tight pixel bbox (l,t,r,b) of the figure on the 1920-wide canvas."""
    f = _font1(e["kind"], e["px"])
    bb = _MEASURE.textbbox((0, 0), e["text"], font=f, anchor=e["anchor"])
    return (e["x"] + bb[0], e["y"] + bb[1], e["x"] + bb[2], e["y"] + bb[3])

# ---------------------------------------------------------------- palette (EU look)
ORANGE = (255, 74, 0)
CORAL  = (255, 107, 53)
LIGHT  = (255, 145, 77)
GREY   = (236, 237, 238)   # canvas background
TRACK  = (220, 221, 223)   # bar track on grey
LINE   = (205, 206, 208)
CHAR   = (28, 37, 41)
DEEP   = (16, 22, 24)
WHITE  = (255, 255, 255)
MUTE   = (120, 124, 128)
YELLOW = (255, 196, 0)
TEAL   = (0, 168, 150)
GREEN  = (46, 125, 50)
RAMP   = [ORANGE, CORAL, LIGHT, (255, 179, 122), (255, 205, 160), CHAR]
ACCENTS = [ORANGE, CORAL, TEAL, YELLOW, GREEN, LIGHT]
ICON = {"people":"E716","org":"E825","globe":"E774","video":"E714","gauge":"EC4A","mic":"E720",
        "clock":"E917","group":"E902","pin":"E707","send":"E724","mail":"E715","click":"E7C9",
        "play":"E768","star":"E735","like":"E8E1","megaphone":"E789","power":"E81D","brief":"E821","phone":"E717"}
_icon_candidates = [os.path.join(FONTS, "SegoeIcons.ttf"), r"C:\Windows\Fonts\SegoeIcons.ttf",
                    r"C:\Windows\Fonts\segmdl2.ttf"]
ICON_FONT_PATH = next((p for p in _icon_candidates if os.path.exists(p)), _icon_candidates[0])
# fixed contact portrait (Cintia) bundled in the repo
CONTACT_PHOTO = os.path.join(_REPO, "assets", "report", "cintia_tight.png")

W, H = 1920, 1080
SCALE = 2

# ---------------------------------------------------------------- fonts
def _f(name, size): return ImageFont.truetype(os.path.join(FONTS, name), size * SCALE)
def black(s): return _f("Montserrat-Black.ttf", s)
def bold(s):  return _f("Montserrat-Bold.ttf", s)
def semi(s):  return _f("Montserrat-SemiBold.ttf", s)
def reg(s):   return _f("Inter-Regular.ttf", s)

# ---------------------------------------------------------------- i18n
STRINGS = {
 "en": {"report":"MARKETING & AUDIENCE REPORT","partner":"In partnership with","moderator":"Moderator",
        "sec_facts":"KEY FACTS","facts_title":"The event in numbers","by_country":"By country",
        "facts_ctx":["professionals signed up","unique organisations","across the world","on Zoom + YouTube","of registrants viewed","expert + moderator"],
        "registrations":"Registrations",
        "companies":"Companies","countries":"Countries","live":"Total viewers","attendance":"Attendance rate",
        "speakers":"Speakers","sec_aud":"AUDIENCE","country_title":"Where they joined from","by_region":"By region",
        "countries_rep":"countries represented","ind_title":"Who registered","ind_note":"other / unspecified",
        "sec_comp":"COMPANIES","comp_title":"Who registered","sel_comp":"Selected registered companies",
        "sec_eng":"ENGAGEMENT","eng_title":"How they watched","ge30":"stayed 30 min or more","avg_conn":"Avg. connection",
        "peak":"Peak concurrent","top_country":"Top country","live_countries":"countries live","minutes":"min",
        "retention":"Time connected (Zoom)","retention_fun":"Almost everyone stayed!","watch_yt":"Watch again on YouTube","yt_pending":"views — to be confirmed","views":"views",
        "sec_reach":"MARKETING REACH","reach_title":"How we filled the room","delivered":"emails sent to our database","opened":"opened email",
        "clicked":"clicked one or more links in the email","campaign":"Campaign","sent":"Sent","opens":"Opens","clicks":"Clicks","call_wa":"Call / WhatsApp",
        "db_note":"Reaching our database of 86,000+ energy professionals","sec_orgs":"AUDIENCE",
        "orgs_title":"Top job titles","orgs_sub":"The most senior roles among the live audience",
        "sec_contact":"LET'S TALK","contact_title":"Shall we talk about your next event?",
        "contact_sub":"For questions, comments or partnership proposals, reach out directly.",
        "watch_rec":"Watch the recording & download materials",
        "seg":{"Developers & IPPs":"Developers & IPPs","Utilities & Energy":"Utilities & Energy",
               "Finance & Advisory":"Finance & Advisory","EPC, Tech & Equipment":"EPC, Tech & Equipment"}},
 "es": {"report":"REPORTE DE MARKETING Y AUDIENCIA","partner":"En colaboración con","moderator":"Moderadora",
        "sec_facts":"DATOS CLAVE","facts_title":"El evento en cifras","by_country":"Por país",
        "facts_ctx":["profesionales inscritos","organizaciones únicas","en todo el mundo","en Zoom + YouTube","de los registrados","experta + moderadora"],
        "registrations":"Registros",
        "companies":"Empresas","countries":"Países","live":"Espectadores totales","attendance":"Tasa de asistencia",
        "speakers":"Ponentes","sec_aud":"AUDIENCIA","country_title":"Desde dónde se conectaron","by_region":"Por región",
        "countries_rep":"países representados","ind_title":"Quién se registró","ind_note":"otros / sin especificar",
        "sec_comp":"EMPRESAS","comp_title":"Quién se registró","sel_comp":"Empresas registradas destacadas",
        "sec_eng":"ENGAGEMENT","eng_title":"Cómo lo siguieron","ge30":"permaneció 30 min o más","avg_conn":"Conexión media",
        "peak":"Pico simultáneo","top_country":"País principal","live_countries":"países en directo","minutes":"min",
        "retention":"Tiempo conectado (Zoom)","retention_fun":"¡Casi todos se quedaron!","watch_yt":"Vuelve a verlo en YouTube","yt_pending":"visualizaciones — por confirmar","views":"visualizaciones",
        "sec_reach":"ALCANCE DE MARKETING","reach_title":"Cómo llenamos la sala","delivered":"emails enviados a nuestra base","opened":"abrieron el email",
        "clicked":"hicieron clic en uno o más enlaces del email","campaign":"Campaña","sent":"Enviados","opens":"Aperturas","clicks":"Clics","call_wa":"Llamada / WhatsApp",
        "db_note":"Llegando a nuestra base de 86.000+ profesionales del sector","sec_orgs":"AUDIENCIA",
        "orgs_title":"Principales cargos","orgs_sub":"Los cargos más senior de la audiencia en directo",
        "sec_contact":"HABLEMOS","contact_title":"¿Hablamos de tu próximo evento?",
        "contact_sub":"Para dudas, comentarios o propuestas de colaboración, contáctanos directamente.",
        "watch_rec":"Ve la grabación y descarga los materiales",
        "seg":{"Developers & IPPs":"Desarrolladores e IPP","Utilities & Energy":"Utilities y energía",
               "Finance & Advisory":"Finanzas y asesoría","EPC, Tech & Equipment":"EPC, tecnología y equipo"}},
 "it": {"report":"REPORT DI MARKETING E AUDIENCE","partner":"In collaborazione con","moderator":"Moderatrice",
        "sec_facts":"DATI CHIAVE","facts_title":"L'evento in numeri","by_country":"Per paese",
        "facts_ctx":["professionisti iscritti","organizzazioni uniche","in tutto il mondo","su Zoom + YouTube","dei registrati","esperto + moderatrice"],
        "registrations":"Registrazioni",
        "companies":"Aziende","countries":"Paesi","live":"Spettatori totali","attendance":"Tasso di partecipazione",
        "speakers":"Relatori","sec_aud":"AUDIENCE","country_title":"Da dove si sono collegati","by_region":"Per regione",
        "countries_rep":"paesi rappresentati","ind_title":"Chi si è registrato","ind_note":"altri / non specificato",
        "sec_comp":"AZIENDE","comp_title":"Chi si è registrato","sel_comp":"Aziende registrate selezionate",
        "sec_eng":"ENGAGEMENT","eng_title":"Come hanno seguito","ge30":"è rimasto 30 min o più","avg_conn":"Connessione media",
        "peak":"Picco simultaneo","top_country":"Paese principale","live_countries":"paesi live","minutes":"min",
        "retention":"Tempo connesso (Zoom)","retention_fun":"Quasi tutti sono rimasti!","watch_yt":"Rivedi su YouTube","yt_pending":"visualizzazioni — da confermare","views":"visualizzazioni",
        "sec_reach":"REACH DI MARKETING","reach_title":"Come abbiamo riempito la sala","delivered":"email inviate al nostro database","opened":"hanno aperto l'email",
        "clicked":"hanno cliccato uno o più link dell'email","campaign":"Campagna","sent":"Inviate","opens":"Aperture","clicks":"Clic","call_wa":"Chiamata / WhatsApp",
        "db_note":"Raggiungendo il nostro database di 86.000+ professionisti del settore","sec_orgs":"AUDIENCE",
        "orgs_title":"Principali ruoli","orgs_sub":"I ruoli più senior tra il pubblico live",
        "sec_contact":"PARLIAMONE","contact_title":"Parliamo del tuo prossimo evento?",
        "contact_sub":"Per domande, commenti o proposte di collaborazione, contattaci direttamente.",
        "watch_rec":"Guarda la registrazione e scarica i materiali",
        "seg":{"Developers & IPPs":"Developer e IPP","Utilities & Energy":"Utility ed energia",
               "Finance & Advisory":"Finanza e consulenza","EPC, Tech & Equipment":"EPC, tech e forniture"}},
 "pl": {"report":"RAPORT MARKETINGOWY I AUDYTORIUM","partner":"We współpracy z","moderator":"Moderatorka",
        "sec_facts":"KLUCZOWE DANE","facts_title":"Wydarzenie w liczbach","by_country":"Według kraju",
        "facts_ctx":["zarejestrowanych profesjonalistów","unikalnych organizacji","na całym świecie","na Zoom + YouTube","zarejestrowanych","ekspert + moderatorka"],
        "registrations":"Rejestracje",
        "companies":"Firmy","countries":"Kraje","live":"Łączna widownia","attendance":"Frekwencja",
        "speakers":"Prelegenci","sec_aud":"AUDYTORIUM","country_title":"Skąd dołączyli","by_region":"Według regionu",
        "countries_rep":"reprezentowanych krajów","ind_title":"Kto się zarejestrował","ind_note":"inne / nieokreślone",
        "sec_comp":"FIRMY","comp_title":"Kto się zarejestrował","sel_comp":"Wybrane zarejestrowane firmy",
        "sec_eng":"ZAANGAŻOWANIE","eng_title":"Jak oglądali","ge30":"pozostało 30 min lub dłużej","avg_conn":"Śr. połączenie",
        "peak":"Szczyt jednoczesny","top_country":"Główny kraj","live_countries":"krajów na żywo","minutes":"min",
        "retention":"Czas połączenia (Zoom)","retention_fun":"Prawie wszyscy zostali!","watch_yt":"Obejrzyj ponownie na YouTube","yt_pending":"wyświetleń — do potwierdzenia","views":"wyświetleń",
        "sec_reach":"ZASIĘG MARKETINGOWY","reach_title":"Jak zapełniliśmy salę","delivered":"e-maili wysłanych do naszej bazy","opened":"otworzyło e-mail",
        "clicked":"kliknęło co najmniej jeden link w e-mailu","campaign":"Kampania","sent":"Wysłane","opens":"Otwarcia","clicks":"Kliknięcia","call_wa":"Telefon / WhatsApp",
        "db_note":"Docierając do bazy 86 000+ specjalistów branży","sec_orgs":"PUBLICZNOŚĆ",
        "orgs_title":"Najczęstsze stanowiska","orgs_sub":"Najbardziej senioralne stanowiska wśród publiczności na żywo",
        "sec_contact":"POROZMAWIAJMY","contact_title":"Porozmawiamy o Twoim kolejnym wydarzeniu?",
        "contact_sub":"W razie pytań lub propozycji współpracy skontaktuj się z nami bezpośrednio.",
        "watch_rec":"Obejrzyj nagranie i pobierz materiały",
        "seg":{"Developers & IPPs":"Deweloperzy i IPP","Utilities & Energy":"Utilities i energia",
               "Finance & Advisory":"Finanse i doradztwo","EPC, Tech & Equipment":"EPC, technologia i sprzęt"}},
}

# ---------------------------------------------------------------- webinar meta + manual data
WEBINAR = {
 "title_lines": ["BESS Procurement", "& QAQC"],
 "subtitle": "Challenges and best practices to maximise ROI",
 "date": "28 May 2026",
 "speakers": [
    {"name":"Lluís Millet Biosca","role":"Senior Consultant, Energy Storage","org":"One Hub Energy","photo":"speaker_lluis.png","mod":False},
    {"name":"Cristina Galán","role":"Content & Event Manager","org":"ATA Insights","photo":"speaker_cristina.png","mod":True},
 ],
 "sponsor_logo":"sponsor_onehub.png",
 "youtube":"youtube.com/live/kYhS0sBKhaY",
 "youtube_url":"https://youtube.com/live/kYhS0sBKhaY",
 "youtube_views":"171",   # auto-fetched via yt-dlp (extract_info view_count)
 "contact": {"name":"Cintia Hernández","role":"Business Development","email":"cintia.hernandez@ata.email"},
 "email": {"delivered":"410K+","opened":"106K+","clicked":"42K+",
    "rows":[("E-shot 1 · Apr 15","11,081","3,783","400"),
            ("E-shot 2 (resend) · Apr 29","7,350","795","98"),
            ("Weekly Spanish ×7","253,786","57,966","22,227"),
            ("Weekly English ×7","138,648","44,028","19,827")]},
 # company showcase grouped by segment (real registrants)
 "segments": [
    ("Developers & IPPs", ["Acciona Energía","Zelestra","Sonnedix","RP-Global","OPDEnergy","Elawan","Q Energy","Grenergy"]),
    ("Utilities & Energy", ["Repsol","Naturgy","Engie","EDP","Galp","TotalEnergies","Voltalia","Alpiq"]),
    ("Finance & Advisory", ["Abanca","BEC Capital","Alantra","Infranity","Fichtner","KPMG","PwC","Quintas Energy"]),
    ("EPC, Tech & Equipment", ["ABB","Hitachi Energy","Huawei","NHOA Energy","Nidec","Elecnor","Ferrovial","Wood"]),
 ],
 # auto-shortlist targets for "most relevant organizations" (display, match keywords)
 "org_targets": [
    ("Repsol",["repsol"]),("Naturgy",["naturgy"]),("Engie",["engie"]),("TotalEnergies",["totalenergies"]),
    ("Hitachi Energy",["hitachienergy"]),("Huawei",["huawei"]),("NHOA Energy",["nhoa.energy"]),
    ("Zelestra",["zelestra"]),("Sonnedix",["sonnedix"]),("RP-Global",["rp-global"]),
    ("Q Energy",["qenergy"]),("Acciona",["acciona"]),("Galp",["galp"]),("Quintas Energy",["quintasenergy"]),
    ("Lightsource bp",["lightsourcebp"]),("ABB",["@it.abb","@abb"]),
 ],
}

# ---------------------------------------------------------------- draw helpers
def canvas(bg=GREY):
    img = Image.new("RGB", (W*SCALE, H*SCALE), bg); return img, ImageDraw.Draw(img)
def T(d, xy, text, font, fill, anchor="la", spacing=4):
    # Register as an EDITABLE native PowerPoint text box (not baked into the PNG).
    # Painted back onto the preview PNG later so the in-app preview is identical.
    if text is None:
        return
    s = text if isinstance(text, str) else str(text)
    if s == "":
        return
    NATIVE_NUMS.append({"slide": _CUR_SLIDE, "x": xy[0], "y": xy[1], "px": _px_of(font),
                        "text": s, "rgb": tuple(fill), "anchor": anchor,
                        "kind": _kind_of(font), "url": None, "spacing": spacing})

def _draw_text(d, xy, text, font, fill, anchor="la", spacing=4):
    """Bake text straight onto the supersampled canvas (used only for the preview)."""
    d.text((xy[0]*SCALE, xy[1]*SCALE), text, font=font, fill=fill, anchor=anchor, spacing=spacing*SCALE)
def tw(d, text, font):
    b = d.textbbox((0,0), text, font=font); return (b[2]-b[0])/SCALE
def rrect(d, box, radius, fill=None, outline=None, width=1):
    d.rounded_rectangle([c*SCALE for c in box], radius=radius*SCALE, fill=fill, outline=outline, width=width*SCALE)
def rect(d, box, fill):
    d.rectangle([c*SCALE for c in box], fill=fill)
def card(d, box, radius=18, fill=WHITE, edge=ORANGE, edge_w=9):
    rrect(d, box, radius, fill=fill)
    x0,y0,x1,y1 = box
    rrect(d, [x0, y0, x0+edge_w+radius, y1], radius, fill=edge)
    rrect(d, [x0+edge_w, y0, x0+edge_w+radius, y1], radius, fill=fill)
def square_crop(path, bias_y=0.30):
    """Center-crop a photo to a square (biased toward the top to keep the
    forehead) and save a _sq.png — used as an editable circle photo in the PPTX."""
    p = Image.open(path).convert("RGB")
    w,h = p.size; s = min(w,h)
    left = (w-s)//2; top = int((h-s)*bias_y)
    out = os.path.splitext(path)[0] + "_sq.png"
    p.crop((left, top, left+s, top+s)).save(out)
    return out

def clean_avatar(path, crop_bottom=0.24, bias_y=0.16):
    """Remove the logo band baked into the bottom of the speaker card, then
    square-crop the photo so it FILLS the circle (person + studio background,
    no white floating). Cached as _avatar.png."""
    out = os.path.splitext(path)[0] + "_avatar.png"
    if os.path.exists(out):
        return out
    p = Image.open(path).convert("RGB")
    w,h = p.size
    p = p.crop((0, 0, w, int(h*(1-crop_bottom))))      # drop the logo band
    w,h = p.size; sq = min(w,h)
    left = (w-sq)//2; top = int((h-sq)*bias_y)
    p.crop((left, top, left+sq, top+sq)).save(out)
    return out
def paste_cover(img, path, box, centering=(0.5,0.08)):
    try: p = Image.open(path).convert("RGBA")
    except Exception: return
    bw,bh = int((box[2]-box[0])*SCALE), int((box[3]-box[1])*SCALE)
    p = ImageOps.fit(p, (bw,bh), Image.LANCZOS, centering=centering)
    img.paste(p, (int(box[0]*SCALE), int(box[1]*SCALE)), p)
def paste_portrait(img, path, box, align="center"):
    try: p = Image.open(path).convert("RGBA")
    except Exception: return
    bw,bh = (box[2]-box[0])*SCALE, (box[3]-box[1])*SCALE
    s = min(bw/p.size[0], bh/p.size[1])
    p = p.resize((max(1,int(p.size[0]*s)), max(1,int(p.size[1]*s))), Image.LANCZOS)
    x = int(box[0]*SCALE + (bw-p.width)//2)
    y = int(box[1]*SCALE + (bh-p.height)) if align=="bottom" else int(box[1]*SCALE + (bh-p.height)//2)
    img.paste(p, (x,y), p)
def paste_logo(img, path, box):
    try: logo = Image.open(path).convert("RGBA")
    except Exception: return
    bw,bh = (box[2]-box[0])*SCALE, (box[3]-box[1])*SCALE
    s = min(bw/logo.size[0], bh/logo.size[1])
    logo = logo.resize((max(1,int(logo.size[0]*s)), max(1,int(logo.size[1]*s))), Image.LANCZOS)
    img.paste(logo, (int(box[0]*SCALE+(bw-logo.width)//2), int(box[1]*SCALE+(bh-logo.height)//2)), logo)
def circle_photo(img, path, cx, cy, d, ring=WHITE, ring_w=8):
    cx,cy,d,rw = cx*SCALE, cy*SCALE, d*SCALE, ring_w*SCALE
    dr = ImageDraw.Draw(img)
    try: ph = Image.open(path).convert("RGB")
    except Exception:
        dr.ellipse([cx-d//2,cy-d//2,cx+d//2,cy+d//2], fill=TRACK); return
    ph = ImageOps.fit(ph, (d,d), Image.LANCZOS, centering=(0.5,0.4))
    mask = Image.new("L",(d,d),0); ImageDraw.Draw(mask).ellipse([0,0,d,d],fill=255)
    dr.ellipse([cx-d//2-rw,cy-d//2-rw,cx+d//2+rw,cy+d//2+rw], fill=ring)
    img.paste(ph,(cx-d//2,cy-d//2),mask)
def hbars(d, rows, x, y, w, row_h, gap, label_w, maxval, fmt, colors=None, lab_size=15):
    bx = x+label_w
    for i,(lab,val) in enumerate(rows):
        yy = y + i*(row_h+gap)
        T(d,(x,yy+row_h/2),lab,reg(lab_size),CHAR,anchor="lm")
        rrect(d,[bx,yy,bx+w,yy+row_h],row_h//2,fill=TRACK)
        bw = max(row_h, w*(val/maxval)) if maxval else row_h
        col = colors[i%len(colors)] if colors else (ORANGE if i==0 else CORAL)
        rrect(d,[bx,yy,bx+bw,yy+row_h],row_h//2,fill=col)
        T(d,(bx+bw+12,yy+row_h/2),fmt(val),bold(14),CHAR,anchor="lm")
def donut(img, d, cx, cy, rad, thick, segments, hole=GREY):
    cx,cy,rad,th = cx*SCALE,cy*SCALE,rad*SCALE,thick*SCALE
    start=-90.0; total=sum(v for _,v,_ in segments) or 1
    for _,v,col in segments:
        end=start+360*v/total; d.pieslice([cx-rad,cy-rad,cx+rad,cy+rad],start,end,fill=col); start=end
    d.ellipse([cx-(rad-th),cy-(rad-th),cx+(rad-th),cy+(rad-th)],fill=hole)
def footer(d, page, lang):
    rect(d,[0,H-46,W,H],CHAR)
    T(d,(70,H-23),"ATA INSIGHTS",bold(13),WHITE,anchor="lm")
    T(d,(W/2,H-23),STRINGS[lang]["report"],reg(12),(205,205,205),anchor="mm")
    T(d,(W-70,H-23),f"{page:02d}",bold(13),ORANGE,anchor="rm")
def kicker(d, x, y, tag, title):
    T(d,(x,y),tag,bold(16),ORANGE,anchor="la")
    rect(d,[x,y+30,x+150,y+34],CHAR)
    T(d,(x,y+58),title,black(40),CHAR,anchor="la")
def save(img, name):
    if SUFFIX:
        base,ext=os.path.splitext(name); name=base+SUFFIX+ext
    img.resize((W,H), Image.LANCZOS).save(os.path.join(OUT,name),"PNG"); print("saved",name)
    return os.path.join(OUT,name)

# ---------------------------------------------------------------- icons & playful bits
def _ifont(size): return ImageFont.truetype(ICON_FONT_PATH, int(size*SCALE))
def icon(d, cx, cy, key, size, color):
    d.text((cx*SCALE, cy*SCALE), chr(int(ICON[key],16)), font=_ifont(size), fill=color, anchor="mm")
def badge(d, cx, cy, r, bg, key, icol, isize):
    d.ellipse([(cx-r)*SCALE,(cy-r)*SCALE,(cx+r)*SCALE,(cy+r)*SCALE], fill=bg)
    icon(d, cx, cy, key, isize, icol)
def star(d, cx, cy, size, color):
    icon(d, cx, cy, "star", size, color)
def halftone(d, x0, y0, x1, y1, color, step=30, r=4):
    yy = y0
    while yy < y1:
        xx = x0
        while xx < x1:
            d.ellipse([(xx-r)*SCALE,(yy-r)*SCALE,(xx+r)*SCALE,(yy+r)*SCALE], fill=color)
            xx += step
        yy += step
def callout(d, x, y, text, key, bg=YELLOW, fg=CHAR):
    w = tw(d, text, bold(15)) + 92
    rrect(d, [x, y, x+w, y+52], 26, fill=bg)
    icon(d, x+28, y+26, key, 22, fg)
    T(d, (x+50, y+26), text, bold(15), fg, anchor="lm")

# --- "read the text if you want" highlight notes (annotated version only) ---
INSIGHTS = {
 "facts": "602 registrations from 490 companies in 65 countries — and a 27% live show-up, above the 20-25% norm.",
 "country": "Europe-led at 75% with a strong 14% from LATAM — 65 countries point to export-ready, cross-border demand.",
 "industries": "Mostly developers, IPPs and EPC contractors — the people who actually buy and build BESS, not just analysts.",
 "companies": "Tier-1 developers, utilities, financiers and EPC / tech suppliers — direct reach into the full BESS value chain.",
 "engagement": "77% stayed 30+ minutes and watched 52 of 60 on average — a high-intent, sticky audience, not casual drop-ins.",
 "reach": "410K+ deliveries and 106K+ opens put the sponsor's topic in front of our 86K database many times over.",
 "orgs": "Senior, technically-relevant roles from leading energy companies watched live — the contacts a sponsor wants.",
}
def ctop(clean_top):
    # in annotated mode, push slide content down to a common top so the
    # highlight line gets a balanced band under the title
    return 272 if ANNOTATE else clean_top
def draw_insight(d, key):
    if not ANNOTATE or key not in INSIGHTS: return
    text = INSIGHTS[key]; y = 222
    fs = 23
    while fs > 18 and tw(d, text, reg(fs)) > (1850-118): fs -= 1
    star(d, 86, y, fs+2, ORANGE)
    T(d, (112, y), text, reg(fs), (38,44,48), anchor="lm")

# ---------------------------------------------------------------- slides
def chips(d, names, cx0, cy0, maxx, max_y, fs=22, pad=52, h=58, vstep=74, maxlen=24):
    """Lay out rounded name-chips; truncates long names, never wider than the card,
    stops when out of vertical room (so nothing overflows the card)."""
    f = semi(fs); cx, cy = cx0, cy0
    for nm in names:
        nm = (nm[:maxlen].rstrip() + "…") if len(nm) > maxlen + 1 else nm
        wpx = min(tw(d, nm, f) + pad, maxx - cx0)
        if cx + wpx > maxx:
            cx = cx0; cy += vstep
        if cy + h > max_y:
            break
        rrect(d, [cx, cy, cx + wpx, cy + h], 12, fill=GREY)
        T(d, (cx + wpx / 2, cy + h / 2), nm, f, CHAR, anchor="mm")
        cx += wpx + 16

def wrap_lines(d, text, font, maxw, max_lines=2):
    words = (text or "").split(); lines = []; cur = ""
    for w in words:
        t = (cur + " " + w).strip()
        if not cur or tw(d, t, font) <= maxw:
            cur = t
        else:
            lines.append(cur); cur = w
    if cur:
        lines.append(cur)
    return lines[:max_lines]

def slide_cover(st, lang):
    _slide(0); s=STRINGS[lang]; img,d=canvas(GREY)
    PX=1170
    rect(d,[PX,0,W,H],ORANGE)              # right orange block
    rect(d,[PX-7,0,PX,H],CHAR)             # charcoal divider
    halftone(d, PX+40, 60, PX+360, 240, (255,120,70), step=30, r=4)   # pop-art texture on orange
    paste_logo(img,ATA_LOGO,[70,64,360,152])
    rect(d,[72,166,240,172],ORANGE)
    T(d,(72,200),s["report"],bold(22),ORANGE,anchor="la")
    star(d, 72+tw(d,s["report"],bold(22))+40, 196, 26, YELLOW)
    avail = PX - 70 - 50
    full = " ".join([l for l in WEBINAR["title_lines"] if l]).strip()
    # Fit the title at the LARGEST font that stays within `avail`, using as FEW
    # lines as possible — but add lines (up to 4) when a long title would otherwise
    # be forced down to a tiny font. (Old behaviour capped at 2 lines -> tiny fonts.)
    FLOOR, GOOD, CAP, MAXL = 46, 72, 112, 5
    lines, fs = [full], CAP
    for L in range(1, MAXL + 1):        # fewest lines that keep the font >= GOOD (>=36pt)
        f = CAP
        while f > FLOOR:
            wl = wrap_lines(d, full, black(f), avail, max_lines=99)
            if len(wl) <= L and all(tw(d, ln, black(f)) <= avail for ln in wl):
                break
            f -= 2
        lines, fs = wrap_lines(d, full, black(f), avail, max_lines=99), f
        if f >= GOOD or L == MAXL:
            break
    y0 = {1: 300, 2: 272, 3: 250}.get(len(lines), 228)
    pitch = int(fs * 1.12)                         # uniform leading for the stacked lines
    for i, ln in enumerate(lines):
        ly = y0 + i * pitch
        NATIVE_NUMS.append({"slide": 0, "x": 70, "y": ly, "px": fs, "text": ln,
                            "rgb": tuple(CHAR), "anchor": "la", "kind": "black",
                            "url": None, "spacing": 4, "vbox": (ly - int(fs * 0.16), pitch)})
    yy = y0 + len(lines) * pitch                    # subtitle/date flow below the title
    sub = WEBINAR.get("subtitle") or ""
    sfs = 32
    while sfs > 18 and tw(d, sub, semi(sfs)) > avail: sfs -= 1
    suby = yy + 24
    T(d,(72,suby),sub,semi(sfs),(48,54,58),anchor="la")
    dy = suby + (sfs + 44 if sub else 6)
    rrect(d,[72,dy,72+tw(d,WEBINAR['date'],bold(20))+56,dy+52],10,fill=CHAR)
    T(d,(100,dy+26),WEBINAR["date"],bold(20),WHITE,anchor="lm")
    global COVER_PHOTOS; COVER_PHOTOS=[]
    sps=WEBINAR["speakers"]; n=max(1,len(sps)); sx=1548
    has_sp=bool(WEBINAR.get("sponsor_logo"))
    top=110; bottom=(885 if has_sp else 1015); block=(bottom-top)/n   # adaptive 1–4 speakers
    D=int(min(248, block*0.52)); name_fs=max(15,int(D*0.115)); role_fs=max(12,int(D*0.083))
    for i,sp in enumerate(sps):
        cy=int(top+block*i+D/2+10)
        # white ring + light disc placeholder; the photo is an EDITABLE native picture in the PPTX
        d.ellipse([(sx-D//2-8)*SCALE,(cy-D//2-8)*SCALE,(sx+D//2+8)*SCALE,(cy+D//2+8)*SCALE],fill=WHITE)
        d.ellipse([(sx-D//2)*SCALE,(cy-D//2)*SCALE,(sx+D//2)*SCALE,(cy+D//2)*SCALE],fill=(228,229,231))
        ph=sp.get("photo")
        if ph:
            COVER_PHOTOS.append({"slide":0,"path":clean_avatar(os.path.join(ASSETS,ph)),
                                 "cx":sx,"cy":cy,"d":D})
        ny=cy+D//2+name_fs+4
        T(d,(sx,ny),sp["name"],bold(name_fs),WHITE,anchor="mm")
        sub2=(sp.get("org") or sp.get("role") or "")
        role=(s["moderator"]+" · " if sp.get("mod") else "")+sub2
        T(d,(sx,ny+role_fs+10),role,reg(role_fs),(255,224,210),anchor="mm")
    sp_logo=WEBINAR.get("sponsor_logo")
    if sp_logo and os.path.exists(os.path.join(ASSETS,sp_logo)):
        rrect(d,[1330,H-150,1850,H-62],12,fill=WHITE)
        T(d,(1356,H-122),s["partner"],reg(13),MUTE,anchor="lm")
        paste_logo(img,os.path.join(ASSETS,sp_logo),[1540,H-138,1828,H-74])
    return save(img,f"01_cover_{lang}.png")

def hero_tile(d, box, num, label, fill=ORANGE, numc=WHITE, labc=(255,224,210)):
    rrect(d,box,18,fill=fill)
    T(d,(box[0]+50,box[1]+44),num,black(72),numc,anchor="la")
    T(d,(box[0]+54,box[1]+162),label,semi(22),labc,anchor="la")

def slide_facts(st, lang):
    _slide(1); s=STRINGS[lang]; img,d=canvas(GREY)
    if not ANNOTATE: halftone(d, 1570, 66, 1885, 210, (220,221,223), step=30, r=4)
    kicker(d,70,64,"01 · "+s["sec_facts"],s["facts_title"])
    star(d, tw(d,s["facts_title"],black(40))+118, 96, 30, YELLOW)
    draw_insight(d,"facts")
    kf=st["key_facts"]; ctx=s["facts_ctx"]
    cw,gx,gy=573,30,34; x0=70; y0=ctop(248)
    ch=340 if not ANNOTATE else int((H-90 - y0 - gy)//2)
    # (num, label, badge-bg, icon, icon-colour, is_hero)
    cells=[(f"{kf['registrations']:,}".replace(',', '.'),s["registrations"],"people",WHITE,ORANGE,True),
           (f"{kf['companies']}",s["companies"],"org",CORAL,WHITE,False),
           (f"{kf['countries']}",s["countries"],"globe",TEAL,WHITE,False),
           (f"{kf['live_attendees']}",s["live"],"video",YELLOW,CHAR,False),
           (f"{kf['attendance_rate_pct']:.0f}%",s["attendance"],"gauge",GREEN,WHITE,False),
           ("2",s["speakers"],"mic",LIGHT,WHITE,False)]
    pos=[(0,0),(1,0),(2,0),(0,1),(1,1),(2,1)]
    for i,((num,lab,key,bg,icol,hero),(c,r)) in enumerate(zip(cells,pos)):
        cx=x0+c*(cw+gx); cy=y0+r*(ch+gy)
        numc, labc, ctxc = (WHITE,WHITE,(255,224,210)) if hero else (CHAR,CHAR,MUTE)
        if hero: rrect(d,[cx,cy,cx+cw,cy+ch],20,fill=ORANGE); badge(d,cx+108,cy+108,52,WHITE,key,ORANGE,46)
        else:    card(d,[cx,cy,cx+cw,cy+ch],edge=bg);        badge(d,cx+108,cy+108,52,bg,key,icol,46)
        nnum(num,cx+58,cy+166,84,numc,"la",1,"black")
        T(d,(cx+62,cy+266),lab,semi(25),labc,anchor="la")
        T(d,(cx+62,cy+304),ctx[i],reg(16),ctxc,anchor="la")
    footer(d,2,lang); return save(img,f"02_facts_{lang}.png")

def slide_country(st, lang):
    _slide(2); s=STRINGS[lang]; img,d=canvas(GREY)
    kicker(d,70,64,"02 · "+s["sec_aud"],s["country_title"])
    draw_insight(d,"country")
    CY0,CY1=ctop(210),H-90
    # left card: by country (bars vertically centred)
    card(d,[70,CY0,1130,CY1])
    T(d,(122,CY0+40),s["by_country"],bold(22),CHAR,anchor="la")
    rows=[(r["label"],r["pct"]) for r in st["countries"]["rows"][:8]]
    rh,gp=52,22; tot=8*rh+7*gp; by=CY0+118+((CY1-CY0-118)-tot)//2
    hbars(d,rows,122,by,560,rh,gp,210,max(v for _,v in rows),lambda v:f"{v:.1f}%",colors=ACCENTS,lab_size=18)
    # right card: by region (donut centred to match)
    card(d,[1180,CY0,1850,CY1])
    T(d,(1232,CY0+40),s["by_region"],bold(22),CHAR,anchor="la")
    regs=st["regions"][:6]
    dcy=CY0+260
    donut(img,d,1380,dcy,150,46,[(r["label"],r["pct"],RAMP[i%len(RAMP)]) for i,r in enumerate(regs)],hole=WHITE)
    nnum(f"{st['countries'].get('distinct', st['key_facts']['countries'])}",1380,dcy-15,46,CHAR,"mm",2,"black")
    T(d,(1380,dcy+30),s["countries_rep"],reg(13),MUTE,anchor="mm")
    ly=CY0+440
    for i,r in enumerate(regs):
        rect(d,[1235,ly+i*42,1259,ly+i*42+24],RAMP[i%len(RAMP)])
        T(d,(1276,ly+i*42+12),r["label"],reg(17),CHAR,anchor="lm")
        T(d,(1800,ly+i*42+12),f"{r['pct']:.1f}%",bold(17),CHAR,anchor="rm")
    footer(d,3,lang); return save(img,f"03_country_{lang}.png")

def slide_industries(st, lang):
    _slide(3); s=STRINGS[lang]; img,d=canvas(GREY)
    kicker(d,70,64,"02 · "+s["sec_aud"],s["ind_title"])
    draw_insight(d,"industries")
    SHORT={"Equipment Supplier or Manufacturer":"Equipment Supplier",
           "Government Institutions":"Government","International Financial Institution (IFI)":"IFI"}
    CY0,CY1=ctop(210),H-90
    card(d,[70,CY0,1080,CY1])
    rows=[(SHORT.get(r["label"],r["label"]),r["pct"]) for r in st["industries"]["rows"] if r["label"]!="Other"][:8]
    rh,gp=50,22; tot=8*rh+7*gp; by=CY0+((CY1-CY0)-tot)//2 - 8
    hbars(d,rows,122,by,440,rh,gp,360,max(v for _,v in rows),lambda v:f"{v:.1f}%",colors=ACCENTS,lab_size=19)
    other=next((r["pct"] for r in st["industries"]["rows"] if r["label"]=="Other"),0)
    T(d,(122,CY1-42),f"+ {other:.0f}% {s['ind_note']}",reg(15),MUTE,anchor="lm")
    # right card: selected companies (larger chips)
    card(d,[1110,CY0,1850,CY1])
    T(d,(1162,CY0+40),s["sel_comp"],bold(22),CHAR,anchor="la")
    chips(d, WEBINAR.get("comp_sample", [])[:48], 1162, CY0+104, 1816, CY1-16, fs=18, h=50, vstep=60)
    footer(d,4,lang); return save(img,f"04_industries_{lang}.png")

def slide_companies(st, lang):
    s=STRINGS[lang]; img,d=canvas(GREY)
    kicker(d,70,64,"03 · "+s["sec_comp"],s["comp_title"])
    draw_insight(d,"companies")
    cols=2; cw=880; gx=40; gy=32; x0=70; y0=ctop(250)
    ch=362 if not ANNOTATE else int((H-90 - y0 - gy)//2)
    seg_style=[(ORANGE,"people"),(CORAL,"power"),(TEAL,"org"),(GREEN,"brief")]
    for i,(seg,names) in enumerate(WEBINAR["segments"]):
        c=i%cols; r=i//cols; x=x0+c*(cw+gx); y=y0+r*(ch+gy)
        ecol,ikey=seg_style[i%len(seg_style)]
        card(d,[x,y,x+cw,y+ch],edge=ecol)
        badge(d,x+82,y+68,34,ecol,ikey,WHITE,30)
        T(d,(x+134,y+68),s["seg"][seg],bold(23),CHAR,anchor="lm")
        chips(d, names, x+48, y+138, x+cw-40, y+ch-24, fs=22)
    footer(d,5,lang); return save(img,f"05_companies_{lang}.png")

def slide_engagement(st, lang):
    _slide(4); s=STRINGS[lang]; img,d=canvas(GREY); live=st["live"]
    kicker(d,70,64,"03 · "+s["sec_eng"],s["eng_title"])
    draw_insight(d,"engagement")
    lt=ctop(210); Hin=(H-90)-lt
    card(d,[70,lt,760,H-90])
    dcy=lt+int(Hin*0.25)
    donut(img,d,415,dcy,140,44,[("",live["pct_ge_30min"],ORANGE),("",live["pct_lt_30min"],TRACK)],hole=WHITE)
    nnum(f"{live['pct_ge_30min']:.0f}%",415,dcy-18,48,CHAR,"mm",4,"black")
    T(d,(415,dcy+28),"30+ "+s["minutes"],reg(13),MUTE,anchor="mm")
    T(d,(415,lt+int(Hin*0.47)),s["ge30"],semi(18),CHAR,anchor="mm")
    T(d,(120,lt+int(Hin*0.55)),s["retention"],bold(15),CHAR,anchor="la")
    b4=live["buckets_4"]; order=[(">45",b4[">45"]),("31-45",b4["31-45"]),("16-30",b4["16-30"]),("0-15",b4["0-15"])]
    mx=max(v["pct"] for _,v in order)
    hbars(d,[(f"{k} {s['minutes']}",v["pct"]) for k,v in order],120,lt+int(Hin*0.60),300,30,16,140,mx,lambda v:f"{v:.1f}%",
          colors=[ORANGE,CORAL,LIGHT,(255,205,160)])
    callout(d,120,lt+int(Hin*0.87),s["retention_fun"],"like",bg=YELLOW)
    tiles=[(f"{live['avg_minutes']:.0f} {s['minutes']}",s["avg_conn"],"clock",CORAL),
           (f"{live['peak_concurrent']}",s["peak"],"group",TEAL),
           (live.get("top_country") or "—",s["top_country"],"pin",ORANGE),
           (f"{live['attendee_countries_represented']}",s["live_countries"],"globe",GREEN)]
    tx,twd,thd,tg=820,500,205,34; ty=lt
    for i,(num,lab,key,col) in enumerate(tiles):
        x=tx+(i%2)*(twd+30); y=ty+(i//2)*(thd+tg); card(d,[x,y,x+twd,y+thd],edge=col)
        badge(d,x+twd-66,y+62,32,col,key,WHITE,28)
        if i==2: T(d,(x+46,y+44),num,black(50),CHAR,anchor="la")   # top country (text, not editable)
        else:    nnum(num,x+46,y+44,50,CHAR,"la",4,"black")
        T(d,(x+50,y+128),lab,semi(21),MUTE,anchor="la")
    yb=ty+2*(thd+tg)+6
    rrect(d,[820,yb,1850,H-90],16,fill=CHAR)
    cyb=(yb+H-90)//2
    badge(d,902,cyb,40,ORANGE,"play",WHITE,34)
    # the link text + views is added as a NATIVE, EDITABLE text box (with a real
    # clickable hyperlink) in build_pptx — so it works and the views can be typed in.
    yv=WEBINAR.get("youtube_views")
    vline=WEBINAR["youtube"]+"    ·    "+(f"{yv} {s['views']}" if yv else s["views"]+": ____")
    global LINK_INFO
    LINK_INFO={"slide":4,"box":(958,yb+20,1840,H-90-20),"url":WEBINAR["youtube_url"],
               "title":s["watch_yt"],"line2":vline}
    footer(d,5,lang); return save(img,f"06_engagement_{lang}.png")

def slide_reach(st, lang):
    _slide(5); s=STRINGS[lang]; img,d=canvas(GREY); em=WEBINAR["email"]
    kicker(d,70,64,"04 · "+s["sec_reach"],s["reach_title"])
    star(d, tw(d,s["reach_title"],black(40))+110, 96, 28, YELLOW)
    draw_insight(d,"reach")
    cw,gx,x0,ch=570,35,70,200; y0=ctop(230)
    stats=[(em["delivered"],s["delivered"],"send",ORANGE,True),
           (em["opened"],s["opened"],"mail",CORAL,False),
           (em["clicked"],s["clicked"],"click",TEAL,False)]
    for i,(num,lab,key,col,hero) in enumerate(stats):
        x=x0+i*(cw+gx)
        numc, labc = (WHITE, (255,224,210)) if hero else (CHAR, MUTE)
        if hero: rrect(d,[x,y0,x+cw,y0+ch],18,fill=ORANGE); badge(d,x+cw-70,y0+66,38,WHITE,key,ORANGE,32)
        else:    card(d,[x,y0,x+cw,y0+ch],edge=col);        badge(d,x+cw-70,y0+66,38,col,key,WHITE,32)
        nnum(num,x+50,y0+44,64,numc,"la",5,"black")
        lns=wrap_lines(d,lab,semi(18),cw-90,2)
        for k,ln in enumerate(lns):
            T(d,(x+54,y0+148+k*26),ln,semi(18),labc,anchor="la")
    ty=470+(y0-230); card(d,[70,ty,1850,ty+430])
    cols=[130,1040,1320,1600]
    for j,lbl in enumerate([s["campaign"],s["sent"],s["opens"],s["clicks"]]):
        T(d,(cols[j],ty+56),lbl,bold(21),MUTE,anchor="lm" if j==0 else "rm")
    rect(d,[120,ty+92,1800,ty+95],LINE)
    for r,(name,sent,op,cl) in enumerate(em["rows"]):
        ry=ty+150+r*66
        T(d,(cols[0],ry),name,semi(21),CHAR,anchor="lm")
        for j,val in zip(range(1,4),(sent,op,cl)): nnum(val,cols[j],ry,21,CHAR,"rm",5,"reg")
    icon(d,82,ty+462,"people",18,MUTE)
    T(d,(106,ty+462),s["db_note"],reg(16),MUTE,anchor="lm")
    footer(d,6,lang); return save(img,f"07_reach_{lang}.png")

def slide_titles(st, lang):
    _slide(6); s=STRINGS[lang]; img,d=canvas(GREY)
    kicker(d,70,64,"05 · "+s["sec_orgs"],s["orgs_title"])
    draw_insight(d,"orgs")
    if not ANNOTATE: T(d,(70,212),s["orgs_sub"],reg(18),MUTE,anchor="la")
    titles=WEBINAR.get("job_titles",[])[:14]
    cols=2; cw=870; gx=30; ch=84; gy=14; x0=70; y0=296
    for i,tt in enumerate(titles):
        c=i%cols; r=i//cols; x=x0+c*(cw+gx); y=y0+r*(ch+gy)
        if y+ch>H-86: break
        card(d,[x,y,x+cw,y+ch],edge=ORANGE if i%2==0 else CORAL)
        ttl=(tt[:46].rstrip()+"…") if len(tt)>47 else tt
        T(d,(x+46,y+ch//2),ttl,bold(24),CHAR,anchor="lm")
    footer(d,7,lang); return save(img,f"08_orgs_{lang}.png")

def slide_contact(st, lang):
    _slide(7); s=STRINGS[lang]; img,d=canvas(DEEP); ct=WEBINAR["contact"]
    rect(d,[0,0,18,H],ORANGE)
    RW=1180
    rect(d,[RW,0,W,H],WHITE)                                   # white right panel
    halftone(d, RW+30, 70, RW+300, 300, (255,210,185), step=30, r=6)     # halftone accent
    star(d, W-150, 120, 66, ORANGE); star(d, W-64, 220, 38, YELLOW)
    paste_cover(img, CONTACT_PHOTO, [RW+20, 150, W, H], centering=(0.5,0.04))
    # left content
    T(d,(120,150),s["sec_contact"],bold(18),ORANGE,anchor="la")
    rect(d,[122,184,272,188],ORANGE)
    T(d,(120,212),s["contact_title"],black(50),WHITE,anchor="la")
    T(d,(122,320),s["contact_sub"],reg(20),(205,200,196),anchor="la")
    rows=[("people",ct["name"],ct["role"]+" · ATA Insights"),
          ("phone",ct.get("phone",""),s.get("call_wa","")),
          ("mail",ct["email"],""),
          ("globe","atainsights.com · my.atainsights.com","")]
    rows=[r for r in rows if r[1]]
    for i,(key,a,b) in enumerate(rows):
        yy=452+i*90
        badge(d,152,yy+14,32,ORANGE,key,WHITE,28)
        if key=="mail":   # clickable mailto: native, editable text box
            nnum(a,214,yy,24,WHITE,"la",7,"bold",url="mailto:"+a)
        else:
            T(d,(214,yy),a,bold(24),WHITE,anchor="la")
        if b: T(d,(214,yy+36),b,reg(19),(205,200,196),anchor="la")
    # sponsor co-brand lockup (only if a sponsor logo is present)
    sp_logo=WEBINAR.get("sponsor_logo")
    if sp_logo and os.path.exists(os.path.join(ASSETS,sp_logo)):
        rrect(d,[120,812,672,906],12,fill=WHITE)
        T(d,(150,859),s["partner"],reg(15),MUTE,anchor="lm")
        paste_logo(img,os.path.join(ASSETS,sp_logo),[150+tw(d,s["partner"],reg(15))+28,828,660,890])
    T(d,(120,H-60),"ATA INSIGHTS",bold(16),ORANGE,anchor="lm")
    return save(img,f"09_contact_{lang}.png")

# ---------------------------------------------------------------- curation
def relevant_orgs(regs, targets):
    out=[]
    for disp, keys in targets:
        for r in regs:
            hay=(r.company+" "+r.email).lower()
            if any(k in hay for k in keys) and r.name and r.job_title:
                out.append((disp, r.name, r.job_title)); break
    return out

# ---------------------------------------------------------------- run
def build_pptx(paths, name, photos=None, link=None, nums=None):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.oxml.ns import qn
        from pptx.enum.text import MSO_ANCHOR
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except Exception as e:
        print("pptx skip:",e); return
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    PXIN=144.0   # render is 1920x1080 over 13.333x7.5in => 144 px/in
    sl=[]
    for p in paths:
        s=prs.slides.add_slide(prs.slide_layouts[6])
        s.shapes.add_picture(p,0,0,width=prs.slide_width,height=prs.slide_height); sl.append(s)
    # editable circle-cropped speaker photos
    for ph in (photos or []):
        s=sl[ph["slide"]]; D=ph["d"]
        pic=s.shapes.add_picture(ph["path"],Inches((ph["cx"]-D/2)/PXIN),Inches((ph["cy"]-D/2)/PXIN),
                                 Inches(D/PXIN),Inches(D/PXIN))
        spPr=pic._element.spPr      # crop-to-circle (still freely editable in PowerPoint)
        geom=spPr.find(qn('a:prstGeom'))
        if geom is not None:
            geom.set('prst','ellipse')
        else:
            g=spPr.makeelement(qn('a:prstGeom'),{'prst':'ellipse'}); g.append(spPr.makeelement(qn('a:avLst'),{})); spPr.append(g)
    # native, editable YouTube text box: clickable hyperlink + type-in views
    if link:
        s=sl[link["slide"]]; x0,y0,x1,y1=link["box"]
        tb=s.shapes.add_textbox(Inches(x0/PXIN),Inches(y0/PXIN),Inches((x1-x0)/PXIN),Inches((y1-y0)/PXIN))
        tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        tf.margin_left=0; tf.margin_top=0; tf.margin_right=0; tf.margin_bottom=0
        r1=tf.paragraphs[0].add_run(); r1.text=link["title"]
        r1.hyperlink.address=link["url"]                 # set link first…
        r1.font.size=Pt(15); r1.font.bold=True; r1.font.name="Montserrat"
        r1.font.underline=False
        r1.font.color.rgb=RGBColor(255,255,255)          # …then force white over theme-blue
        # belt-and-braces: drop any theme hyperlink colour on the run
        _rPr=r1._r.get_or_add_rPr()
        for _u in _rPr.findall(qn('a:uFill'))+_rPr.findall(qn('a:uLn')): _rPr.remove(_u)
        r2=tf.add_paragraph().add_run(); r2.text=link["line2"]
        r2.font.size=Pt(11); r2.font.name="Inter"; r2.font.color.rgb=RGBColor(255,200,180)
        # force the THEME hyperlink colour to white (run-level colour gets overridden by theme-blue)
        try:
            from pptx.opc.constants import RELATIONSHIP_TYPE as RT
            from lxml import etree
            NS='{http://schemas.openxmlformats.org/drawingml/2006/main}'
            tpart=prs.slide_masters[0].part.part_related_by(RT.THEME)
            root=etree.fromstring(tpart.blob); cs=root.find('.//%sclrScheme'%NS)
            for tag in ('hlink','folHlink'):
                el=cs.find('%s%s'%(NS,tag))
                if el is not None:
                    for ch in list(el): el.remove(ch)
                    etree.SubElement(el,'%ssrgbClr'%NS).set('val','FFFFFF')
            tpart._blob=etree.tostring(root)
        except Exception as e:
            print('theme hlink skip',e)
    # native, editable figures (numbers / email) — typed over in PowerPoint
    _ALIGN = {"l": PP_ALIGN.LEFT, "m": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
    for e in (nums or []):
        if e["slide"] >= len(sl):
            continue
        s = sl[e["slide"]]
        l, t, r, b = _num_bbox(e)
        cx = (l + r) / 2.0; cy = (t + b) / 2.0
        boxw = (r - l) + 44
        ax = e["anchor"][0]
        bl = l if ax == "l" else (r - boxw if ax == "r" else cx - boxw / 2.0)
        vbox = e.get("vbox")
        if vbox:                                   # stacked lines (e.g. title): fixed,
            bt, boxh = vbox[0], vbox[1]            # uniform vertical box -> even leading
            _valign = MSO_ANCHOR.TOP
        else:                                      # single line: centre on its glyph bbox
            boxh = e["px"] * 1.7
            bt = cy - boxh / 2.0
            _valign = MSO_ANCHOR.MIDDLE
        tb = s.shapes.add_textbox(Inches(bl/PXIN), Inches(bt/PXIN), Inches(boxw/PXIN), Inches(boxh/PXIN))
        tf = tb.text_frame; tf.word_wrap = False; tf.auto_size = None
        tf.vertical_anchor = _valign
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        fam, bd = _KIND2PPTX[e["kind"]]
        lines = e["text"].split("\n")
        for li, line in enumerate(lines):
            para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            para.alignment = _ALIGN.get(ax, PP_ALIGN.LEFT)
            run = para.add_run(); run.text = line
            run.font.size = Pt(e["px"] / 2.0); run.font.name = fam; run.font.bold = bd
            run.font.color.rgb = RGBColor(*e["rgb"])
            if e.get("url"):
                run.hyperlink.address = e["url"]
                run.font.underline = False
                _rPr = run._r.get_or_add_rPr()
                for _u in _rPr.findall(qn('a:uFill')) + _rPr.findall(qn('a:uLn')):
                    _rPr.remove(_u)
    prs.save(os.path.join(OUT,name))
    _embed_fonts(os.path.join(OUT,name))
    print("PPTX:",os.path.join(OUT,name))


# families used by the editable text -> the font file(s) that back each PowerPoint slot
_EMBED_FONTS = [
    ("Montserrat Black",    {"regular": "Montserrat-Black.ttf"}),
    ("Montserrat SemiBold", {"regular": "Montserrat-SemiBold.ttf"}),
    ("Montserrat",          {"regular": "Montserrat-Regular.ttf", "bold": "Montserrat-Bold.ttf"}),
    ("Inter",               {"regular": "Inter-Regular.ttf"}),
]

def _embed_fonts(path):
    """Embed the brand TrueType fonts in the .pptx so the editable text renders
    correctly even on machines that don't have Montserrat / Inter installed.
    Fully defensive: on any problem the original (non-embedded) file is kept."""
    import zipfile, shutil
    try:
        from lxml import etree
    except Exception as e:
        print("font embed skipped (no lxml):", e); return
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    CT = "http://schemas.openxmlformats.org/package/2006/content-types"
    RELNS = "http://schemas.openxmlformats.org/package/2006/relationships"
    FONT_RT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"

    # gather the embed plan: only slots whose font file actually exists
    plan, parts, idx = [], [], 1   # plan: (typeface, [(slot, partname, relid)])
    for typeface, slots in _EMBED_FONTS:
        got = []
        for slot, fname in slots.items():
            fpath = os.path.join(FONTS, fname)
            if not os.path.exists(fpath):
                continue
            partname = "fonts/font%d.fntdata" % idx
            relid = "rIdFont%d" % idx
            got.append((slot, partname, relid)); parts.append((partname, fpath)); idx += 1
        if got:
            plan.append((typeface, got))
    if not plan:
        return

    backup = path + ".bak"
    shutil.copyfile(path, backup)
    tmp = path + ".tmp"
    try:
        zin = zipfile.ZipFile(path, "r")
        names = zin.namelist()
        ct_xml = etree.fromstring(zin.read("[Content_Types].xml"))
        # 1) content type for .fntdata
        if not any(el.get("Extension") == "fntdata" for el in ct_xml if el.tag == "{%s}Default" % CT):
            d = etree.SubElement(ct_xml, "{%s}Default" % CT)
            d.set("Extension", "fntdata"); d.set("ContentType", "application/x-fontdata")
        # 2) presentation.xml.rels — add a relationship per font part
        rels = etree.fromstring(zin.read("ppt/_rels/presentation.xml.rels"))
        for _typeface, got in plan:
            for _slot, partname, relid in got:
                rel = etree.SubElement(rels, "{%s}Relationship" % RELNS)
                rel.set("Id", relid); rel.set("Type", FONT_RT); rel.set("Target", partname)
        # 3) presentation.xml — embedTrueTypeFonts + <p:embeddedFontLst>
        pres = etree.fromstring(zin.read("ppt/presentation.xml"))
        pres.set("embedTrueTypeFonts", "1")
        lst = etree.Element("{%s}embeddedFontLst" % P)
        for typeface, got in plan:
            ef = etree.SubElement(lst, "{%s}embeddedFont" % P)
            f = etree.SubElement(ef, "{%s}font" % P); f.set("typeface", typeface)
            for slot, _partname, relid in got:
                sl = etree.SubElement(ef, "{%s}%s" % (P, slot)); sl.set("{%s}id" % R, relid)
        # insert in schema order: right after <p:notesSz> (fallback: before defaultTextStyle/extLst/end)
        notes = pres.find("{%s}notesSz" % P)
        if notes is not None:
            notes.addnext(lst)
        else:
            anchor = pres.find("{%s}defaultTextStyle" % P)
            if anchor is not None: anchor.addprevious(lst)
            else: pres.append(lst)

        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
            for n in names:
                if n == "[Content_Types].xml":
                    zout.writestr(n, etree.tostring(ct_xml, xml_declaration=True, encoding="UTF-8", standalone=True))
                elif n == "ppt/_rels/presentation.xml.rels":
                    zout.writestr(n, etree.tostring(rels, xml_declaration=True, encoding="UTF-8", standalone=True))
                elif n == "ppt/presentation.xml":
                    zout.writestr(n, etree.tostring(pres, xml_declaration=True, encoding="UTF-8", standalone=True))
                else:
                    zout.writestr(zin.getinfo(n), zin.read(n))
            for partname, fpath in parts:
                with open(fpath, "rb") as fh:
                    zout.writestr("ppt/" + partname, fh.read())
        zin.close()
        # validate: it must still open as a valid package
        from pptx import Presentation as _P
        _P(tmp)
        shutil.move(tmp, path)
        os.remove(backup)
        print("fonts embedded: %d slots" % len(parts))
    except Exception as e:
        # restore the good, non-embedded file
        for f in (tmp,):
            try: os.remove(f)
            except Exception: pass
        try: shutil.move(backup, path)
        except Exception: pass
        print("font embed skipped:", e)

def _bake_preview(lang, deck=None):
    """After the PPTX is saved, flatten the native speaker photos + YouTube text
    + editable figures into the slide PNGs so the in-app preview matches the file."""
    # editable figures: paint each onto its slide PNG (exact match to the layout)
    if deck:
        by_slide = {}
        for e in NATIVE_NUMS:
            by_slide.setdefault(e["slide"], []).append(e)
        for si, items in by_slide.items():
            if si >= len(deck):
                continue
            try:
                im = Image.open(deck[si]).convert("RGB"); dr = ImageDraw.Draw(im)
                for e in items:
                    dr.text((e["x"], e["y"]), e["text"], font=_font1(e["kind"], e["px"]),
                            fill=tuple(e["rgb"]), anchor=e["anchor"], spacing=e.get("spacing", 4))
                im.save(deck[si])
            except Exception as ex:
                print("preview nums skip", ex)
    cov = os.path.join(OUT, "01_cover_%s.png" % lang)
    try:
        im = Image.open(cov).convert("RGB")
        for ph in COVER_PHOTOS:
            d = ph["d"]; av = ImageOps.fit(Image.open(ph["path"]).convert("RGB"), (d, d), Image.LANCZOS)
            mask = Image.new("L", (d, d), 0); ImageDraw.Draw(mask).ellipse([0, 0, d, d], fill=255)
            im.paste(av, (int(ph["cx"] - d / 2), int(ph["cy"] - d / 2)), mask)
        im.save(cov)
    except Exception as e:
        print("preview cover skip", e)
    if LINK_INFO:
        eng = os.path.join(OUT, "06_engagement_%s.png" % lang)
        try:
            im = Image.open(eng).convert("RGB"); dr = ImageDraw.Draw(im)
            bf = ImageFont.truetype(os.path.join(FONTS, "Montserrat-Bold.ttf"), 30)
            rf = ImageFont.truetype(os.path.join(FONTS, "Inter-Regular.ttf"), 22)
            y0, y1 = LINK_INFO["box"][1], LINK_INFO["box"][3]; cyb = (y0 + y1) // 2
            dr.text((968, cyb - 20), LINK_INFO["title"], font=bf, fill=(255, 255, 255), anchor="lm")
            dr.text((968, cyb + 20), LINK_INFO["line2"], font=rf, fill=(255, 200, 180), anchor="lm")
            im.save(eng)
        except Exception as e:
            print("preview eng skip", e)


def generate_report(webinar, insights, stats, orgs, lang="en", annotate=False,
                    assets_dir=None, out_dir=None, filename=None):
    """Render the 9-slide editable PPTX for an arbitrary webinar.

    webinar  : data dict (title_lines, subtitle, date, speakers[], youtube*, email,
               segments, comp_sample, contact) — same shape as the module WEBINAR demo.
    insights : {slide_key: one-line highlight text} (used only when annotate=True).
    stats    : output of engine.ingest.build_stats().
    orgs     : list of (company, name, job_title) for the 'most relevant orgs' slide.
    Returns the path to the saved .pptx.
    """
    global WEBINAR, INSIGHTS, ANNOTATE, SUFFIX, ASSETS, OUT, COVER_PHOTOS, LINK_INFO, NATIVE_NUMS
    WEBINAR = webinar
    INSIGHTS = insights or {}
    ANNOTATE = annotate
    SUFFIX = ""
    NATIVE_NUMS = []
    if assets_dir:
        ASSETS = assets_dir
    if out_dir:
        OUT = out_dir
        os.makedirs(OUT, exist_ok=True)
    # total audience = Zoom live attendees + YouTube views (recompute attendance %)
    import copy as _copy, re as _re
    _yt = int(_re.sub(r"[^\d]", "", (webinar.get("youtube_views") or "")) or 0)
    if _yt:
        stats = _copy.deepcopy(stats); _kf = stats["key_facts"]
        _kf["live_attendees"] = _kf["live_attendees"] + _yt
        if _kf.get("registrations"):
            _kf["attendance_rate_pct"] = round(100.0 * _kf["live_attendees"] / _kf["registrations"], 1)
    deck = [slide_cover(stats, lang), slide_facts(stats, lang), slide_country(stats, lang),
            slide_industries(stats, lang), slide_engagement(stats, lang),
            slide_reach(stats, lang), slide_titles(stats, lang), slide_contact(stats, lang)]
    if not filename:
        filename = "ATA_Webinar_Report_%s%s.pptx" % (lang.upper(), "_annotated" if annotate else "")
    build_pptx(deck, filename, photos=list(COVER_PHOTOS), link=LINK_INFO, nums=list(NATIVE_NUMS))
    _bake_preview(lang, deck)
    return os.path.join(OUT, filename)
