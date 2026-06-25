# -*- coding: utf-8 -*-
"""RENMAD Events 2026-2027 infographic (orange + navy), 16:9, rendered to PNG
then placed on a single slide."""
import numpy as np
from PIL import Image, ImageDraw, ImageFont, ImageFilter
import math, os

W, H = 3200, 1800
FONTS = r"C:\Users\Belén\renmad_generator\assets\fonts"
LOGOS = r"C:\Users\Belén\renmad_generator\assets\logos"
IMG_OUT = r"C:\Users\Belén\renmad_generator\assets\renmad_infographic.png"

# ---------- palette ----------
NAVY   = (13, 38, 62)
NAVY_D = (8, 24, 42)
NAVY_L = (24, 58, 90)
ORANGE = (240, 84, 18)
ORANGE_L = (255, 122, 40)
CYAN   = (70, 190, 235)
WHITE  = (255, 255, 255)
INK    = (28, 32, 40)
GREY   = (108, 116, 126)
PANEL  = (245, 247, 249)
ROW_A  = (236, 240, 244)
LINE   = (214, 220, 226)

def F(name, size):
    return ImageFont.truetype(os.path.join(FONTS, name), size)

mB  = lambda s: F("Montserrat-Black.ttf", s)
mb  = lambda s: F("Montserrat-Bold.ttf", s)
ms  = lambda s: F("Montserrat-SemiBold.ttf", s)
mr  = lambda s: F("Montserrat-Regular.ttf", s)
ir  = lambda s: F("Inter-Regular.ttf", s)

# ================= ICON RENDERER (supersampled) =================
def render_icon(fn, px, color, **kw):
    T = px * 3
    im = Image.new("RGBA", (T, T), (0, 0, 0, 0))
    d = ImageDraw.Draw(im)
    fn(d, T, color, **kw)
    return im.resize((px, px), Image.LANCZOS)

def _w(T, k=0.085):
    return max(2, int(T * k))

def ic_server(d, T, c):
    m = T * 0.18; aw = T - 2 * m; bh = (T - 2 * m) * 0.26; gap = bh * 0.42
    y = m
    for _ in range(3):
        d.rounded_rectangle([m, y, m + aw, y + bh], radius=bh * 0.28,
                            outline=c, width=_w(T, 0.07))
        d.ellipse([m + aw * 0.10 - bh*0.12, y + bh/2 - bh*0.12,
                   m + aw * 0.10 + bh*0.12, y + bh/2 + bh*0.12], fill=c)
        d.line([m + aw*0.30, y + bh/2, m + aw*0.85, y + bh/2], fill=c, width=_w(T, 0.05))
        y += bh + gap

def ic_battery(d, T, c):
    m = T * 0.16; w = T - 2*m; h = w * 0.56; x0 = m; y0 = (T - h)/2
    d.rounded_rectangle([x0, y0, x0 + w*0.9, y0 + h], radius=h*0.18,
                        outline=c, width=_w(T, 0.075))
    d.rounded_rectangle([x0 + w*0.9, y0 + h*0.3, x0 + w, y0 + h*0.7],
                        radius=h*0.1, fill=c)
    cx = x0 + w*0.45
    bolt = [(cx+h*0.10, y0+h*0.18), (cx-h*0.18, y0+h*0.56),
            (cx+h*0.02, y0+h*0.56), (cx-h*0.10, y0+h*0.86),
            (cx+h*0.22, y0+h*0.44), (cx+h*0.02, y0+h*0.44)]
    d.polygon(bolt, fill=c)

def ic_flame(d, T, c):
    cx = T/2; m = T*0.16
    pts = [(cx, m), (cx+T*0.20, T*0.36), (cx+T*0.26, T*0.60),
           (cx+T*0.16, T*0.80), (cx, T-m), (cx-T*0.16, T*0.80),
           (cx-T*0.26, T*0.60), (cx-T*0.18, T*0.40)]
    d.polygon(pts, fill=c)
    inner = [(cx, T*0.46), (cx+T*0.11, T*0.64), (cx, T*0.82), (cx-T*0.11, T*0.64)]
    d.polygon(inner, fill=(255,255,255,235))

def ic_globe(d, T, c):
    m = T*0.16; r = (T-2*m)/2; cx=cy=T/2; wln=_w(T,0.06)
    d.ellipse([cx-r, cy-r, cx+r, cy+r], outline=c, width=wln)
    d.line([cx-r, cy, cx+r, cy], fill=c, width=wln)
    d.ellipse([cx-r*0.5, cy-r, cx+r*0.5, cy+r], outline=c, width=wln)
    d.arc([cx-r, cy-r*0.55, cx+r, cy+r*1.55], 200, 340, fill=c, width=wln)
    d.arc([cx-r, cy-r*1.55, cx+r, cy+r*0.55], 20, 160, fill=c, width=wln)

def ic_chart(d, T, c):
    m = T*0.18; wln=_w(T,0.07)
    d.line([m, m, m, T-m], fill=c, width=wln)
    d.line([m, T-m, T-m, T-m], fill=c, width=wln)
    pts = [(m+ (T-2*m)*0.05, T-m-(T-2*m)*0.18),
           (m+(T-2*m)*0.38, T-m-(T-2*m)*0.46),
           (m+(T-2*m)*0.60, T-m-(T-2*m)*0.30),
           (T-m-(T-2*m)*0.04, m+(T-2*m)*0.05)]
    d.line(pts, fill=c, width=_w(T,0.075), joint="curve")
    ax, ay = pts[-1]
    d.polygon([(ax,ay),(ax-T*0.12,ay-T*0.01),(ax-T*0.02,ay+T*0.12)], fill=c)

def ic_calendar(d, T, c):
    m=T*0.17; wln=_w(T,0.065)
    x0,y0,x1,y1=m,m+T*0.06,T-m,T-m
    d.rounded_rectangle([x0,y0,x1,y1], radius=T*0.06, outline=c, width=wln)
    d.line([x0,y0+T*0.16,x1,y0+T*0.16], fill=c, width=wln)
    d.line([x0+T*0.22,m-T*0.0,x0+T*0.22,y0+T*0.10], fill=c, width=wln)
    d.line([x1-T*0.22,m-T*0.0,x1-T*0.22,y0+T*0.10], fill=c, width=wln)
    r=T*0.035
    for ix in range(3):
        for iy in range(2):
            cxp=x0+T*0.18+ix*T*0.24; cyp=y0+T*0.34+iy*T*0.22
            d.ellipse([cxp-r,cyp-r,cxp+r,cyp+r], fill=c)

def ic_people(d, T, c):
    def fig(cx, scale, col):
        hr=T*0.13*scale; hy=T*0.34
        d.ellipse([cx-hr,hy-hr,cx+hr,hy+hr], fill=col)
        bw=T*0.30*scale
        d.pieslice([cx-bw/2,hy+hr*0.4,cx+bw/2,hy+hr*0.4+bw], 180,360, fill=col)
    fig(T*0.62, 0.92, tuple(int(v*0.55+255*0.45) for v in c[:3]))
    fig(T*0.42, 1.0, c)

def ic_pin(d, T, c):
    cx=T/2; cy=T*0.40; r=T*0.22; wln=_w(T,0.08)
    d.arc([cx-r,cy-r,cx+r,cy+r], 0, 360, fill=c, width=wln)
    tip=(cx, T*0.86)
    d.line([(cx-r*0.86, cy+r*0.5), tip], fill=c, width=wln)
    d.line([(cx+r*0.86, cy+r*0.5), tip], fill=c, width=wln)
    d.ellipse([cx-r*0.34,cy-r*0.34,cx+r*0.34,cy+r*0.34], fill=c)

def ic_network(d, T, c):
    nodes=[(T*0.30,T*0.32),(T*0.72,T*0.28),(T*0.5,T*0.68),(T*0.78,T*0.66)]
    for a in range(len(nodes)):
        for b in range(a+1,len(nodes)):
            d.line([nodes[a],nodes[b]], fill=c, width=_w(T,0.04))
    for (x,y) in nodes:
        rr=T*0.075
        d.ellipse([x-rr,y-rr,x+rr,y+rr], fill=c)

def ic_mic(d, T, c):
    cx=T/2; wln=_w(T,0.07)
    d.rounded_rectangle([cx-T*0.13,T*0.16,cx+T*0.13,T*0.58], radius=T*0.13,
                        fill=c)
    d.arc([cx-T*0.22,T*0.30,cx+T*0.22,T*0.66], 20,160, fill=c, width=wln)
    d.line([cx,T*0.66,cx,T*0.80], fill=c, width=wln)
    d.line([cx-T*0.16,T*0.82,cx+T*0.16,T*0.82], fill=c, width=wln)

def ic_bulb(d, T, c):
    cx=T/2; r=T*0.24; wln=_w(T,0.07)
    d.arc([cx-r,T*0.14,cx+r,T*0.14+2*r], 0,360, fill=c, width=wln)
    d.line([cx-r*0.5,T*0.14+r*1.5,cx-r*0.32,T*0.74], fill=c, width=wln)
    d.line([cx+r*0.5,T*0.14+r*1.5,cx+r*0.32,T*0.74], fill=c, width=wln)
    for yy in (0.76,0.83):
        d.line([cx-T*0.13,T*yy,cx+T*0.13,T*yy], fill=c, width=wln)
    d.line([cx-T*0.09,T*0.90,cx+T*0.09,T*0.90], fill=c, width=wln)

# ================= FLAGS =================
def flag(kind, w=120, h=80):
    im = Image.new("RGBA",(w,h),(255,255,255,255))
    d = ImageDraw.Draw(im)
    if kind=="es":
        d.rectangle([0,0,w,h], fill=(198,11,30))
        d.rectangle([0,h*0.25,w,h*0.75], fill=(255,196,0))
    elif kind=="it":
        d.rectangle([0,0,w/3,h], fill=(0,140,69))
        d.rectangle([w/3,0,2*w/3,h], fill=(255,255,255))
        d.rectangle([2*w/3,0,w,h], fill=(206,43,55))
    elif kind=="pl":
        d.rectangle([0,0,w,h/2], fill=(255,255,255))
        d.rectangle([0,h/2,w,h], fill=(220,20,60))
    elif kind=="cl":
        d.rectangle([0,0,w,h/2], fill=(255,255,255))
        d.rectangle([0,h/2,w,h], fill=(213,43,30))
        d.rectangle([0,0,h/2,h/2], fill=(0,57,166))
        # star
        cx,cy,r=h/4,h/4,h*0.17
        pts=[]
        for i in range(10):
            ang=-math.pi/2+i*math.pi/5
            rr=r if i%2==0 else r*0.42
            pts.append((cx+math.cos(ang)*rr, cy+math.sin(ang)*rr))
        d.polygon(pts, fill=(255,255,255))
    elif kind=="mx":
        d.rectangle([0,0,w/3,h], fill=(0,104,71))
        d.rectangle([w/3,0,2*w/3,h], fill=(255,255,255))
        d.rectangle([2*w/3,0,w,h], fill=(206,17,38))
        d.ellipse([w/2-h*0.12,h/2-h*0.12,w/2+h*0.12,h/2+h*0.12], outline=(120,80,30), width=3)
    d.rectangle([0,0,w-1,h-1], outline=(150,156,162), width=2)
    return im

# ================= CANVAS =================
img = Image.new("RGB",(W,H), WHITE)
draw = ImageDraw.Draw(img)

def text(x,y,s,font,fill,anchor="la",spacing=0):
    draw.text((x,y),s,font=font,fill=fill,anchor=anchor)

def paste_icon(fn,px,color,x,y,**kw):
    ic=render_icon(fn,px,color,**kw)
    img.paste(ic,(int(x),int(y)),ic)

# ---------- HEADER ----------
HH=648
hdr=np.zeros((HH,W,3))
ty=np.linspace(0,1,HH)[:,None,None]
hdr[:]= np.array(NAVY_D)[None,None,:]*(1-ty)+np.array(NAVY)[None,None,:]*ty
# orange glow bottom-left + cyan glow right
yy,xx=np.mgrid[0:HH,0:W]
g1=np.clip(1-np.sqrt(((xx-W*0.18)/ (W*0.5))**2+((yy-HH*1.0)/(HH*0.9))**2),0,1)**2
hdr+=np.array(ORANGE)[None,None,:]*g1[...,None]*0.55
g2=np.clip(1-np.sqrt(((xx-W*0.9)/(W*0.34))**2+((yy-HH*0.5)/(HH*0.9))**2),0,1)**2
hdr+=np.array(CYAN)[None,None,:]*g2[...,None]*0.30
header=Image.fromarray(np.clip(hdr,0,255).astype(np.uint8))
hd=ImageDraw.Draw(header,"RGBA")

# energy frieze (left, subtle orange)
def h_turbine(cx,base,s):
    col=(255,150,90,120); wln=int(6*s)
    hh=160*s; top=base-hh
    hd.line([(cx,base),(cx,top)], fill=col, width=wln)
    for a in (90,210,330):
        ar=math.radians(a)
        hd.line([(cx,top),(cx+math.cos(ar)*70*s, top-math.sin(ar)*70*s)], fill=col, width=wln)
def h_solar(cx,base,s):
    col=(255,150,90,110); wln=int(5*s)
    w=150*s
    pts=[(cx-w/2,base),(cx+w/2,base-30*s),(cx+w/2,base-66*s),(cx-w/2,base-36*s)]
    hd.line(pts+[pts[0]], fill=col, width=wln)
for i,cx in enumerate([180,360,540,720]):
    if i%2==0: h_turbine(cx,HH-70,1.0)
    else: h_solar(cx,HH-90,1.0)
# AI mesh right
mesh=[(2760,250),(2980,200),(3080,330),(2880,400),(3060,470),(2720,420)]
for a in range(len(mesh)):
    for b in range(a+1,len(mesh)):
        if (a+b)%2==0:
            hd.line([mesh[a],mesh[b]], fill=(90,200,240,90), width=3)
for (x,y) in mesh:
    hd.ellipse([x-9,y-9,x+9,y+9], fill=(120,210,245,200))

img.paste(header,(0,0))
draw=ImageDraw.Draw(img)

# logo
logo=Image.open(os.path.join(LOGOS,"logo_renmad_events.png")).convert("RGBA")
lw=380; lh=int(lw*logo.height/logo.width)
logo=logo.resize((lw,lh),Image.LANCZOS)
img.paste(logo,(56,48),logo)

# AI big
text(2980,300,"AI", mB(150), (150,220,250), anchor="mm")

# title block (centered ~1560)
CXT=1580
t1=mB(96)
# line1: RENMAD EVENTS —
s_ren="RENMAD "; s_eve="EVENTS"; s_dash=" —"
w_all=draw.textlength(s_ren,font=t1)+draw.textlength(s_eve,font=t1)+draw.textlength(s_dash,font=t1)
x0=CXT-w_all/2; yb=58
draw.text((x0,yb),s_ren,font=t1,fill=WHITE)
x0+=draw.textlength(s_ren,font=t1)
draw.text((x0,yb),s_eve,font=t1,fill=ORANGE_L)
x0+=draw.textlength(s_eve,font=t1)
draw.text((x0,yb),s_dash,font=t1,fill=WHITE)
text(CXT,yb+108,"EVENTS CALENDAR",t1,WHITE,anchor="ma")
text(CXT,yb+232,"(International Series)", ms(46), ORANGE_L, anchor="ma")
text(CXT,yb+312,"RENEWABLES · STORAGE · HYDROGEN · BIOMETHANE · DATA CENTERS · AI",
     ms(34), WHITE, anchor="ma")
text(CXT,yb+360,"EVERY FRONTIER OF ENERGY", mB(60), ORANGE_L, anchor="ma")

# stat row
sy=yb+470
stats=[(ic_calendar,"2026 – 2027"),(ic_people,"14 summits"),
       (ic_globe,"5 countries"),(ic_pin,"2 continents")]
total_w=1700; sx=CXT-total_w/2; step=total_w/len(stats)
for i,(icon,label) in enumerate(stats):
    cx=sx+step*i+60
    paste_icon(icon,58,ORANGE_L,cx-29,sy-29)
    text(cx+44,sy,label, ms(42), WHITE, anchor="lm")
    if i>0:
        draw.line([(sx+step*i-10,sy-34),(sx+step*i-10,sy+34)], fill=(120,140,160), width=2)

# ---------- BODY ----------
BX=48; gap=40
Lx0=BX; Lw=1480
Rx0=Lx0+Lw+gap; Rw=W-Rx0-BX
by=HH+44
barh=78

def panel_header(x,w,y,h,color,label):
    draw.rounded_rectangle([x,y,x+w,y+h], radius=16, fill=color)
    text(x+w/2,y+h/2,label, mB(40), WHITE, anchor="mm")

panel_header(Lx0,Lw,by,barh,ORANGE,"STILL TO COME IN 2026")
panel_header(Rx0,Rw,by,barh,NAVY,"THE 2027 SEASON  (12 SUMMITS)")

# --- LEFT: 2026 events ---
ev2026=[("server","Data Centers Italia","Nov 11–12","Milan, Italy","it"),
        ("h2","Hidrógeno","Nov 18–19","Zaragoza, Spain","es")]
iconmap={"server":ic_server,"battery":ic_battery,"flame":ic_flame,
         "globe":ic_globe,"chart":ic_chart,"network":ic_network}

py=by+barh+24
rowh=128
draw.rounded_rectangle([Lx0,py,Lx0+Lw,py+rowh*2+24], radius=18, fill=PANEL)
def event_row_L(y,kind,name,date,loc,fl):
    ts=86
    tx=Lx0+34; ty2=y+(rowh-ts)/2
    draw.rounded_rectangle([tx,ty2,tx+ts,ty2+ts], radius=16, fill=ORANGE)
    if kind=="h2":
        text(tx+ts/2,ty2+ts/2-4,"H", mB(48), WHITE, anchor="mm")
        text(tx+ts/2+30,ty2+ts/2+18,"2", mb(30), WHITE, anchor="mm")
    else:
        paste_icon(iconmap[kind],int(ts*0.62),WHITE,tx+ts*0.19,ty2+ts*0.19)
    cy=y+rowh/2
    text(tx+ts+34,cy,name, mb(42), INK, anchor="lm")
    draw.line([(Lx0+660,cy-30),(Lx0+660,cy+30)], fill=LINE, width=2)
    paste_icon(ic_calendar,46,ORANGE,Lx0+690,cy-23)
    text(Lx0+748,cy,date, ms(36), INK, anchor="lm")
    draw.line([(Lx0+930,cy-30),(Lx0+930,cy+30)], fill=LINE, width=2)
    paste_icon(ic_pin,46,ORANGE,Lx0+958,cy-23)
    text(Lx0+1012,cy,loc, ms(32), INK, anchor="lm")
    fim=flag(fl).resize((92,61),Image.LANCZOS)
    img.paste(fim,(Lx0+Lw-116,int(cy-30)))
for i,(k,n,dt,lc,fl) in enumerate(ev2026):
    event_row_L(py+12+i*rowh,k,n,dt,lc,fl)

# --- value props ---
vy=py+rowh*2+24+70
vals=[(ic_globe,"GLOBAL","PERSPECTIVE"),(ic_people,"SENIOR","AUDIENCE"),
      (ic_network,"STRATEGIC","NETWORKING"),(ic_mic,"WORLD-CLASS","SPEAKERS"),
      (ic_bulb,"REAL SOLUTIONS","FOR THE FUTURE")]
vstep=Lw/len(vals)
for i,(icon,l1,l2) in enumerate(vals):
    cx=Lx0+vstep*i+vstep/2
    paste_icon(icon,96,ORANGE,cx-48,vy)
    text(cx,vy+118,l1, mb(28), INK, anchor="ma")
    text(cx,vy+154,l2, mb(28), INK, anchor="ma")

# --- RIGHT: 2027 season ---
ev2027=[("server","Data Centers","27–28 Jan","Zaragoza, Spain","es"),
        ("battery","Storage Italia","9–10 Feb","Bologna, Italy","it"),
        ("flame","Biometano","10–11 Feb","Toledo, Spain","es"),
        ("battery","Almacenamiento","31 Mar – 1 Apr","Seville, Spain","es"),
        ("globe","Chile","29–30 Jul","Santiago, Chile","cl"),
        ("h2","Hidrógeno","November","Zaragoza, Spain","es"),
        ("server","Data Centers Italia","November","Milan, Italy","it"),
        ("chart","Invest","Dates soon","Madrid, Spain","es"),
        ("chart","Invest Italia","Dates soon","Milan, Italy","it"),
        ("globe","México","Dates soon","Mexico","mx"),
        ("battery","Storage Polska","Dates soon","Warsaw, Poland","pl"),
        ("ai","Useful AI","Dates soon","Madrid, Spain","es")]
ry0=by+barh+18
ravail=(H-116)-ry0          # stop above the footer band
rh=ravail/12
for i,(k,n,dt,lc,fl) in enumerate(ev2027):
    y=ry0+i*rh
    if i%2==1:
        draw.rectangle([Rx0,y,Rx0+Rw,y+rh], fill=ROW_A)
    cy=y+rh/2
    ts=64; tx=Rx0+22; ty2=cy-ts/2
    tile_col = NAVY if k in ("server",) else ORANGE
    draw.rounded_rectangle([tx,ty2,tx+ts,ty2+ts], radius=12, fill=tile_col)
    if k=="h2":
        text(tx+ts/2-6,ty2+ts/2-2,"H", mB(34), WHITE, anchor="mm")
        text(tx+ts/2+18,ty2+ts/2+12,"2", mb(22), WHITE, anchor="mm")
    elif k=="ai":
        text(tx+ts/2,ty2+ts/2,"AI", mB(30), WHITE, anchor="mm")
    else:
        paste_icon(iconmap[k],int(ts*0.6),WHITE,tx+ts*0.2,ty2+ts*0.2)
    text(tx+ts+26,cy,n, mb(34), INK, anchor="lm")
    draw.line([(Rx0+640,cy-26),(Rx0+640,cy+26)], fill=LINE, width=2)
    paste_icon(ic_calendar,40,ORANGE,Rx0+666,cy-20)
    text(Rx0+716,cy,dt, ms(30), INK, anchor="lm")
    draw.line([(Rx0+980,cy-26),(Rx0+980,cy+26)], fill=LINE, width=2)
    paste_icon(ic_pin,40,ORANGE,Rx0+1006,cy-20)
    text(Rx0+1052,cy,lc, ms(28), INK, anchor="lm")
    fim=flag(fl).resize((84,56),Image.LANCZOS)
    img.paste(fim,(Rx0+Rw-104,int(cy-28)))

# ---------- FOOTER ----------
FY=H-104
draw.rectangle([0,FY,W,H], fill=NAVY)
fx=60; fcy=FY+52
draw.text((fx,fcy),"CONNECTING MARKETS. ",font=mB(40),fill=ORANGE_L,anchor="lm")
wM=draw.textlength("CONNECTING MARKETS. ",font=mB(40))
draw.text((fx+wM,fcy),"ACCELERATING TRANSITION.",font=mB(40),fill=WHITE,anchor="lm")
# right cluster
paste_icon(ic_globe,46,WHITE,W-960,fcy-23)
text(W-900,fcy,"www.renmad.com", ms(36), WHITE, anchor="lm")
draw.line([(W-560,fcy-26),(W-560,fcy+26)], fill=(120,140,160), width=2)
draw.rounded_rectangle([W-520,fcy-26,W-468,fcy+26], radius=10, fill=ORANGE_L)
text(W-494,fcy,"in", mB(40), NAVY, anchor="mm")
text(W-446,fcy,"RENMAD EVENTS", mb(36), WHITE, anchor="lm")

img.save(IMG_OUT, quality=95)
print("Saved image:", IMG_OUT)

# ================= place on slide =================
from pptx import Presentation
from pptx.util import Inches
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
s=prs.slides.add_slide(prs.slide_layouts[6])
s.shapes.add_picture(IMG_OUT,0,0,width=prs.slide_width,height=prs.slide_height)
OUT=r"C:\Users\Belén\RENMAD_Events_Infographic.pptx"
prs.save(OUT)
print("Saved:", OUT)
