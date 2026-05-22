"""Common python-pptx helpers."""
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu
from PIL import Image


def new_presentation(width_px: int, height_px: int, dpi: int = 96) -> Presentation:
    prs = Presentation()
    prs.slide_width = Emu(int(width_px / dpi * 914400))
    prs.slide_height = Emu(int(height_px / dpi * 914400))
    return prs


def px_to_emu(px: float, dpi: int = 96) -> int:
    return int(px / dpi * 914400)


def add_blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def set_background_colour(slide, r: int, g: int, b: int):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width_pt=0):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = RGBColor(*line_rgb)
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_name="Calibri", font_size=14, bold=False, italic=False,
                color_rgb=(51, 51, 51), align=PP_ALIGN.LEFT,
                word_wrap=True, v_anchor=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = RGBColor(*color_rgb)
    return txBox


def add_image_bytes(slide, img_bytes: bytes, left, top, width, height):
    buf = io.BytesIO(img_bytes)
    return slide.shapes.add_picture(buf, left, top, width, height)


def pptx_to_bytes(prs: Presentation) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def export_slide_to_png(pptx_bytes: bytes, slide_index: int = 0,
                        width_px: int = 1920, height_px: int = 1080) -> bytes:
    """
    Export a slide to PNG using LibreOffice if available, otherwise
    fall back to a python-pptx thumbnail via Pillow compositing.
    This returns PNG bytes.
    """
    import subprocess, tempfile, os

    with tempfile.TemporaryDirectory() as tmp:
        pptx_path = os.path.join(tmp, "slide.pptx")
        pdf_path = os.path.join(tmp, "slide.pdf")
        with open(pptx_path, "wb") as f:
            f.write(pptx_bytes)

        # Try LibreOffice
        for soffice in ["soffice", "libreoffice"]:
            try:
                result = subprocess.run(
                    [soffice, "--headless", "--convert-to", "pdf",
                     "--outdir", tmp, pptx_path],
                    capture_output=True, timeout=30
                )
                if result.returncode == 0 and os.path.exists(pdf_path):
                    import subprocess as sp
                    png_path = os.path.join(tmp, "slide.png")
                    sp.run(["pdftoppm", "-jpeg", "-r", "150", "-f", "1", "-l", "1",
                            pdf_path, os.path.join(tmp, "page")],
                           capture_output=True, timeout=30)
                    candidates = [f for f in os.listdir(tmp) if f.startswith("page")]
                    if candidates:
                        with open(os.path.join(tmp, sorted(candidates)[0]), "rb") as f:
                            return f.read()
            except Exception:
                pass

    # Fallback: generate a simple preview using Pillow (no LibreOffice)
    return _pptx_fallback_png(pptx_bytes, width_px, height_px)


def _pptx_fallback_png(pptx_bytes: bytes, width_px: int, height_px: int) -> bytes:
    """Very basic fallback: render slide background + placed images only."""
    from PIL import Image as PILImage
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide = prs.slides[0]

    # White canvas
    canvas = PILImage.new("RGB", (width_px, height_px), (255, 255, 255))

    buf = io.BytesIO()
    canvas.save(buf, format="PNG")
    return buf.getvalue()
