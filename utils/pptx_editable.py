"""Helpers for adding EDITABLE picture shapes to PPTX slides.

add_shaped_picture() embeds the ORIGINAL photo (downscaled, never pre-cropped)
as a picture shape whose geometry is an ellipse / rounded rectangle, with
PowerPoint crop values replicating the ImageOps.fit centre crop used for the
PNG render. In PowerPoint the user can move/resize the shape AND use
Picture Format → Crop to drag the photo inside the frame — fixing
"face too high / cut by the circle" without regenerating the slide.
"""
import io
from PIL import Image
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

_EMU_PX = 9525          # EMU per px at 96 dpi
_MAX_SRC_SIDE = 1600    # downscale huge source photos before embedding


def _px(p) -> Emu:
    return Emu(int(p * _EMU_PX))


def _set_geometry(pic, prst: str, adj_val: int | None = None):
    """Replace the picture's preset geometry (rect → ellipse/roundRect)."""
    spPr = pic._element.spPr
    for g in list(spPr.findall(qn('a:prstGeom'))) + list(spPr.findall(qn('a:custGeom'))):
        spPr.remove(g)
    geom = etree.SubElement(spPr, qn('a:prstGeom'))
    geom.set('prst', prst)
    av = etree.SubElement(geom, qn('a:avLst'))
    if adj_val is not None:
        gd = etree.SubElement(av, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {adj_val}')
    # a:prstGeom must come right after a:xfrm inside spPr
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is not None:
        spPr.remove(geom)
        xfrm.addnext(geom)


def add_shaped_picture(slide, photo_img: Image.Image,
                       x: float, y: float, w: float, h: float,
                       shape: str = "ellipse", radius_px: float = 0,
                       border_rgb=None, border_px: float = 0):
    """Add `photo_img` as a movable/re-croppable picture shape.

    x/y/w/h are in canvas pixels (96 dpi). The full photo is embedded and the
    crop is set to the centre box-fit region, so the initial view matches the
    PNG render while the rest of the photo stays available to the Crop tool.
    shape: 'ellipse' | 'roundRect' | 'rect'.
    """
    img = photo_img.convert("RGB")
    if max(img.size) > _MAX_SRC_SIDE:
        ratio = _MAX_SRC_SIDE / max(img.size)
        img = img.resize((max(1, int(img.width * ratio)),
                          max(1, int(img.height * ratio))), Image.LANCZOS)

    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=90)
    buf.seek(0)

    pic = slide.shapes.add_picture(buf, _px(x), _px(y), _px(w), _px(h))

    # Crop to the same centre region ImageOps.fit(w, h) shows
    iw, ih = img.size
    target_ratio = w / h
    region_w = min(iw, ih * target_ratio)
    region_h = region_w / target_ratio
    if region_h > ih:
        region_h = ih
        region_w = region_h * target_ratio
    crop_x = (iw - region_w) / 2 / iw
    crop_y = (ih - region_h) / 2 / ih
    pic.crop_left = crop_x
    pic.crop_right = crop_x
    pic.crop_top = crop_y
    pic.crop_bottom = crop_y

    if shape == "ellipse":
        _set_geometry(pic, "ellipse")
    elif shape == "roundRect":
        adj = int(max(0.0, min(0.5, radius_px / max(1.0, min(w, h)))) * 100000)
        _set_geometry(pic, "roundRect", adj)

    if border_rgb and border_px:
        pic.line.color.rgb = RGBColor(*[int(c) for c in border_rgb[:3]])
        pic.line.width = Pt(border_px * 0.75)   # px → pt at 96 dpi

    return pic


def add_picture_el(slide, img: Image.Image, x: float, y: float,
                   w: float, h: float):
    """Add a plain (already rendered, transparent-PNG) image as a movable shape."""
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return slide.shapes.add_picture(buf, _px(x), _px(y), _px(w), _px(h))
