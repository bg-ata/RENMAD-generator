# -*- coding: utf-8 -*-
"""RENMAD Events 2026-2027  ·  'Painting the dream' hero poster (16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
RED    = RGBColor(0xC8, 0x28, 0x3A)
ORANGE = RGBColor(0xFF, 0x7A, 0x3C)   # ember orange (accents / glow text)
BLUE   = RGBColor(0x3F, 0xB6, 0xE8)
GREEN  = RGBColor(0x5A, 0xB6, 0x3C)
PURPLE = RGBColor(0x9A, 0x6C, 0xD0)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MIST   = RGBColor(0xCF, 0xCB, 0xCE)   # soft light text
DIM    = RGBColor(0x9A, 0x96, 0x9C)

F_BLACK = "Montserrat Black"
F_BOLD  = "Montserrat"
F_SEMI  = "Montserrat SemiBold"
F_BODY  = "Inter"

BG   = r"C:\Users\Belén\renmad_generator\assets\calendar_hero_bg.png"
LOGO = r"C:\Users\Belén\renmad_generator\assets\logos\logo_renmad_events.png"
OUT  = r"C:\Users\Belén\RENMAD_Events_Calendar_2026-2027.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def no_shadow(shape):
    spPr = shape._element.spPr
    for e in spPr.findall(qn("a:effectLst")):
        spPr.remove(e)
    spPr.append(spPr.makeelement(qn("a:effectLst"), {}))


def fill_alpha(shape, rgb, alpha_pct):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    sF = shape.fill._xPr.find(qn("a:solidFill"))
    srgb = sF.find(qn("a:srgbClr"))
    srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha_pct * 1000))}))


def chip(x, y, w, h, rad=0.10):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    fill_alpha(s, RGBColor(0x12, 0x10, 0x16), 52)
    s.line.color.rgb = RGBColor(0x7A, 0x76, 0x80)
    s.line.width = Pt(0.75)
    try:
        s.adjustments[0] = rad
    except Exception:
        pass
    no_shadow(s)
    return s


def solid(x, y, w, h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE, rad=0.18, line=False):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if not line:
        s.line.fill.background()
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = rad
        except Exception:
            pass
    no_shadow(s)
    return s


def textbox(x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.0)
    tf.vertical_anchor = anchor
    return tf


def par(tf, text, font, size, color, bold=False, align=PP_ALIGN.LEFT,
        sa=0, sb=0, ls=1.0, first=False, tracking=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(sa); p.space_before = Pt(sb); p.line_spacing = ls
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    if tracking is not None:
        r._r.get_or_add_rPr().set("spc", str(tracking))
    return p


def add_run(p, text, font, size, color, bold=False, tracking=None):
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    if tracking is not None:
        r._r.get_or_add_rPr().set("spc", str(tracking))
    return r


# ================= background =================
slide.shapes.add_picture(BG, 0, 0, width=prs.slide_width, height=prs.slide_height)

# ================= header / headline =================
slide.shapes.add_picture(LOGO, Inches(0.6), Inches(0.52), width=Inches(2.25))

tf = textbox(0.64, 1.46, 7.8, 0.3)
par(tf, "RENMAD EVENTS   ·   THE SEASON AHEAD", F_SEMI, 11.5, ORANGE,
    first=True, tracking=300)

# vertical roll-call: white words, ember dots
tf = textbox(0.62, 1.80, 7.8, 0.34)
verticals = ["RENEWABLES", "STORAGE", "HYDROGEN", "BIOMETHANE", "DATA CENTERS", "AI"]
p = par(tf, verticals[0], F_SEMI, 11.5, WHITE, first=True, tracking=10)
for w in verticals[1:]:
    add_run(p, "   .   ", F_SEMI, 11.5, ORANGE, tracking=10)
    add_run(p, w, F_SEMI, 11.5, WHITE, tracking=10)

# payoff headline
tf = textbox(0.6, 2.16, 7.8, 1.3)
par(tf, "EVERY FRONTIER", F_BLACK, 37, ORANGE, first=True, ls=0.98)
par(tf, "OF ENERGY", F_BLACK, 37, WHITE, ls=0.98, sb=2)

tf = textbox(0.64, 3.28, 7.8, 0.3)
p = par(tf, "2026 – 2027", F_SEMI, 12, WHITE, first=True, tracking=120)
add_run(p, "      14 summits      5 countries      2 continents", F_SEMI, 12, MIST)

# ================= 2026 finale (top-right) =================
tf = textbox(8.45, 0.58, 4.4, 0.3, anchor=MSO_ANCHOR.MIDDLE)
par(tf, "STILL TO COME IN 2026", F_SEMI, 11, ORANGE, align=PP_ALIGN.RIGHT,
    first=True, tracking=260)


def finale_chip(x, y, w, h, color, dd, mm, name, city):
    chip(x, y, w, h, rad=0.10)
    bw = 1.02
    solid(x + 0.12, y + 0.13, bw, h - 0.26, color, rad=0.14)
    tf = textbox(x + 0.12, y + 0.13, bw, h - 0.26, anchor=MSO_ANCHOR.MIDDLE)
    par(tf, dd, F_BLACK, 19, WHITE, align=PP_ALIGN.CENTER, first=True, ls=0.9)
    par(tf, mm, F_SEMI, 10, WHITE, align=PP_ALIGN.CENTER, sb=1, tracking=200)
    tf = textbox(x + bw + 0.28, y + 0.13, w - bw - 0.42, h - 0.26,
                 anchor=MSO_ANCHOR.MIDDLE)
    par(tf, name, F_BOLD, 12.5, WHITE, bold=True, first=True, ls=0.95)
    par(tf, city, F_BODY, 9.5, MIST, sb=3)


fw, fh = 4.38, 0.95
finale_chip(8.47, 0.98, fw, fh, BLUE,  "11–12", "NOV", "Data Centers Italia", "Milan, Italy")
finale_chip(8.47, 2.05, fw, fh, GREEN, "18–19", "NOV", "Hidrógeno", "Zaragoza, Spain")

# ================= 2027 lineup (lower band) =================
tf = textbox(0.62, 3.62, 9.0, 0.36)
p = par(tf, "THE 2027 SEASON", F_BLACK, 17, WHITE, first=True, tracking=120)
add_run(p, "    ·    12 summits across the energy transition", F_SEMI, 11.5, MIST)

events = [
    ("Data Centers",        "27–28 Jan",  "Zaragoza, Spain", BLUE),
    ("Storage Italia",      "9–10 Feb",   "Bologna, Italy",  ORANGE),
    ("Biometano",           "10–11 Feb",  "Toledo, Spain",   GREEN),
    ("Almacenamiento",      "31 Mar–1 Apr", "Seville, Spain", ORANGE),
    ("Chile",               "29–30 Jul",  "Santiago, Chile", ORANGE),
    ("Hidrógeno",           "November",    "Zaragoza, Spain", GREEN),
    ("Data Centers Italia", "November",    "Milan, Italy",    BLUE),
    ("Invest",              "Dates soon",  "Madrid, Spain",   PURPLE),
    ("Invest Italia",       "Dates soon",  "Milan, Italy",    PURPLE),
    ("México",              "Dates soon",  "Mexico",          RED),
    ("Storage Polska",      "Dates soon",  "Warsaw, Poland",  ORANGE),
    ("Useful AI",           "Dates soon",  "Madrid, Spain",   RED),
]

cols, rows = 4, 3
gx, gy = 0.20, 0.16
left, right = 0.55, 12.85
cw = (right - left - (cols - 1) * gx) / cols
top, ch = 4.06, 0.78
for i, (name, date, city, color) in enumerate(events):
    r, c = divmod(i, cols)
    x = left + c * (cw + gx)
    y = top + r * (ch + gy)
    chip(x, y, cw, ch)
    # glowing category dot
    dot = solid(x + 0.20, y + ch / 2 - 0.085, 0.17, 0.17, color,
                shape=MSO_SHAPE.OVAL)
    tf = textbox(x + 0.52, y, cw - 0.62, ch, anchor=MSO_ANCHOR.MIDDLE)
    par(tf, name, F_BOLD, 12, WHITE, bold=True, first=True, ls=0.95)
    p = par(tf, date, F_SEMI, 9.5, color, sb=2)
    add_run(p, "   ·   " + city, F_BODY, 9, MIST)

# ================= footer =================
tf = textbox(0.6, 7.04, 6.0, 0.3, anchor=MSO_ANCHOR.MIDDLE)
par(tf, "WWW.RENMAD.COM", F_SEMI, 10, WHITE, first=True, tracking=200)

prs.save(OUT)
print("Saved:", OUT)
