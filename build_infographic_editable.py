# -*- coding: utf-8 -*-
"""RENMAD Events 2026-2027 infographic — EDITABLE (native pptx shapes/text),
orange + charcoal, NO blue. Icons & flags are small PNGs; everything else
is editable PowerPoint."""
import os, math
from PIL import Image, ImageDraw

# ---------------- palette (NO blue) ----------------
ORANGE   = (240, 84, 18)
ORANGE_L = (255, 120, 38)
DARK     = (26, 26, 30)      # charcoal (replaces navy)
DARK_2   = (38, 38, 44)
WHITE    = (255, 255, 255)
INK      = (28, 32, 40)
GREY     = (110, 117, 126)
PANEL    = (245, 247, 249)
ROW_A    = (236, 240, 244)
LINE     = (214, 220, 226)

ICDIR = r"C:\Users\Belén\renmad_generator\assets\ic"
LOGOS = r"C:\Users\Belén\renmad_generator\assets\logos"
os.makedirs(ICDIR, exist_ok=True)

# ================= ICON DRAWING =================
def render_icon(fn, px, color):
    T = px * 3
    im = Image.new("RGBA", (T, T), (0, 0, 0, 0))
    d = ImageDraw.Draw(im)
    fn(d, T, color)
    return im.resize((px, px), Image.LANCZOS)

def _w(T, k=0.085):
    return max(2, int(T * k))

def ic_server(d,T,c):
    m=T*0.18; aw=T-2*m; bh=(T-2*m)*0.26; gap=bh*0.42; y=m
    for _ in range(3):
        d.rounded_rectangle([m,y,m+aw,y+bh], radius=bh*0.28, outline=c, width=_w(T,0.07))
        d.ellipse([m+aw*0.10-bh*0.12,y+bh/2-bh*0.12,m+aw*0.10+bh*0.12,y+bh/2+bh*0.12], fill=c)
        d.line([m+aw*0.30,y+bh/2,m+aw*0.85,y+bh/2], fill=c, width=_w(T,0.05)); y+=bh+gap

def ic_battery(d,T,c):
    m=T*0.16; w=T-2*m; h=w*0.56; x0=m; y0=(T-h)/2
    d.rounded_rectangle([x0,y0,x0+w*0.9,y0+h], radius=h*0.18, outline=c, width=_w(T,0.075))
    d.rounded_rectangle([x0+w*0.9,y0+h*0.3,x0+w,y0+h*0.7], radius=h*0.1, fill=c)
    cx=x0+w*0.45
    d.polygon([(cx+h*0.10,y0+h*0.18),(cx-h*0.18,y0+h*0.56),(cx+h*0.02,y0+h*0.56),
               (cx-h*0.10,y0+h*0.86),(cx+h*0.22,y0+h*0.44),(cx+h*0.02,y0+h*0.44)], fill=c)

def ic_flame(d,T,c):
    cx=T/2; m=T*0.16
    d.polygon([(cx,m),(cx+T*0.20,T*0.36),(cx+T*0.26,T*0.60),(cx+T*0.16,T*0.80),
               (cx,T-m),(cx-T*0.16,T*0.80),(cx-T*0.26,T*0.60),(cx-T*0.18,T*0.40)], fill=c)
    d.polygon([(cx,T*0.46),(cx+T*0.11,T*0.64),(cx,T*0.82),(cx-T*0.11,T*0.64)], fill=(255,255,255,235))

def ic_globe(d,T,c):
    m=T*0.16; r=(T-2*m)/2; cx=cy=T/2; wln=_w(T,0.06)
    d.ellipse([cx-r,cy-r,cx+r,cy+r], outline=c, width=wln)
    d.line([cx-r,cy,cx+r,cy], fill=c, width=wln)
    d.ellipse([cx-r*0.5,cy-r,cx+r*0.5,cy+r], outline=c, width=wln)
    d.arc([cx-r,cy-r*0.55,cx+r,cy+r*1.55], 200,340, fill=c, width=wln)
    d.arc([cx-r,cy-r*1.55,cx+r,cy+r*0.55], 20,160, fill=c, width=wln)

def ic_chart(d,T,c):
    m=T*0.18; wln=_w(T,0.07)
    d.line([m,m,m,T-m], fill=c, width=wln); d.line([m,T-m,T-m,T-m], fill=c, width=wln)
    a=T-2*m
    pts=[(m+a*0.05,T-m-a*0.18),(m+a*0.38,T-m-a*0.46),(m+a*0.60,T-m-a*0.30),(T-m-a*0.04,m+a*0.05)]
    d.line(pts, fill=c, width=_w(T,0.075), joint="curve")
    ax,ay=pts[-1]; d.polygon([(ax,ay),(ax-T*0.12,ay-T*0.01),(ax-T*0.02,ay+T*0.12)], fill=c)

def ic_calendar(d,T,c):
    m=T*0.17; wln=_w(T,0.065); x0,y0,x1,y1=m,m+T*0.06,T-m,T-m
    d.rounded_rectangle([x0,y0,x1,y1], radius=T*0.06, outline=c, width=wln)
    d.line([x0,y0+T*0.16,x1,y0+T*0.16], fill=c, width=wln)
    d.line([x0+T*0.22,m,x0+T*0.22,y0+T*0.10], fill=c, width=wln)
    d.line([x1-T*0.22,m,x1-T*0.22,y0+T*0.10], fill=c, width=wln)
    r=T*0.035
    for ix in range(3):
        for iy in range(2):
            cxp=x0+T*0.18+ix*T*0.24; cyp=y0+T*0.34+iy*T*0.22
            d.ellipse([cxp-r,cyp-r,cxp+r,cyp+r], fill=c)

def ic_people(d,T,c):
    def fig(cx,scale,col):
        hr=T*0.13*scale; hy=T*0.34
        d.ellipse([cx-hr,hy-hr,cx+hr,hy+hr], fill=col)
        bw=T*0.30*scale; d.pieslice([cx-bw/2,hy+hr*0.4,cx+bw/2,hy+hr*0.4+bw],180,360,fill=col)
    fig(T*0.62,0.92, tuple(int(v*0.55+255*0.45) for v in c[:3])); fig(T*0.42,1.0,c)

def ic_pin(d,T,c):
    cx=T/2; cy=T*0.40; r=T*0.22; wln=_w(T,0.08)
    d.arc([cx-r,cy-r,cx+r,cy+r],0,360,fill=c,width=wln)
    tip=(cx,T*0.86)
    d.line([(cx-r*0.86,cy+r*0.5),tip],fill=c,width=wln); d.line([(cx+r*0.86,cy+r*0.5),tip],fill=c,width=wln)
    d.ellipse([cx-r*0.34,cy-r*0.34,cx+r*0.34,cy+r*0.34],fill=c)

def ic_network(d,T,c):
    nodes=[(T*0.30,T*0.32),(T*0.72,T*0.28),(T*0.5,T*0.68),(T*0.78,T*0.66)]
    for a in range(len(nodes)):
        for b in range(a+1,len(nodes)):
            d.line([nodes[a],nodes[b]],fill=c,width=_w(T,0.04))
    for (x,y) in nodes:
        rr=T*0.075; d.ellipse([x-rr,y-rr,x+rr,y+rr],fill=c)

def ic_mic(d,T,c):
    cx=T/2; wln=_w(T,0.07)
    d.rounded_rectangle([cx-T*0.13,T*0.16,cx+T*0.13,T*0.58],radius=T*0.13,fill=c)
    d.arc([cx-T*0.22,T*0.30,cx+T*0.22,T*0.66],20,160,fill=c,width=wln)
    d.line([cx,T*0.66,cx,T*0.80],fill=c,width=wln); d.line([cx-T*0.16,T*0.82,cx+T*0.16,T*0.82],fill=c,width=wln)

def ic_bulb(d,T,c):
    cx=T/2; r=T*0.24; wln=_w(T,0.07)
    d.arc([cx-r,T*0.14,cx+r,T*0.14+2*r],0,360,fill=c,width=wln)
    d.line([cx-r*0.5,T*0.14+r*1.5,cx-r*0.32,T*0.74],fill=c,width=wln)
    d.line([cx+r*0.5,T*0.14+r*1.5,cx+r*0.32,T*0.74],fill=c,width=wln)
    for yy in (0.76,0.83): d.line([cx-T*0.13,T*yy,cx+T*0.13,T*yy],fill=c,width=wln)
    d.line([cx-T*0.09,T*0.90,cx+T*0.09,T*0.90],fill=c,width=wln)

def ic_turbine(d,T,c):
    cx=T*0.5; base=T*0.86; wln=_w(T,0.07); top=T*0.30
    d.line([(cx,base),(cx,top)],fill=c,width=wln)
    for a in (90,210,330):
        ar=math.radians(a); d.line([(cx,top),(cx+math.cos(ar)*T*0.22,top-math.sin(ar)*T*0.22)],fill=c,width=wln)
    d.ellipse([cx-T*0.03,top-T*0.03,cx+T*0.03,top+T*0.03],fill=c)

def ic_solar(d,T,c):
    cx=T*0.5; base=T*0.82; wln=_w(T,0.06); w=T*0.5
    d.line([(cx,base),(cx,base-T*0.16)],fill=c,width=wln)
    pts=[(cx-w/2,base-T*0.16),(cx+w/2,base-T*0.26),(cx+w/2,base-T*0.42),(cx-w/2,base-T*0.32)]
    d.line(pts+[pts[0]],fill=c,width=wln)

# ================= FLAGS =================
def flag(kind,w=300,h=200):
    im=Image.new("RGBA",(w,h),(255,255,255,255)); d=ImageDraw.Draw(im)
    if kind=="es":
        d.rectangle([0,0,w,h],fill=(198,11,30)); d.rectangle([0,h*0.25,w,h*0.75],fill=(255,196,0))
    elif kind=="it":
        d.rectangle([0,0,w/3,h],fill=(0,140,69)); d.rectangle([w/3,0,2*w/3,h],fill=(255,255,255)); d.rectangle([2*w/3,0,w,h],fill=(206,43,55))
    elif kind=="pl":
        d.rectangle([0,0,w,h/2],fill=(255,255,255)); d.rectangle([0,h/2,w,h],fill=(220,20,60))
    elif kind=="cl":
        d.rectangle([0,0,w,h/2],fill=(255,255,255)); d.rectangle([0,h/2,w,h],fill=(213,43,30)); d.rectangle([0,0,h/2,h/2],fill=(0,57,166))
        cx,cy,r=h/4,h/4,h*0.17; pts=[]
        for i in range(10):
            ang=-math.pi/2+i*math.pi/5; rr=r if i%2==0 else r*0.42
            pts.append((cx+math.cos(ang)*rr,cy+math.sin(ang)*rr))
        d.polygon(pts,fill=(255,255,255))
    elif kind=="mx":
        d.rectangle([0,0,w/3,h],fill=(0,104,71)); d.rectangle([w/3,0,2*w/3,h],fill=(255,255,255)); d.rectangle([2*w/3,0,w,h],fill=(206,17,38))
        d.ellipse([w/2-h*0.12,h/2-h*0.12,w/2+h*0.12,h/2+h*0.12],outline=(120,80,30),width=4)
    d.rectangle([0,0,w-1,h-1],outline=(150,156,162),width=3)
    return im

# ---------- generate PNGs ----------
white_icons=["server","battery","flame","globe","chart"]
orange_icons=["calendar","pin","globe","people","network","mic","bulb","turbine","solar"]
fns={"server":ic_server,"battery":ic_battery,"flame":ic_flame,"globe":ic_globe,
     "chart":ic_chart,"calendar":ic_calendar,"pin":ic_pin,"people":ic_people,
     "network":ic_network,"mic":ic_mic,"bulb":ic_bulb,"turbine":ic_turbine,"solar":ic_solar}
for n in white_icons:
    render_icon(fns[n],360,WHITE).save(os.path.join(ICDIR,f"w_{n}.png"))
for n in orange_icons:
    render_icon(fns[n],360,ORANGE).save(os.path.join(ICDIR,f"o_{n}.png"))
for fl in ["es","it","cl","mx","pl"]:
    flag(fl).save(os.path.join(ICDIR,f"flag_{fl}.png"))
print("icons generated")

# ================= BUILD PPTX (editable) =================
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
slide=prs.slides.add_slide(prs.slide_layouts[6])

DPI=240.0
def IN(px): return Inches(px/DPI)
def PT(px): return Pt(px*0.30)
RGB=lambda c: RGBColor(*c)
FB="Montserrat Black"; FM="Montserrat"; FS="Montserrat SemiBold"

def no_shadow(sh):
    sp=sh._element.spPr
    for e in sp.findall(qn("a:effectLst")): sp.remove(e)
    sp.append(sp.makeelement(qn("a:effectLst"),{}))

def shape(kind,px_x,px_y,px_w,px_h,color,rad=None,line=None,lw=1.0):
    s=slide.shapes.add_shape(kind,IN(px_x),IN(px_y),IN(px_w),IN(px_h))
    if color is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb=RGB(color)
    if line is None: s.line.fill.background()
    else: s.line.color.rgb=RGB(line); s.line.width=Pt(lw)
    if kind==MSO_SHAPE.ROUNDED_RECTANGLE and rad is not None:
        try: s.adjustments[0]=rad
        except Exception: pass
    no_shadow(s); return s

def rrect(x,y,w,h,color,rad=0.12,line=None,lw=1.0): return shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h,color,rad,line,lw)
def rect(x,y,w,h,color,line=None,lw=1.0): return shape(MSO_SHAPE.RECTANGLE,x,y,w,h,color,None,line,lw)
def vline(x,y1,y2,color=LINE,wpx=2): rect(x,y1,wpx,y2-y1,color)
def oval(x,y,w,h,color): return shape(MSO_SHAPE.OVAL,x,y,w,h,color)

def tbox(px_x,px_y,px_w,px_h,anchor=MSO_ANCHOR.TOP):
    tb=slide.shapes.add_textbox(IN(px_x),IN(px_y),IN(px_w),IN(px_h))
    tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
    tf.vertical_anchor=anchor; return tf

def para(tf,align=PP_ALIGN.LEFT,first=False,sb=0,sa=0,ls=1.0):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_before=Pt(sb); p.space_after=Pt(sa); p.line_spacing=ls; return p

def run(p,text,font,size_px,color,bold=False,tracking=None,sub=False):
    r=p.add_run(); r.text=text; r.font.name=font; r.font.size=PT(size_px)
    r.font.bold=bold; r.font.color.rgb=RGB(color)
    rPr=r._r.get_or_add_rPr()
    if tracking is not None: rPr.set("spc",str(tracking))
    if sub: rPr.set("baseline","-25000")
    return r

def img(path,px_x,px_y,px_w,px_h=None):
    if px_h is None: slide.shapes.add_picture(path,IN(px_x),IN(px_y),width=IN(px_w))
    else: slide.shapes.add_picture(path,IN(px_x),IN(px_y),width=IN(px_w),height=IN(px_h))

def ic(name,px_x,px_y,px_sz,white=False):
    pre="w_" if white else "o_"
    img(os.path.join(ICDIR,pre+name+".png"),px_x,px_y,px_sz,px_sz)

W,H=3200,1800
# ---------------- HEADER ----------------
HH=648
rect(0,0,W,HH,DARK)
rect(0,HH-6,W,6,ORANGE)                       # accent rule
# logo
lg=Image.open(os.path.join(LOGOS,"logo_renmad_events.png")); lr=lg.height/lg.width
img(os.path.join(LOGOS,"logo_renmad_events.png"),56,52,380)
# left energy deco
ic("turbine",70,360,150); ic("solar",250,400,150); ic("turbine",430,360,150)
# AI cluster (orange) right
nodes=[(2770,250),(2980,205),(3070,330),(2880,395),(3060,470)]
for a in range(len(nodes)):
    for b in range(a+1,len(nodes)):
        if (a+b)%2==0:
            x1,y1=nodes[a]; x2,y2=nodes[b]
            cn=slide.shapes.add_connector(2,IN(x1),IN(y1),IN(x2),IN(y2))
            cn.line.color.rgb=RGB(ORANGE); cn.line.width=Pt(1.2); no_shadow(cn)
for (x,y) in nodes: oval(x-9,y-9,18,18,ORANGE_L)
tf=tbox(2870,250,260,120,MSO_ANCHOR.MIDDLE); p=para(tf,PP_ALIGN.CENTER,first=True)
run(p,"AI",FB,128,ORANGE_L)

# title
CXT=1580; BW=2200; bx=CXT-BW/2
tf=tbox(bx,52,BW,130,MSO_ANCHOR.MIDDLE); p=para(tf,PP_ALIGN.CENTER,first=True)
run(p,"RENMAD ",FB,96,WHITE); run(p,"EVENTS",FB,96,ORANGE_L); run(p," —",FB,96,WHITE)
tf=tbox(bx,170,BW,130,MSO_ANCHOR.MIDDLE); p=para(tf,PP_ALIGN.CENTER,first=True)
run(p,"EVENTS CALENDAR",FB,96,WHITE)
tf=tbox(bx,300,BW,60,MSO_ANCHOR.MIDDLE); p=para(tf,PP_ALIGN.CENTER,first=True)
run(p,"(International Series)",FS,46,ORANGE_L)
tf=tbox(bx,372,BW,48,MSO_ANCHOR.MIDDLE); p=para(tf,PP_ALIGN.CENTER,first=True)
run(p,"RENEWABLES · STORAGE · HYDROGEN · BIOMETHANE · DATA CENTERS · AI",FS,34,WHITE,tracking=40)
tf=tbox(bx,420,BW,70,MSO_ANCHOR.MIDDLE); p=para(tf,PP_ALIGN.CENTER,first=True)
run(p,"EVERY FRONTIER OF ENERGY",FB,60,ORANGE_L)

# stat row
sy=528; stats=[("calendar","2026 – 2027"),("people","14 summits"),("globe","5 countries"),("pin","2 continents")]
total_w=1820; sx=CXT-total_w/2; step=total_w/4
for i,(icon,label) in enumerate(stats):
    cx=sx+step*i+50
    ic(icon,cx-29,sy-29,58)
    tf=tbox(cx+44,sy-30,step-110,60,MSO_ANCHOR.MIDDLE); p=para(tf,first=True)
    run(p,label,FS,42,WHITE)
    if i>0: vline(sx+step*i-20,sy-34,sy+34,(120,120,128))

# ---------------- BODY ----------------
Lx0=48; Lw=1480; by=HH+44; barh=78
Rx0=1568; Rw=1584
def panel_header(x,w,y,h,color,label):
    rrect(x,y,w,h,color,rad=0.18)
    tf=tbox(x,y,w,h,MSO_ANCHOR.MIDDLE); p=para(tf,PP_ALIGN.CENTER,first=True)
    run(p,label,FB,40,WHITE)
panel_header(Lx0,Lw,by,barh,ORANGE,"STILL TO COME IN 2026")
panel_header(Rx0,Rw,by,barh,DARK,"THE 2027 SEASON  (12 SUMMITS)")

# LEFT events
py=by+barh+24; rowh=128
rrect(Lx0,py,Lw,rowh*2+24,PANEL,rad=0.06)
ev2026=[("server","Data Centers Italia","Nov 11–12","Milan, Italy","it",False),
        ("h2","Hidrógeno","Nov 18–19","Zaragoza, Spain","es",False)]
def h2_tile(tx,ty2,ts,big,small):
    tf=tbox(tx,ty2,ts,ts,MSO_ANCHOR.MIDDLE); p=para(tf,PP_ALIGN.CENTER,first=True)
    run(p,"H",FB,big,WHITE); run(p,"2",FB,small,WHITE,sub=True)
for i,(k,n,dt,lc,fl,dark) in enumerate(ev2026):
    y=py+12+i*rowh; cy=y+rowh/2; ts=86; tx=Lx0+34; ty2=y+(rowh-ts)/2
    rrect(tx,ty2,ts,ts,ORANGE,rad=0.18)
    if k=="h2": h2_tile(tx,ty2,ts,48,30)
    else: ic(k,tx+ts*0.19,ty2+ts*0.19,ts*0.62,white=True)
    tf=tbox(tx+ts+34,y,560,rowh,MSO_ANCHOR.MIDDLE); p=para(tf,first=True); run(p,n,FM,42,INK,bold=True)
    vline(Lx0+660,cy-30,cy+30); ic("calendar",Lx0+690,cy-23,46)
    tf=tbox(Lx0+748,y,200,rowh,MSO_ANCHOR.MIDDLE); p=para(tf,first=True); run(p,dt,FS,36,INK)
    vline(Lx0+930,cy-30,cy+30); ic("pin",Lx0+958,cy-23,46)
    tf=tbox(Lx0+1012,y,300,rowh,MSO_ANCHOR.MIDDLE); p=para(tf,first=True); run(p,lc,FS,32,INK)
    img(os.path.join(ICDIR,f"flag_{fl}.png"),Lx0+Lw-116,cy-30,92,61)

# value props
vy=py+rowh*2+24+70
vals=[("globe","GLOBAL","PERSPECTIVE"),("people","SENIOR","AUDIENCE"),
      ("network","STRATEGIC","NETWORKING"),("mic","WORLD-CLASS","SPEAKERS"),
      ("bulb","REAL SOLUTIONS","FOR THE FUTURE")]
vstep=Lw/5
for i,(icon,l1,l2) in enumerate(vals):
    cx=Lx0+vstep*i+vstep/2
    ic(icon,cx-48,vy,96)
    tf=tbox(cx-vstep/2,vy+116,vstep,90,MSO_ANCHOR.TOP)
    p=para(tf,PP_ALIGN.CENTER,first=True); run(p,l1,FM,28,INK,bold=True)
    p=para(tf,PP_ALIGN.CENTER,sb=2); run(p,l2,FM,28,INK,bold=True)

# RIGHT events
ev2027=[("server","Data Centers","27–28 Jan","Zaragoza, Spain","es",True),
        ("battery","Storage Italia","9–10 Feb","Bologna, Italy","it",False),
        ("flame","Biometano","10–11 Feb","Toledo, Spain","es",False),
        ("battery","Almacenamiento","31 Mar – 1 Apr","Seville, Spain","es",False),
        ("globe","Chile","29–30 Jul","Santiago, Chile","cl",False),
        ("h2","Hidrógeno","November","Zaragoza, Spain","es",False),
        ("server","Data Centers Italia","November","Milan, Italy","it",True),
        ("chart","Invest","Dates soon","Madrid, Spain","es",False),
        ("chart","Invest Italia","Dates soon","Milan, Italy","it",False),
        ("globe","México","Dates soon","Mexico","mx",False),
        ("battery","Storage Polska","Dates soon","Warsaw, Poland","pl",False),
        ("ai","Useful AI","Dates soon","Madrid, Spain","es",True)]
ry0=by+barh+18; ravail=(H-116)-ry0; rh=ravail/12
for i,(k,n,dt,lc,fl,dark) in enumerate(ev2027):
    y=ry0+i*rh; cy=y+rh/2
    if i%2==1: rect(Rx0,y,Rw,rh,ROW_A)
    ts=64; tx=Rx0+22; ty2=cy-ts/2
    rrect(tx,ty2,ts,ts,DARK if dark else ORANGE,rad=0.16)
    if k=="h2":
        tf=tbox(tx,ty2,ts,ts,MSO_ANCHOR.MIDDLE); p=para(tf,PP_ALIGN.CENTER,first=True)
        run(p,"H",FB,34,WHITE); run(p,"2",FB,22,WHITE,sub=True)
    elif k=="ai":
        tf=tbox(tx,ty2,ts,ts,MSO_ANCHOR.MIDDLE); p=para(tf,PP_ALIGN.CENTER,first=True); run(p,"AI",FB,30,WHITE)
    else: ic(k,tx+ts*0.2,ty2+ts*0.2,ts*0.6,white=True)
    tf=tbox(tx+ts+26,y,520,rh,MSO_ANCHOR.MIDDLE); p=para(tf,first=True); run(p,n,FM,34,INK,bold=True)
    vline(Rx0+640,cy-26,cy+26); ic("calendar",Rx0+666,cy-20,40)
    tf=tbox(Rx0+716,y,270,rh,MSO_ANCHOR.MIDDLE); p=para(tf,first=True); run(p,dt,FS,30,INK)
    vline(Rx0+980,cy-26,cy+26); ic("pin",Rx0+1006,cy-20,40)
    tf=tbox(Rx0+1052,y,360,rh,MSO_ANCHOR.MIDDLE); p=para(tf,first=True); run(p,lc,FS,28,INK)
    img(os.path.join(ICDIR,f"flag_{fl}.png"),Rx0+Rw-104,cy-28,84,56)

# ---------------- FOOTER ----------------
FY=H-104
rect(0,FY,W,H-FY,DARK)
fcy=FY+52
tf=tbox(60,FY,1400,H-FY,MSO_ANCHOR.MIDDLE); p=para(tf,first=True)
run(p,"CONNECTING MARKETS. ",FB,40,ORANGE_L); run(p,"ACCELERATING TRANSITION.",FB,40,WHITE)
ic("globe",W-960,fcy-23,46)
tf=tbox(W-900,FY,360,H-FY,MSO_ANCHOR.MIDDLE); p=para(tf,first=True); run(p,"www.renmad.com",FS,36,WHITE)
vline(W-560,fcy-26,fcy+26,(120,120,128))
rrect(W-520,fcy-26,52,52,ORANGE_L,rad=0.2)
tf=tbox(W-520,fcy-26,52,52,MSO_ANCHOR.MIDDLE); p=para(tf,PP_ALIGN.CENTER,first=True); run(p,"in",FB,38,DARK)
tf=tbox(W-446,FY,420,H-FY,MSO_ANCHOR.MIDDLE); p=para(tf,first=True); run(p,"RENMAD EVENTS",FM,36,WHITE,bold=True)

OUT=r"C:\Users\Belén\RENMAD_Events_Infographic_Editable.pptx"
prs.save(OUT)
print("Saved:",OUT)
